"""
金錢小達人 競賽申請 PPT 生成腳本
視覺風格：教科書圖解 × 教育知識科普 × 圖表型資訊設計
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
import copy
from lxml import etree

# ── 顏色系統 ──────────────────────────────────────────────────
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK       = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
C_GRAY       = RGBColor(0x64, 0x74, 0x8B)   # slate-500
C_LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFC)   # nearly white
C_PANEL      = RGBColor(0xF1, 0xF5, 0xF9)   # slate-100

C_BLUE       = RGBColor(0x22, 0x63, 0xEB)   # F series
C_BLUE_L     = RGBColor(0xDB, 0xEA, 0xFE)   # F series light
C_GREEN      = RGBColor(0x05, 0x96, 0x69)   # C series
C_GREEN_L    = RGBColor(0xD1, 0xFA, 0xE5)   # C series light
C_ORANGE     = RGBColor(0xEA, 0x58, 0x0C)   # A series
C_ORANGE_L   = RGBColor(0xFF, 0xED, 0xD5)   # A series light
C_PURPLE     = RGBColor(0x7C, 0x3A, 0xED)   # accent
C_PURPLE_L   = RGBColor(0xED, 0xE9, 0xFE)   # accent light
C_AMBER      = RGBColor(0xD9, 0x77, 0x06)   # highlight
C_AMBER_L    = RGBColor(0xFE, 0xF3, 0xC7)   # highlight light
C_RED        = RGBColor(0xDC, 0x26, 0x26)
C_RED_L      = RGBColor(0xFE, 0xE2, 0xE2)
C_TEAL       = RGBColor(0x0D, 0x94, 0x88)
C_TEAL_L     = RGBColor(0xCC, 0xFB, 0xF1)

# ── 投影片尺寸（16:9 寬螢幕）────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # 完全空白

# ══════════════════════════════════════════════════════════════
#  工具函數
# ══════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill=None, border_color=None, border_pt=0, radius=None):
    """新增色塊，可選圓角"""
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color and border_pt > 0:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    if radius is not None:
        # 設定圓角
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=16, bold=False, color=None, align=PP_ALIGN.LEFT,
             font_name='微軟正黑體', italic=False, wrap=True):
    """新增文字方塊"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def add_multiline(slide, lines, x, y, w, h,
                  font_size=14, color=None, line_spacing=1.3,
                  font_name='微軟正黑體', align=PP_ALIGN.LEFT):
    """新增多行文字（list of strings）"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as PT2
    from pptx.oxml.ns import qn as QN2
    from lxml import etree as ET2
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = PT2(font_size)
        if color:
            run.font.color.rgb = color
    return txBox

def add_pill(slide, text, x, y, w, h, bg, fg, font_size=12, bold=True):
    """圓角藥丸標籤"""
    r = add_rect(slide, x, y, w, h, fill=bg, radius=20000)
    t = add_text(slide, text, x, y, w, h,
                 font_size=font_size, bold=bold, color=fg,
                 align=PP_ALIGN.CENTER, font_name='微軟正黑體')
    return r, t

def horizontal_divider(slide, y, x_start=Inches(0.5), x_end=Inches(12.83),
                        color=C_PANEL, pt=1.5):
    """水平分隔線"""
    line = slide.shapes.add_connector(1, x_start, y, x_end, y)
    line.line.color.rgb = color
    line.line.width = Pt(pt)

def add_section_header(slide, title_text, subtitle_text='', bg=C_BLUE, x=Inches(0), w=Inches(13.33)):
    """頁面頂部標題條"""
    add_rect(slide, x, Inches(0), w, Inches(1.15), fill=bg)
    add_text(slide, title_text,
             x + Inches(0.45), Inches(0.12), w - Inches(0.9), Inches(0.6),
             font_size=26, bold=True, color=C_WHITE, font_name='微軟正黑體')
    if subtitle_text:
        add_text(slide, subtitle_text,
                 x + Inches(0.45), Inches(0.68), w - Inches(0.9), Inches(0.4),
                 font_size=14, color=RGBColor(0xBF, 0xDB, 0xFF), font_name='微軟正黑體')

def add_callout_box(slide, icon, title, body, x, y, w, h,
                    bg=C_BLUE_L, border=C_BLUE, title_c=C_BLUE, body_c=C_DARK):
    """資訊框：左側圖示 + 標題 + 內文"""
    add_rect(slide, x, y, w, h, fill=bg, border_color=border, border_pt=1.5, radius=15000)
    add_text(slide, icon,  x + Inches(0.15), y + Inches(0.1),
             Inches(0.55), h - Inches(0.2), font_size=28, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.75), y + Inches(0.08),
             w - Inches(0.9), Inches(0.38), font_size=14, bold=True, color=title_c)
    add_multiline(slide, body.split('\n'), x + Inches(0.75), y + Inches(0.42),
                  w - Inches(0.9), h - Inches(0.52), font_size=12, color=body_c)

# ══════════════════════════════════════════════════════════════
#  SLIDE 01  封面
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=RGBColor(0x0F, 0x17, 0x2A))

# 左色條
add_rect(sl, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=C_ORANGE)

# 右側裝飾圓
for r_data in [(Inches(10.8), Inches(-0.5), Inches(3.2), C_BLUE, 60),
               (Inches(11.5), Inches(5.2), Inches(2.4), C_GREEN, 50),
               (Inches(9.5), Inches(3.5), Inches(1.6), C_PURPLE, 30)]:
    rx, ry, rsize, rc, alpha = r_data
    shape = add_rect(sl, rx, ry, rsize, rsize, fill=RGBColor(rc[0]//3, rc[1]//3, rc[2]//3))
    # 設圓形
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'ellipse')

# 競賽標籤
add_pill(sl, '第八屆教育部獎助研發特殊教育教材教具競賽',
         Inches(0.55), Inches(1.0), Inches(6.2), Inches(0.42),
         bg=C_ORANGE, fg=C_WHITE, font_size=13)

# 主標題
add_text(sl, '金 錢 小 達 人',
         Inches(0.55), Inches(1.65), Inches(9), Inches(1.2),
         font_size=60, bold=True, color=C_WHITE, font_name='微軟正黑體')

# 副標題
add_text(sl, '互動式特殊教育生活數學暨金融素養教學軟體',
         Inches(0.55), Inches(2.80), Inches(9), Inches(0.55),
         font_size=22, color=RGBColor(0xFD, 0xBA, 0x74), font_name='微軟正黑體')

add_text(sl, 'Money Tutor: Interactive Software for Financial Literacy in Special Education',
         Inches(0.55), Inches(3.35), Inches(9.5), Inches(0.45),
         font_size=13, color=RGBColor(0x94, 0xA3, 0xB8), font_name='Arial')

# 橫線
horizontal_divider(sl, Inches(3.9), Inches(0.55), Inches(9.2),
                   color=RGBColor(0x33, 0x4A, 0x6E), pt=1.5)

# 統計卡片
stats = [
    ('18', '互動\n教學單元'),
    ('150,000+', '程式碼行數'),
    ('4', '研發月份'),
    ('0', '授權費\n（新台幣）'),
]
for i, (num, label) in enumerate(stats):
    sx = Inches(0.55 + i * 2.15)
    add_rect(sl, sx, Inches(4.1), Inches(1.95), Inches(0.95),
             fill=RGBColor(0x1E, 0x3A, 0x5C), radius=12000)
    add_text(sl, num, sx, Inches(4.1), Inches(1.95), Inches(0.48),
             font_size=28, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_text(sl, label, sx, Inches(4.52), Inches(1.95), Inches(0.5),
             font_size=11, color=RGBColor(0x94, 0xA3, 0xB8),
             align=PP_ALIGN.CENTER, font_name='微軟正黑體')

# 色系說明
series_info = [
    (C_BLUE, 'F 系列  數學基礎'),
    (C_GREEN, 'C 系列  貨幣認知'),
    (C_ORANGE, 'A 系列  情境類化'),
]
for i, (c, t) in enumerate(series_info):
    sx = Inches(0.55 + i * 2.9)
    add_rect(sl, sx, Inches(5.35), Inches(0.18), Inches(0.32), fill=c)
    add_text(sl, t, sx + Inches(0.25), Inches(5.35), Inches(2.5), Inches(0.32),
             font_size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

add_text(sl, '適用對象：國中至高中特教班・資源班\n輕度至中度智能障礙・自閉症・學習障礙・多重障礙',
         Inches(0.55), Inches(5.9), Inches(9), Inches(0.7),
         font_size=13, color=RGBColor(0x7C, 0x8D, 0xA8), font_name='微軟正黑體')

add_text(sl, '中華民國 115 年 3 月',
         Inches(0.55), Inches(6.75), Inches(4), Inches(0.4),
         font_size=12, color=RGBColor(0x47, 0x5A, 0x73))

# ══════════════════════════════════════════════════════════════
#  SLIDE 02  目錄
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '簡報綱要', 'Table of Contents', bg=C_DARK)

items = [
    (C_RED,    '01', '設計動機', '特教現場的三重結構性困境'),
    (C_AMBER,  '02', '設計理念', 'UDL × 工作分析法 × 直接教學法'),
    (C_BLUE,   '03', 'F 系列', '數學基礎先備（6 個單元）'),
    (C_GREEN,  '04', 'C 系列', '貨幣認知建立（6 個單元）'),
    (C_ORANGE, '05', 'A 系列', '真實情境類化（6 個單元）'),
    (C_PURPLE, '06', '特色功能', '輔助點擊・語音・深色模式・作業單'),
    (C_TEAL,   '07', '三段學習歷程', '數位練習 → 紙本評量 → 社區類化'),
    (C_BLUE,   '08', '課程法規依據', 'IEP / 特教法 / 課綱對應'),
    (C_GREEN,  '09', '試教觀察與預期成效', '個案轉化 · 成效指標'),
]

cols = [0, 1, 2]
rows = [0, 1, 2]
items_per_col = 3
for idx, (c, num, title, sub) in enumerate(items):
    col = idx % 3
    row = idx // 3
    bx = Inches(0.5 + col * 4.27)
    by = Inches(1.35 + row * 1.95)
    add_rect(sl, bx, by, Inches(4.0), Inches(1.75),
             fill=C_WHITE, border_color=RGBColor(0xE2, 0xE8, 0xF0), border_pt=1.2, radius=12000)
    # 色條
    add_rect(sl, bx, by, Inches(4.0), Inches(0.12), fill=c, radius=0)
    add_text(sl, num, bx + Inches(0.18), by + Inches(0.2),
             Inches(0.6), Inches(0.5), font_size=28, bold=True,
             color=c, font_name='Arial')
    add_text(sl, title, bx + Inches(0.82), by + Inches(0.2),
             Inches(3.0), Inches(0.42), font_size=16, bold=True, color=C_DARK)
    add_text(sl, sub, bx + Inches(0.82), by + Inches(0.62),
             Inches(3.0), Inches(0.9), font_size=12, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
#  SLIDE 03  設計動機 — 三重困境
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '設計動機 ─── 特教現場的三重結構性困境', bg=C_RED)

problems = [
    (C_RED, C_RED_L,    '🧠', '困境一：交易序列的意義理解障礙',
     '能計算 100−35=65，卻無法判斷「給哪幾枚硬幣」\n'
     '靜態紙本缺乏即時個別化回饋\n'
     '幣值辨識 ≠ 交易序列理解'),
    (C_AMBER, C_AMBER_L, '⚠️', '困境二：安全練習場域的長期缺位',
     '真實場景試誤 → 消費糾紛 + 情緒失調風險\n'
     '一次挫折 → 長期情境迴避（小俊・速食店、小雯・販賣機）\n'
     '需要「失敗不帶社會後果」的數位練習環境'),
    (C_BLUE, C_BLUE_L,   '📚', '困境三：系統性課程工具的十年缺口',
     '課綱已列財務管理為核心主題（生管-Ⅳ-E-1～E-4）\n'
     '市場鮮有系統性・個別化可調整的中學特教電腦輔助教材\n'
     '教師普遍依賴自製圖卡，備課耗時且課程銜接脆弱'),
]

for i, (bc, bl, icon, title, body) in enumerate(problems):
    bx = Inches(0.4 + i * 4.3)
    add_callout_box(sl, icon, title, body,
                    bx, Inches(1.3), Inches(4.15), Inches(5.85),
                    bg=bl, border=bc, title_c=bc)

# 底部引用
add_rect(sl, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35),
         fill=RGBColor(0xFE, 0xE2, 0xE2), radius=8000)
add_text(sl, '《十二年國教特殊需求領域課程綱要》生活管理・財務管理主題 ／《特殊教育法》第 28 條 ／《身心障礙者權利公約》第 9、19 條',
         Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.32),
         font_size=11, color=C_RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  SLIDE 04  個案故事
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '真實個案──為何需要數位安全練習場域', bg=C_RED)

cases = [
    ('🧒 小涵', C_PURPLE, C_PURPLE_L,
     '輕度智能障礙・14 歲',
     '能穩定計算 100－35＝65，卻在模擬購物中反覆困惑\n「我有多少錢」與「這個東西多少錢」的連結\n錯誤模式在缺乏即時回饋的環境中持續固化'),
    ('🧒 小俊', C_AMBER, C_AMBER_L,
     '輕度自閉症・中學生',
     '在速食店自助點餐機前，因「漢堡要選哪個類別」\n陷入選擇迴路，後方催促誘發情緒失調\n此後長達一學期拒絕接觸相關學習活動'),
    ('🧒 小雯', C_ORANGE, C_ORANGE_L,
     '特殊需求學生',
     '在走廊販賣機反覆因硬幣面額選錯而失敗\n主動繞道迴避所有自助設備\n告訴老師：「那個機器我不會用」'),
]

for i, (name, bc, bl, type_str, body) in enumerate(cases):
    bx = Inches(0.4 + i * 4.3)
    add_rect(sl, bx, Inches(1.3), Inches(4.15), Inches(4.8),
             fill=bl, border_color=bc, border_pt=2, radius=15000)
    add_text(sl, name, bx, Inches(1.3), Inches(4.15), Inches(0.6),
             font_size=22, bold=True, color=bc, align=PP_ALIGN.CENTER)
    add_text(sl, type_str, bx, Inches(1.85), Inches(4.15), Inches(0.35),
             font_size=12, color=bc, align=PP_ALIGN.CENTER, font_name='微軟正黑體')
    horizontal_divider(sl, Inches(2.25), bx + Inches(0.3), bx + Inches(3.85), color=bc, pt=1)
    add_multiline(sl, body.split('\n'), bx + Inches(0.25), Inches(2.35),
                  Inches(3.65), Inches(3.5), font_size=13, color=C_DARK)

# 轉化
add_rect(sl, Inches(0.4), Inches(6.25), Inches(12.5), Inches(1.0),
         fill=RGBColor(0xEC, 0xFD, 0xF5), border_color=C_GREEN, border_pt=1.5, radius=12000)
add_text(sl, '✅  小雯完成 A1 十五次練習後，第一次在真實販賣機前完整購買，並說：「我下次自己來。」',
         Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.9),
         font_size=15, bold=True, color=C_GREEN, font_name='微軟正黑體')

# ══════════════════════════════════════════════════════════════
#  SLIDE 05  設計理念 ── UDL + 工作分析法
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '設計理念 ─── UDL × 工作分析法 × 直接教學法', bg=C_DARK)

# UDL 三欄
udl = [
    (C_BLUE, C_BLUE_L,   '📖', '多元表徵',  'Multiple Means of\nRepresentation',
     '全程中文語音引導（雅婷優先）\n高解析商品／服務圖片\n介面仿真真實設備\n短句操作提示（≤5 字）'),
    (C_GREEN, C_GREEN_L,  '✋', '多元行動\n與表達', 'Action &\nExpression',
     '觸控拖曳（平板）\n滑鼠點擊（電腦）\n輔助點擊模式（首創）\n三種難度對應褪除序列'),
    (C_PURPLE, C_PURPLE_L,'🎯', '多元參與',  'Engagement',
     '支援自訂商品圖片上傳\n即時語音讚美 + 煙火動畫\n獎勵積分系統整合 IEP\n深色模式（視覺敏感適用）'),
]

add_text(sl, 'UDL 通用設計學習三大原則（CAST, 2018）',
         Inches(0.5), Inches(1.2), Inches(8), Inches(0.4),
         font_size=15, bold=True, color=C_DARK)

for i, (bc, bl, icon, title, sub, body) in enumerate(udl):
    bx = Inches(0.4 + i * 3.3)
    add_rect(sl, bx, Inches(1.65), Inches(3.1), Inches(4.0),
             fill=bl, border_color=bc, border_pt=2, radius=15000)
    add_text(sl, icon, bx, Inches(1.7), Inches(3.1), Inches(0.65),
             font_size=30, align=PP_ALIGN.CENTER)
    add_text(sl, title, bx, Inches(2.3), Inches(3.1), Inches(0.55),
             font_size=17, bold=True, color=bc, align=PP_ALIGN.CENTER)
    add_text(sl, sub, bx, Inches(2.8), Inches(3.1), Inches(0.45),
             font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
    horizontal_divider(sl, Inches(3.3), bx + Inches(0.3), bx + Inches(2.8), color=bc, pt=1)
    add_multiline(sl, body.split('\n'), bx + Inches(0.2), Inches(3.42),
                  Inches(2.7), Inches(2.1), font_size=12, color=C_DARK)

# 工作分析法 + 直接教學法
add_rect(sl, Inches(10.15), Inches(1.2), Inches(3.0), Inches(5.95),
         fill=C_WHITE, border_color=RGBColor(0xE2, 0xE8, 0xF0), border_pt=1.2, radius=12000)
add_text(sl, '⚙️ 工作分析法', Inches(10.25), Inches(1.28), Inches(2.8), Inches(0.42),
         font_size=14, bold=True, color=C_DARK)
add_multiline(sl, [
    'A 系列 6 情境均以工作',
    '分析法（Gold, 1980）',
    '拆解真實設備操作流程',
    '步數：5～8 步/單元',
    '每步驟均可單獨評量',
    '完整對應 IEP 目標',
], Inches(10.25), Inches(1.72), Inches(2.75), Inches(2.2), font_size=12, color=C_DARK)

add_rect(sl, Inches(10.15), Inches(3.95), Inches(3.0), Inches(0.08), fill=C_PANEL)
add_text(sl, '📊 直接教學法三階段', Inches(10.25), Inches(4.05), Inches(2.8), Inches(0.42),
         font_size=14, bold=True, color=C_DARK)
diff_rows = [
    (C_BLUE_L,   C_BLUE,   '簡單', '明示教學\n光暈+語音引導'),
    (C_AMBER_L,  C_AMBER,  '普通', '引導練習\n三次錯誤後提示'),
    (C_RED_L,    C_RED,    '困難', '獨立練習\n完全無提示'),
]
for j, (bg, bc2, lvl, desc) in enumerate(diff_rows):
    ry = Inches(4.52 + j * 0.78)
    add_rect(sl, Inches(10.25), ry, Inches(2.75), Inches(0.7),
             fill=bg, border_color=bc2, border_pt=1, radius=8000)
    add_text(sl, lvl, Inches(10.32), ry + Inches(0.08),
             Inches(0.65), Inches(0.55), font_size=13, bold=True, color=bc2)
    add_multiline(sl, desc.split('\n'), Inches(11.0), ry + Inches(0.04),
                  Inches(1.9), Inches(0.65), font_size=11, color=C_DARK)

# ══════════════════════════════════════════════════════════════
#  SLIDE 06  三層次學習序階圖
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '三層次螺旋序階課程架構', '數學先備 → 貨幣認知 → 真實情境類化', bg=C_DARK)

# 金字塔（三層梯形）
pyramid_layers = [
    # (bg, border, text_main, text_sub, y, w)
    (C_ORANGE, C_ORANGE, 'A 系列  ·  真實情境類化', '社區設備操作 · 類化應用', Inches(1.5), Inches(5.8)),
    (C_GREEN,  C_GREEN,  'C 系列  ·  貨幣認知建立',  '認識幣值 · 付款找零',   Inches(3.0), Inches(7.8)),
    (C_BLUE,   C_BLUE,   'F 系列  ·  數學基礎先備',  '對應 · 唱數 · 排序',   Inches(4.5), Inches(9.8)),
]

for i, (bg, bc, main, sub, y, w) in enumerate(pyramid_layers):
    cx = (SLIDE_W - w) / 2
    alpha_bg = RGBColor(int(bg[0]*0.15+0xF8*0.85)//1,
                        int(bg[1]*0.15+0xFA*0.85)//1,
                        int(bg[2]*0.15+0xFC*0.85)//1)
    add_rect(sl, cx, y, w, Inches(1.35), fill=bg, border_color=bc, border_pt=2, radius=8000)
    add_text(sl, main, cx + Inches(0.3), y + Inches(0.12), w - Inches(0.6), Inches(0.6),
             font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, sub, cx + Inches(0.3), y + Inches(0.68), w - Inches(0.6), Inches(0.5),
             font_size=13, color=RGBColor(0xFF, 0xFF, 0xE0), align=PP_ALIGN.CENTER)

# 箭頭標籤（左側）
arrows = [
    (C_ORANGE, '6 個單元', Inches(0.4), Inches(1.9)),
    (C_GREEN,  '6 個單元', Inches(0.4), Inches(3.4)),
    (C_BLUE,   '6 個單元', Inches(0.4), Inches(4.9)),
]
for c, t, ax, ay in arrows:
    add_rect(sl, ax, ay, Inches(1.55), Inches(0.38), fill=c, radius=10000)
    add_text(sl, t, ax, ay, Inches(1.55), Inches(0.38),
             font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 右側說明
right_x = Inches(9.2)
add_rect(sl, right_x, Inches(1.35), Inches(3.8), Inches(5.8),
         fill=C_WHITE, border_color=RGBColor(0xE2,0xE8,0xF0), border_pt=1.2, radius=12000)
add_text(sl, '課程設計原則', right_x + Inches(0.2), Inches(1.45),
         Inches(3.4), Inches(0.45), font_size=15, bold=True, color=C_DARK)

principles = [
    (C_BLUE,   '螺旋序階', '由具體到抽象，由簡單到\n複雜的漸進式學習序列'),
    (C_GREEN,  '個別差異', '三種難度 × 輔助點擊，\n確保不同能力學生成功'),
    (C_ORANGE, '情境類化', '最終目標：社區真實場景\n獨立完成消費交易'),
    (C_PURPLE, 'IEP 對應', '各單元提供評量基準，\n直接對應轉銜目標'),
]
for j, (c, title, desc) in enumerate(principles):
    ry2 = Inches(1.95 + j * 1.15)
    add_rect(sl, right_x + Inches(0.2), ry2, Inches(0.06), Inches(0.75), fill=c)
    add_text(sl, title, right_x + Inches(0.38), ry2,
             Inches(1.3), Inches(0.35), font_size=13, bold=True, color=c)
    add_multiline(sl, desc.split('\n'), right_x + Inches(0.38), ry2 + Inches(0.35),
                  Inches(3.15), Inches(0.7), font_size=11, color=C_GRAY)

# 底部
add_rect(sl, Inches(0.4), Inches(6.1), Inches(8.5), Inches(1.1),
         fill=C_PANEL, border_color=RGBColor(0xCB,0xD5,0xE1), border_pt=1, radius=10000)
add_text(sl, '學習路徑建議', Inches(0.6), Inches(6.15), Inches(8), Inches(0.38),
         font_size=13, bold=True, color=C_DARK)
add_text(sl, 'F1→F2→F3→F4→F5→F6  ➡  C1→C2→C3→C4→C5→C6  ➡  A1→A2→A3→A4→A5→A6',
         Inches(0.6), Inches(6.52), Inches(8), Inches(0.5),
         font_size=12, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
#  SLIDE 07  F 系列
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, 'F 系列 ─── 數學基礎先備（6 個單元）',
                   '適用：國小中高年級至國中 ／ 先備能力：能執行觸控動作', bg=C_BLUE)

f_units = [
    ('F1', '一對一對應', '拖曳物件配對\n建立數量等量直觀', '能執行觸控'),
    ('F2', '唱數',      '依序點擊計數\n語音同步播報', '能識別圖片'),
    ('F3', '數字認讀',  '看圖選數字\n支援自訂主題上傳', '能辨識 1–10'),
    ('F4', '數字排序',  '拖曳數字由小到大\n或由大到小', '能完成 F3'),
    ('F5', '量比較',    '比較兩組多寡\n建立大小概念', '能完成 F2'),
    ('F6', '數的組成',  '合成與分解概念\n理解數的結構', '能完成 F4/F5'),
]

for i, (code, name, skill, prereq) in enumerate(f_units):
    col = i % 3
    row = i // 3
    bx = Inches(0.4 + col * 4.27)
    by = Inches(1.35 + row * 2.85)
    add_rect(sl, bx, by, Inches(4.0), Inches(2.65),
             fill=C_WHITE, border_color=C_BLUE_L, border_pt=1.5, radius=12000)
    add_rect(sl, bx, by, Inches(4.0), Inches(0.55), fill=C_BLUE, radius=0)
    add_text(sl, code, bx + Inches(0.15), by + Inches(0.06),
             Inches(0.7), Inches(0.42), font_size=18, bold=True, color=C_WHITE)
    add_text(sl, name, bx + Inches(0.85), by + Inches(0.09),
             Inches(3.0), Inches(0.38), font_size=17, bold=True, color=C_WHITE)
    add_multiline(sl, skill.split('\n'), bx + Inches(0.18), by + Inches(0.65),
                  Inches(3.64), Inches(1.3), font_size=13, color=C_DARK)
    add_rect(sl, bx + Inches(0.18), by + Inches(2.05), Inches(3.64), Inches(0.4),
             fill=C_BLUE_L, radius=8000)
    add_text(sl, f'先備：{prereq}', bx + Inches(0.28), by + Inches(2.07),
             Inches(3.44), Inches(0.36), font_size=11, color=C_BLUE)

# ══════════════════════════════════════════════════════════════
#  SLIDE 08  C 系列
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, 'C 系列 ─── 貨幣認知建立（6 個單元）',
                   '適用：國中至高中 ／ 台灣 9 種面額・拖曳操作・三難度', bg=C_GREEN)

c_units = [
    ('C1', '認識錢幣', '辨識台灣 9 種面額\n看圖選字 / 看字選幣 / 聽音選幣', '能完成 F3'),
    ('C2', '數錢',    '點擊錢幣依序累加\n語音播報累計金額', '能完成 C1'),
    ('C3', '換錢',    '以不同面額組合等值\n大換小 / 小換大 / 隨機', '能完成 C2'),
    ('C4', '付款',    '拖曳正確金額完成付款\n自動判定超額/不足', '能完成 C3'),
    ('C5', '夠不夠',  '判斷錢包金額能否購買\n差額視覺化顯示', '能完成 C4'),
    ('C6', '找零',    '拖曳找零錢幣至對應格\n驗收完整交易閉環', '能完成 C5'),
]

for i, (code, name, skill, prereq) in enumerate(c_units):
    col = i % 3
    row = i // 3
    bx = Inches(0.4 + col * 4.27)
    by = Inches(1.35 + row * 2.85)
    add_rect(sl, bx, by, Inches(4.0), Inches(2.65),
             fill=C_WHITE, border_color=C_GREEN_L, border_pt=1.5, radius=12000)
    add_rect(sl, bx, by, Inches(4.0), Inches(0.55), fill=C_GREEN, radius=0)
    add_text(sl, code, bx + Inches(0.15), by + Inches(0.06),
             Inches(0.7), Inches(0.42), font_size=18, bold=True, color=C_WHITE)
    add_text(sl, name, bx + Inches(0.85), by + Inches(0.09),
             Inches(3.0), Inches(0.38), font_size=17, bold=True, color=C_WHITE)
    add_multiline(sl, skill.split('\n'), bx + Inches(0.18), by + Inches(0.65),
                  Inches(3.64), Inches(1.3), font_size=13, color=C_DARK)
    add_rect(sl, bx + Inches(0.18), by + Inches(2.05), Inches(3.64), Inches(0.4),
             fill=C_GREEN_L, radius=8000)
    add_text(sl, f'先備：{prereq}', bx + Inches(0.28), by + Inches(2.07),
             Inches(3.44), Inches(0.36), font_size=11, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
#  SLIDE 09  A 系列（工作分析法）
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, 'A 系列 ─── 真實情境類化（6 個單元）',
                   '工作分析法拆解 5～8 個操作步驟 ／ 台灣在地化仿真介面', bg=C_ORANGE)

a_units = [
    ('A1', '自動販賣機', 5, '選商品→投幣→確認→取貨→核對找零', '幣值累計・找零驗收'),
    ('A2', '理髮廳自助機', 6, '選類別→選服務→確認→付款→取號→等候', '服務辨識・等候規範'),
    ('A3', '速食店點餐機', 7, '瀏覽→加購物車→確認→選付款→付款→等號→取餐', '分類導覽・多品項加總'),
    ('A4', '超市購物', 6, '選商品→確認價格→選付款→付款→找零→完成', '購物清單執行'),
    ('A5', 'ATM 操作', 7, '插卡→密碼→選功能→輸金額→取款→退卡→取收據', '密碼保護・金融流程'),
    ('A6', '火車票購票機', 8, '選出發站→目的地→票種→確認票價→付款→取票→找零→完成', '路線規劃・多步驟整合'),
]

for i, (code, name, steps, seq, skill) in enumerate(a_units):
    by2 = Inches(1.35 + i * 0.98)
    add_rect(sl, Inches(0.35), by2, Inches(12.6), Inches(0.88),
             fill=C_WHITE, border_color=C_ORANGE_L, border_pt=1.2, radius=8000)
    # 代碼標籤
    add_rect(sl, Inches(0.35), by2, Inches(0.88), Inches(0.88), fill=C_ORANGE, radius=8000)
    add_text(sl, code, Inches(0.35), by2, Inches(0.88), Inches(0.88),
             font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # 名稱
    add_text(sl, name, Inches(1.32), by2 + Inches(0.05),
             Inches(2.1), Inches(0.42), font_size=15, bold=True, color=C_ORANGE)
    # 步數標籤
    add_pill(sl, f'{steps}步', Inches(1.32), by2 + Inches(0.5),
             Inches(0.75), Inches(0.28), bg=C_ORANGE_L, fg=C_ORANGE, font_size=11)
    # 核心技能
    add_text(sl, f'★ {skill}', Inches(1.32), by2 + Inches(0.5),
             Inches(2.1), Inches(0.28), font_size=10, color=C_GRAY)
    # 步驟序列
    add_text(sl, seq, Inches(3.6), by2 + Inches(0.17),
             Inches(9.2), Inches(0.55), font_size=11.5, color=C_DARK)

# ══════════════════════════════════════════════════════════════
#  SLIDE 10  特色功能
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '核心特色功能', '四大創新設計 × 全 18 單元整合', bg=C_PURPLE)

features = [
    (C_PURPLE, C_PURPLE_L, '🖱️',
     '輔助點擊模式（首創）',
     '精細動作障礙學生\n任意位置點擊 → 自動執行正確步驟\n全 18 單元支援\n「位置無關」點擊推進流程'),
    (C_BLUE, C_BLUE_L, '🔊',
     '全程語音引導',
     'Web Speech API\nMicrosoft 雅婷優先（Windows）\n台灣語境金額唸法\n即時讚美音效 + 語音回饋'),
    (C_TEAL, C_TEAL_L, '🌙',
     '深色模式',
     '全 18 單元支援\n適合視覺敏感學生\n一鍵切換，無需重載\n自動偵測系統偏好'),
    (C_AMBER, C_AMBER_L, '📝',
     '配套作業單系統',
     '18 單元 × 10 種題型\n一鍵列印 / PDF 下載\n看圖填空・找零計算\n完整對應 IEP 評量記錄'),
    (C_GREEN, C_GREEN_L, '🏆',
     '獎勵積分系統',
     '獨立模組（reward/）\n學生照片 + 積分管理\n整合 IEP 行為正增強\n教師課堂即時操作'),
    (C_ORANGE, C_ORANGE_L, '📱',
     '跨裝置支援',
     'Windows 電腦（推薦）\niPad / Android 平板\n觸控拖曳 × 滑鼠點擊\n完全離線 · 授權費 0 元'),
]

for i, (bc, bl, icon, title, body) in enumerate(features):
    col = i % 3
    row = i // 3
    bx = Inches(0.35 + col * 4.3)
    by3 = Inches(1.35 + row * 2.85)
    add_rect(sl, bx, by3, Inches(4.0), Inches(2.65),
             fill=bl, border_color=bc, border_pt=1.8, radius=15000)
    add_text(sl, icon, bx, by3 + Inches(0.12), Inches(4.0), Inches(0.65),
             font_size=30, align=PP_ALIGN.CENTER)
    add_text(sl, title, bx, by3 + Inches(0.75), Inches(4.0), Inches(0.45),
             font_size=15, bold=True, color=bc, align=PP_ALIGN.CENTER)
    horizontal_divider(sl, by3 + Inches(1.22), bx + Inches(0.3), bx + Inches(3.7), color=bc, pt=0.8)
    add_multiline(sl, body.split('\n'), bx + Inches(0.25), by3 + Inches(1.3),
                  Inches(3.5), Inches(1.2), font_size=12, color=C_DARK)

# ══════════════════════════════════════════════════════════════
#  SLIDE 11  三段式學習歷程
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '三段式學習歷程設計', '數位練習 → 紙本評量 → 社區類化', bg=C_TEAL)

stages = [
    (C_BLUE, C_BLUE_L, '第一段', '數位模擬練習',
     '課堂主軸',
     ['軟體情境操作（18 個互動單元）',
      '即時觀察錯誤類型',
      '三種難度對應個別需求',
      '輔助點擊支援精細動作障礙',
      '調整教學策略依據']),
    (C_GREEN, C_GREEN_L, '第二段', '紙本形成性評量',
     '課後作業',
     ['列印對應作業單',
      '各單元備有 10 種題型',
      'PDF 下載 / 直接列印',
      '評估數位學習是否遷移',
      '至紙筆形式（IEP 記錄）']),
    (C_ORANGE, C_ORANGE_L, '第三段', '實地類化應用',
     '社區轉銜',
     ['真實場景操作演練',
      '行為檢核表記錄達成程度',
      '教師逐步褪除支持',
      '最終目標：教師 3 公尺外觀察',
      'IEP 轉銜評估依據']),
]

for i, (bc, bl, num, title, sub, items) in enumerate(stages):
    bx = Inches(0.5 + i * 4.2)
    add_rect(sl, bx, Inches(1.3), Inches(3.9), Inches(5.85),
             fill=bl, border_color=bc, border_pt=2, radius=15000)
    # 序號圓形
    add_rect(sl, bx + Inches(1.45), Inches(1.35), Inches(1.0), Inches(1.0), fill=bc, radius=50000)
    add_text(sl, num, bx + Inches(1.45), Inches(1.35), Inches(1.0), Inches(1.0),
             font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, bx, Inches(2.45), Inches(3.9), Inches(0.55),
             font_size=20, bold=True, color=bc, align=PP_ALIGN.CENTER)
    add_text(sl, sub, bx, Inches(2.95), Inches(3.9), Inches(0.38),
             font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
    horizontal_divider(sl, Inches(3.38), bx + Inches(0.3), bx + Inches(3.6), color=bc, pt=1)
    for j, item in enumerate(items):
        iy = Inches(3.5 + j * 0.52)
        add_rect(sl, bx + Inches(0.25), iy + Inches(0.09), Inches(0.12), Inches(0.28), fill=bc)
        add_text(sl, item, bx + Inches(0.48), iy, Inches(3.25), Inches(0.48),
                 font_size=12.5, color=C_DARK)

# 箭頭連接
for ax in [Inches(4.55), Inches(8.75)]:
    add_text(sl, '➡', ax, Inches(4.1), Inches(0.6), Inches(0.6),
             font_size=32, color=C_GRAY, align=PP_ALIGN.CENTER)

# 底部
add_rect(sl, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
         fill=C_TEAL_L, radius=8000)
add_text(sl, '三段設計確保學習由「數位情境」遷移至「真實社區」，完整支援 IEP 轉銜目標的評量、介入與成效記錄',
         Inches(0.6), Inches(7.12), Inches(12.0), Inches(0.28),
         font_size=12, color=C_TEAL, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  SLIDE 12  課程法規依據
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '課程法規依據與 IEP 對應', bg=C_DARK)

laws = [
    (C_RED,    '🏛️', '特殊教育法 第 28 條',
     '應為身心障礙學生訂定個別化教育計畫（IEP）',
     '本軟體各單元提供\nIEP 年度目標範本及評量基準'),
    (C_ORANGE, '📋', '特殊需求領域課綱\n生活管理・財務管理',
     '生管-Ⅳ-E-1 認識貨幣\n生管-Ⅳ-E-2 選取適當幣值完成付款\n生管-Ⅳ-E-3 驗收找零\n生管-Ⅳ-F-1 使用自動化設備',
     'F/C/A 三系列 18 單元\n完整對應財務管理學習序階'),
    (C_BLUE,   '🌏', '身心障礙者權利公約\n第 9・19 條',
     '無障礙近用設施\n獨立生活於社區',
     '輔助點擊消除動作障礙\nA 系列六種真實場景\n服務社區轉銜'),
    (C_GREEN,  '📚', '十二年國教課綱總綱',
     '自主行動・社會參與核心素養',
     'A 系列六種真實場景\nF 系列數概念建構'),
    (C_PURPLE, '⚖️', '高中以下特教課程\n教材教法及評量辦法 第 6 條',
     '教材應符合學生個別學習需求',
     '三種難度 × 輔助點擊\n支援個別化差異教學'),
    (C_AMBER,  '📖', '十二年國教\n特殊需求領域課綱（2019）',
     '財務管理主題\n（頁 45–52）',
     '系統性對應課綱架構\n提供教學目標與評量規準'),
]

for i, (bc, icon, law_title, law_text, response) in enumerate(laws):
    col = i % 2
    row = i // 2
    bx = Inches(0.35 + col * 6.4)
    by4 = Inches(1.35 + row * 2.0)
    add_rect(sl, bx, by4, Inches(6.15), Inches(1.85),
             fill=C_WHITE, border_color=RGBColor(0xE2,0xE8,0xF0), border_pt=1.2, radius=10000)
    add_rect(sl, bx, by4, Inches(0.12), Inches(1.85), fill=bc, radius=0)
    add_text(sl, icon, bx + Inches(0.2), by4 + Inches(0.12),
             Inches(0.55), Inches(0.55), font_size=22)
    add_text(sl, law_title, bx + Inches(0.8), by4 + Inches(0.08),
             Inches(5.1), Inches(0.52), font_size=13, bold=True, color=bc)
    add_multiline(sl, law_text.split('\n'), bx + Inches(0.8), by4 + Inches(0.6),
                  Inches(3.1), Inches(1.1), font_size=11, color=C_GRAY)
    add_rect(sl, bx + Inches(3.95), by4 + Inches(0.4), Inches(2.05), Inches(1.3),
             fill=RGBColor(0xF8,0xFA,0xFC), border_color=bc, border_pt=0.8, radius=8000)
    add_text(sl, '▶ 本軟體對應', bx + Inches(4.0), by4 + Inches(0.42),
             Inches(1.9), Inches(0.32), font_size=10, bold=True, color=bc)
    add_multiline(sl, response.split('\n'), bx + Inches(4.0), by4 + Inches(0.72),
                  Inches(1.9), Inches(0.9), font_size=11, color=C_DARK)

# ══════════════════════════════════════════════════════════════
#  SLIDE 13  試教觀察與預期成效
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '試教觀察與預期成效', '類化窗口的關鍵時刻 × 三階段成效指標', bg=C_GREEN)

# 試教觀察（左）
add_rect(sl, Inches(0.35), Inches(1.3), Inches(5.5), Inches(5.85),
         fill=C_WHITE, border_color=C_GREEN, border_pt=1.5, radius=12000)
add_rect(sl, Inches(0.35), Inches(1.3), Inches(5.5), Inches(0.52),
         fill=C_GREEN, radius=0)
add_text(sl, '🔬  試教觀察：類化窗口', Inches(0.5), Inches(1.34),
         Inches(5.2), Inches(0.44), font_size=16, bold=True, color=C_WHITE)

observations = [
    '🎯 共同轉折點',
    '學生在第 5～7 次軟體操作後，開始在啟動前\n主動低聲說出下一步名稱',
    '→ 程序性記憶語言化的信號',
    '→ 類化窗口即將開啟的前兆',
    '',
    '✅ A2 自閉症學生',
    '社區實地操作時，面對真實理髮廳觸控機：\n完成完整流程 + 主動說出等待規則\n「我是 7 號，還沒叫到我，我要繼續等」',
    '',
    '✅ A1 個案（小雯）',
    '完成 15 次軟體練習後：\n第一次在真實販賣機前完整購買\n主動說：「我下次自己來。」',
]
obs_y = Inches(1.9)
for obs in observations:
    if obs == '':
        obs_y += Inches(0.12)
        continue
    fy = 11 if not obs.startswith('🎯') and not obs.startswith('✅') else 13
    bold2 = obs.startswith('🎯') or obs.startswith('✅')
    col2 = C_GREEN if bold2 else C_DARK
    add_text(sl, obs, Inches(0.55), obs_y, Inches(5.1), Inches(0.55),
             font_size=fy, bold=bold2, color=col2)
    obs_y += Inches(0.45 if bold2 else 0.55)

# 預期成效（右）
add_rect(sl, Inches(6.15), Inches(1.3), Inches(6.85), Inches(5.85),
         fill=C_WHITE, border_color=C_BLUE, border_pt=1.5, radius=12000)
add_rect(sl, Inches(6.15), Inches(1.3), Inches(6.85), Inches(0.52),
         fill=C_BLUE, radius=0)
add_text(sl, '📊  預期成效指標', Inches(6.3), Inches(1.34),
         Inches(6.5), Inches(0.44), font_size=16, bold=True, color=C_WHITE)

outcomes = [
    (C_GREEN, C_GREEN_L, '短期（1～2 個月）',
     '認知：9 種面額正確辨識 ≥80%\n技能：A1～A3 普通模式完成\n情意：面對自助設備焦慮降低'),
    (C_AMBER, C_AMBER_L, '中期（3～6 個月）',
     '認知：找零計算正確率 ≥70%\n技能：A4～A6 普通模式；ATM 獨立操作\n情意：不再迴避社區場景'),
    (C_ORANGE, C_ORANGE_L, '長期（6 個月以上）',
     '認知：能向他人說明消費流程\n技能：真實場景（理髮廳/超市/ATM）\n情意：建立社區消費自主信心'),
]

for j, (bc, bl, period, desc) in enumerate(outcomes):
    oy = Inches(1.95 + j * 1.72)
    add_rect(sl, Inches(6.3), oy, Inches(6.5), Inches(1.55),
             fill=bl, border_color=bc, border_pt=1.2, radius=10000)
    add_text(sl, period, Inches(6.45), oy + Inches(0.1),
             Inches(6.2), Inches(0.38), font_size=13, bold=True, color=bc)
    add_multiline(sl, desc.split('\n'), Inches(6.45), oy + Inches(0.5),
                  Inches(6.2), Inches(0.9), font_size=12, color=C_DARK)

# ══════════════════════════════════════════════════════════════
#  SLIDE 14  技術規格
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_LIGHT_BG)
add_section_header(sl, '技術規格與無障礙設計', 'Technical Specifications & Accessibility', bg=C_DARK)

specs = [
    ('🖥️', '開發技術',  'HTML5 / CSS3 / JavaScript（ES6+）\n無需安裝任何外掛'),
    ('📏', '程式規模',  '18 個主程式，超過 15 萬行程式碼\n研發歷時約 4 個月'),
    ('📱', '支援裝置',  'Windows 電腦（推薦）\niPad / Android 平板（9.7 吋以上）'),
    ('🌐', '支援瀏覽器','Google Chrome / Microsoft Edge\n（最新版）'),
    ('📵', '離線使用',  '完整支援\n首次語音啟動後可完全離線'),
    ('🔊', '語音引擎',  'Web Speech API\n優先 Microsoft 雅婷（Windows）\n台灣語境金額唸法'),
    ('🖱️', '輔助點擊',  '全 18 單元支援\n精細動作障礙學生任意位置點擊\n自動執行正確步驟'),
    ('🌙', '深色模式',  '全 18 單元支援\n適合視覺敏感或視覺疲勞學生'),
    ('💰', '製作費用',  '新臺幣 0 元\n教師自主研發，無任何商業授權費用'),
]

for i, (icon, title, desc) in enumerate(specs):
    col = i % 3
    row = i // 3
    bx = Inches(0.35 + col * 4.3)
    by5 = Inches(1.35 + row * 1.95)
    add_rect(sl, bx, by5, Inches(4.0), Inches(1.78),
             fill=C_WHITE, border_color=RGBColor(0xE2,0xE8,0xF0), border_pt=1.2, radius=12000)
    add_text(sl, icon, bx + Inches(0.15), by5 + Inches(0.12),
             Inches(0.65), Inches(0.65), font_size=26)
    add_text(sl, title, bx + Inches(0.85), by5 + Inches(0.12),
             Inches(2.9), Inches(0.4), font_size=14, bold=True, color=C_DARK)
    add_multiline(sl, desc.split('\n'), bx + Inches(0.85), by5 + Inches(0.55),
                  Inches(2.9), Inches(1.1), font_size=12, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
#  SLIDE 15  總結
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=RGBColor(0x0F, 0x17, 0x2A))
add_rect(sl, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=C_GREEN)

# 標題
add_text(sl, '核心創新摘要', Inches(0.55), Inches(0.5),
         Inches(9), Inches(0.65),
         font_size=32, bold=True, color=C_WHITE, font_name='微軟正黑體')
add_text(sl, 'Key Innovations Summary',
         Inches(0.55), Inches(1.1), Inches(9), Inches(0.4),
         font_size=15, color=RGBColor(0x94,0xA3,0xB8), font_name='Arial', italic=True)

innovations = [
    (C_BLUE,   '①', '全台首套三層次序階\n數位金融素養課程',
     'F→C→A 18 個互動教學單元，系統覆蓋\n「數學先備→貨幣認知→情境類化」完整序列'),
    (C_GREEN,  '②', '輔助點擊模式\n（特教數位教材首創）',
     '精細動作障礙學生任意位置點擊\n即自動執行正確步驟，全 18 單元支援'),
    (C_ORANGE, '③', '六種台灣在地化\n真實情境仿真介面',
     '自動販賣機・理髮廳・速食店\n超市・ATM・台鐵售票機'),
    (C_PURPLE, '④', '完整配套系統\n對應 IEP 各功能向度',
     '作業單（18 × 10 題型）\n獎勵積分・深色模式・語音引導'),
]

for i, (bc, num, title, desc) in enumerate(innovations):
    col = i % 2
    row = i // 2
    bx = Inches(0.55 + col * 6.3)
    by6 = Inches(1.75 + row * 2.5)
    add_rect(sl, bx, by6, Inches(5.9), Inches(2.2),
             fill=RGBColor(0x1E, 0x3A, 0x5C), border_color=bc, border_pt=1.5, radius=12000)
    add_rect(sl, bx, by6, Inches(0.6), Inches(2.2), fill=bc, radius=0)
    add_text(sl, num, bx, by6, Inches(0.6), Inches(2.2),
             font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, bx + Inches(0.7), by6 + Inches(0.15),
             Inches(5.0), Inches(0.75), font_size=16, bold=True, color=bc, font_name='微軟正黑體')
    add_multiline(sl, desc.split('\n'), bx + Inches(0.7), by6 + Inches(0.95),
                  Inches(5.0), Inches(1.1), font_size=13, color=RGBColor(0xCB,0xD5,0xE1))

# 底部
add_rect(sl, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.45),
         fill=RGBColor(0x1E, 0x3A, 0x5C), radius=10000)
add_text(sl, '身心障礙類・國民中學及高級中等學校組　｜　製作費用 NT$0　｜　完全離線・免安裝　｜　中華民國 115 年 3 月',
         Inches(0.55), Inches(6.87), Inches(12.2), Inches(0.4),
         font_size=12, color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  儲存
# ══════════════════════════════════════════════════════════════
import sys
sys.stdout.reconfigure(encoding='utf-8')
out_path = 'doc/金錢小達人_競賽申請簡報.pptx'
prs.save(out_path)
print(f'OK saved: {out_path}  slides={len(prs.slides)}')
