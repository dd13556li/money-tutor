# -*- coding: utf-8 -*-
"""
金錢小達人 競賽申請 PPT V3
修正：S07 金字塔版面；新增 S19 B系列、S20 IEP目標、S21 三難度比較、S22 結語
共 22 張投影片
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

W   = RGBColor(0xFF,0xFF,0xFF)
BK  = RGBColor(0x0F,0x17,0x2A)
DK  = RGBColor(0x1E,0x29,0x3B)
GR  = RGBColor(0x64,0x74,0x8B)
BG  = RGBColor(0xF8,0xFA,0xFC)
PNL = RGBColor(0xF1,0xF5,0xF9)
PNL2= RGBColor(0xE2,0xE8,0xF0)
FB  = RGBColor(0x22,0x63,0xEB)
FBL = RGBColor(0xDB,0xEA,0xFE)
FBD = RGBColor(0x1D,0x4E,0xD8)
CG  = RGBColor(0x05,0x96,0x69)
CGL = RGBColor(0xD1,0xFA,0xE5)
CGD = RGBColor(0x04,0x7A,0x55)
AO  = RGBColor(0xEA,0x58,0x0C)
AOL = RGBColor(0xFF,0xED,0xD5)
AOD = RGBColor(0xC2,0x41,0x0C)
PU  = RGBColor(0x7C,0x3A,0xED)
PUL = RGBColor(0xED,0xE9,0xFE)
AM  = RGBColor(0xD9,0x77,0x06)
AML = RGBColor(0xFE,0xF3,0xC7)
RD  = RGBColor(0xDC,0x26,0x26)
RDL = RGBColor(0xFE,0xE2,0xE2)
TE  = RGBColor(0x0D,0x94,0x88)
TEL = RGBColor(0xCC,0xFB,0xF1)
IND = RGBColor(0x44,0x38,0xCA)
INDL= RGBColor(0xE0,0xE7,0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

def rect(sl, x, y, w, h, fill=None, bc=None, bpt=0, rad=None):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if bc and bpt > 0:
        sh.line.color.rgb = bc; sh.line.width = Pt(bpt)
    if rad is not None:
        sp = sh._element; spPr = sp.find(qn('p:spPr'))
        pg = spPr.find(qn('a:prstGeom'))
        if pg is not None:
            pg.set('prst','roundRect')
            av = pg.find(qn('a:avLst'))
            if av is None: av = etree.SubElement(pg, qn('a:avLst'))
            gd = etree.SubElement(av, qn('a:gd'))
            gd.set('name','adj'); gd.set('fmla',f'val {rad}')
    return sh

def ellipse(sl, x, y, w, h, fill=None, bc=None, bpt=0):
    sh = sl.shapes.add_shape(9, x, y, w, h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if bc and bpt>0: sh.line.color.rgb=bc; sh.line.width=Pt(bpt)
    return sh

def txt(sl, s, x, y, w, h, sz=14, bold=False, col=None,
        align=PP_ALIGN.LEFT, fn='Microsoft JhengHei', italic=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.name = fn; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    if col: r.font.color.rgb = col
    return tb

def mtxt(sl, lines, x, y, w, h, sz=13, col=None,
         align=PP_ALIGN.LEFT, fn='Microsoft JhengHei', bold=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = fn; r.font.size = Pt(sz)
        r.font.bold = bold
        if col: r.font.color.rgb = col
    return tb

def hline(sl, y, x0=Inches(0.5), x1=Inches(12.83), col=PNL2, pt=1.2):
    ln = sl.shapes.add_connector(1, x0, y, x1, y)
    ln.line.color.rgb = col; ln.line.width = Pt(pt)

def vline(sl, x, y0, y1, col=PNL2, pt=1.2):
    ln = sl.shapes.add_connector(1, x, y0, x, y1)
    ln.line.color.rgb = col; ln.line.width = Pt(pt)

def pill(sl, s, x, y, w, h, bg, fg, sz=12, bold=True):
    rect(sl, x, y, w, h, fill=bg, rad=20000)
    txt(sl, s, x, y, w, h, sz=sz, bold=bold, col=fg, align=PP_ALIGN.CENTER)

def dot_badge(sl, s, cx, cy, r_inch, bg, fg, sz=16):
    ellipse(sl, cx-r_inch, cy-r_inch, r_inch*2, r_inch*2, fill=bg)
    txt(sl, s, cx-r_inch, cy-r_inch*0.75, r_inch*2, r_inch*1.5,
        sz=sz, bold=True, col=fg, align=PP_ALIGN.CENTER)

def page_num(sl, n, total=22):
    txt(sl, f'{n} / {total}',
        SLIDE_W - Inches(1.1), SLIDE_H - Inches(0.38),
        Inches(0.9), Inches(0.32),
        sz=10, col=GR, align=PP_ALIGN.RIGHT)

def top_bar(sl, title, subtitle='', bg=DK, h=Inches(1.05)):
    rect(sl, Inches(0), Inches(0), SLIDE_W, h, fill=bg)
    rect(sl, Inches(0), Inches(0), Inches(0.06), h, fill=W)
    txt(sl, title, Inches(0.35), Inches(0.08), SLIDE_W-Inches(0.7), Inches(0.58),
        sz=24, bold=True, col=W)
    if subtitle:
        txt(sl, subtitle, Inches(0.35), Inches(0.65), SLIDE_W-Inches(0.7), Inches(0.34),
            sz=12.5, col=RGBColor(0xBA,0xD8,0xFF), italic=True)

def stat_card(sl, val, label, x, y, w=Inches(2.0), h=Inches(1.05),
              val_c=AM, bg=RGBColor(0x1E,0x3A,0x5C), lbl_c=RGBColor(0x94,0xA3,0xB8)):
    rect(sl, x, y, w, h, fill=bg, rad=12000)
    txt(sl, val, x, y+Inches(0.08), w, Inches(0.52),
        sz=26, bold=True, col=val_c, align=PP_ALIGN.CENTER)
    txt(sl, label, x, y+Inches(0.58), w, Inches(0.4),
        sz=11, col=lbl_c, align=PP_ALIGN.CENTER)

def info_card(sl, icon, title, body_lines, x, y, w, h,
              bg=FBL, bord=FB, tc=FB, ic_sz=28):
    rect(sl, x, y, w, h, fill=bg, bc=bord, bpt=1.5, rad=14000)
    txt(sl, icon, x+Inches(0.15), y+Inches(0.1),
        Inches(0.65), h-Inches(0.2), sz=ic_sz, align=PP_ALIGN.CENTER)
    txt(sl, title, x+Inches(0.88), y+Inches(0.1),
        w-Inches(1.0), Inches(0.42), sz=13, bold=True, col=tc)
    mtxt(sl, body_lines, x+Inches(0.88), y+Inches(0.5),
         w-Inches(1.0), h-Inches(0.6), sz=12, col=DK)

def step_circle(sl, n, x, y, r=Inches(0.28), bg=FB, fg=W):
    ellipse(sl, x-r, y-r, r*2, r*2, fill=bg)
    txt(sl, str(n), x-r, y-r*0.8, r*2, r*1.6,
        sz=12, bold=True, col=fg, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# S01  封面
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BK)
rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill=AO)
for rx,ry,rw,rh,c in [
    (Inches(10.2), Inches(-0.6), Inches(3.8), Inches(3.8), RGBColor(0x1E,0x40,0x7A)),
    (Inches(11.2), Inches(4.6),  Inches(2.8), Inches(2.8), RGBColor(0x06,0x40,0x38)),
    (Inches(9.1),  Inches(3.2),  Inches(1.8), Inches(1.8), RGBColor(0x30,0x26,0x80)),
]:
    sh = rect(sl, rx, ry, rw, rh, fill=c)
    sp = sh._element; spPr=sp.find(qn('p:spPr'))
    pg=spPr.find(qn('a:prstGeom'))
    if pg is not None: pg.set('prst','ellipse')
pill(sl, '第八屆教育部獎助研發特殊教育教材教具競賽  ·  身心障礙類  ·  國民中學及高級中等學校組',
     Inches(0.55), Inches(0.55), Inches(8.2), Inches(0.44),
     bg=RGBColor(0x7C,0x2D,0x12), fg=RGBColor(0xFF,0xC5,0x87), sz=12)
txt(sl, '金 錢 小 達 人', Inches(0.55), Inches(1.18), Inches(9.5), Inches(1.35),
    sz=62, bold=True, col=W)
txt(sl, '互動式特殊教育生活數學暨金融素養教學軟體',
    Inches(0.55), Inches(2.52), Inches(9.5), Inches(0.55),
    sz=22, col=RGBColor(0xFD,0xBA,0x74), bold=True)
txt(sl, 'Money Tutor: Interactive Software for Financial Literacy in Special Education',
    Inches(0.55), Inches(3.1), Inches(9.5), Inches(0.42),
    sz=13, col=RGBColor(0x7C,0x8D,0xA8), fn='Arial', italic=True)
hline(sl, Inches(3.65), Inches(0.55), Inches(9.2),
      col=RGBColor(0x33,0x4A,0x6E), pt=1.5)
stat_card(sl, '18', '互動教學單元', Inches(0.55), Inches(3.85))
stat_card(sl, '150,000+', '程式碼行數',  Inches(2.7),  Inches(3.85))
stat_card(sl, '4 個月', '研發歷時',    Inches(4.85), Inches(3.85))
stat_card(sl, 'NT$ 0', '授權費用',    Inches(7.0),  Inches(3.85))
rect(sl, Inches(0.55), Inches(5.2), Inches(8.7), Inches(0.06),
     fill=RGBColor(0x33,0x4A,0x6E))
txt(sl, '適用對象：國中～高中特教班・資源班',
    Inches(0.55), Inches(5.35), Inches(8.7), Inches(0.4),
    sz=13, col=RGBColor(0xCB,0xD5,0xE1))
txt(sl, '輕度至中度智能障礙・自閉症・學習障礙・多重障礙',
    Inches(0.55), Inches(5.72), Inches(8.7), Inches(0.38),
    sz=12, col=RGBColor(0x94,0xA3,0xB8))
for i,(c,t) in enumerate([(FB,'F 系列  數學基礎先備'),(CG,'C 系列  貨幣認知建立'),(AO,'A 系列  真實情境類化')]):
    sx = Inches(0.55+i*2.95)
    rect(sl, sx, Inches(6.25), Inches(0.2), Inches(0.38), fill=c, rad=6000)
    txt(sl, t, sx+Inches(0.28), Inches(6.25), Inches(2.55), Inches(0.38),
        sz=12, col=RGBColor(0xCB,0xD5,0xE1))
txt(sl, '中華民國 115 年 3 月',
    Inches(0.55), Inches(6.85), Inches(4), Inches(0.38),
    sz=12, col=RGBColor(0x47,0x5A,0x73))
page_num(sl, 1)

# ═══════════════════════════════════════════════════════════════
# S02  目錄（V3 更新版，15 項 3×5）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '簡報綱要', 'Table of Contents', bg=DK)

toc = [
    (RD,  '01','設計動機',    '特教現場三重結構性困境'),
    (AO,  '02','設計理念',    'UDL × 工作分析法 × 直接教學'),
    (FB,  '03','F 系列',      '數學基礎先備 6 單元'),
    (CG,  '04','C 系列',      '貨幣認知建立 6 單元'),
    (AO,  '05','A 系列',      '真實情境類化 6 單元'),
    (PU,  '06','輔助點擊',    '精細動作障礙首創解方'),
    (TE,  '07','三段學習',    '數位→紙本評量→社區類化'),
    (AM,  '08','課程法規',    'IEP / 特教法 / 課綱對應'),
    (IND, '09','試教觀察',    '個案轉化 · 類化窗口'),
    (CG,  '10','預期成效',    '三層面 × 短中長期指標'),
    (AM,  '11','B 系列延伸',  '預算規劃 6 個單元介紹'),
    (FB,  '12','IEP 目標',    '各系列年度目標範本'),
    (PU,  '13','三難度比較',  '簡單/普通/困難/輔助點擊'),
    (GR,  '14','技術規格',    '無障礙設計 · 配套系統'),
    (DK,  '15','核心創新',    '四大首創特色 · 結語'),
]
cols = 3
for idx,(c,num,t,s) in enumerate(toc):
    col = idx % cols
    row = idx // cols
    bx = Inches(0.32 + col*4.35)
    by = Inches(1.18 + row*1.22)
    rect(sl, bx, by, Inches(4.12), Inches(1.1),
         fill=W, bc=PNL2, bpt=1.2, rad=10000)
    rect(sl, bx, by, Inches(4.12), Inches(0.08), fill=c, rad=0)
    txt(sl, num, bx+Inches(0.12), by+Inches(0.12),
        Inches(0.62), Inches(0.48), sz=22, bold=True, col=c, fn='Arial')
    txt(sl, t, bx+Inches(0.82), by+Inches(0.12),
        Inches(3.15), Inches(0.38), sz=14, bold=True, col=DK)
    txt(sl, s, bx+Inches(0.82), by+Inches(0.54),
        Inches(3.15), Inches(0.46), sz=10.5, col=GR)
page_num(sl, 2)

# ═══════════════════════════════════════════════════════════════
# S03  設計動機 — 三重困境
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '01  設計動機 ── 特教現場的三重結構性困境', bg=RD)

cards3 = [
    (RD, RDL, '🧠', '困境一\n交易序列的意義理解障礙',
     ['台灣 9 種面額貨幣──',
      '挑戰不在「辨識」，在「序列理解」',
      '找零 = 付出 → 接受 → 驗收 邏輯鏈',
      '夠不夠 = 幣值辨識 + 數量比較同步',
      '靜態紙本無法提供即時個別化回饋',
      '錯誤模式在缺乏回饋中持續固化']),
    (AM, AML, '⚠️', '困境二\n安全練習場域的長期缺位',
     ['真實試誤 → 消費糾紛 + 情緒失調風險',
      '【小俊・速食店】陷入選擇迴路',
      '→ 後方催促誘發情緒失調',
      '→ 一學期拒絕接觸相關學習',
      '【小雯・販賣機】反覆失敗後',
      '→ 主動繞道迴避所有自助設備']),
    (PU, PUL, '🧩', '困境三\n個別差異無法兼顧的教材斷層',
     ['市售教材：一難度・一題型・無回饋',
      '認知差距 3～5 年的班級共用同份',
      '精細動作障礙學生無法操作觸控',
      '教師自製：耗時且難以系統性',
      '每次上課重複製作 → 教師疲憊',
      '缺乏系統化序階課程整合']),
]
for i,(bc,bl,icon,title,body) in enumerate(cards3):
    bx = Inches(0.35+i*4.32)
    rect(sl, bx, Inches(1.12), Inches(4.12), Inches(6.05),
         fill=bl, bc=bc, bpt=2, rad=16000)
    txt(sl, icon, bx, Inches(1.18), Inches(4.12), Inches(0.75),
        sz=36, align=PP_ALIGN.CENTER)
    mtxt(sl, title.split('\n'), bx+Inches(0.2), Inches(1.95),
         Inches(3.72), Inches(0.72), sz=15, bold=True, col=bc)
    hline(sl, Inches(2.72), bx+Inches(0.2), bx+Inches(3.92), col=bc, pt=1)
    for k,line in enumerate(body):
        c3 = bc if line.startswith('→') else DK
        txt(sl, line, bx+Inches(0.22), Inches(2.82+k*0.48),
            Inches(3.68), Inches(0.45), sz=12, col=c3,
            italic=line.startswith('→'))
page_num(sl, 3)

# ═══════════════════════════════════════════════════════════════
# S04  設計理念 — UDL
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '02  設計理念 ── 三大教育理論架構',
        'UDL 通用學習設計 × 工作分析法 × 直接教學法', bg=DK)

theories4 = [
    (PU, PUL, 'UDL\n通用學習設計',
     ['多元表徵',
      '多元行動與表達',
      '多元參與方式'],
     ['三種難度對應認知差異',
      '語音+視覺+觸控多管道',
      '輔助點擊降低動作門檻',
      '深色模式視覺無障礙',
      '自訂圖片個人化學習']),
    (CG, CGL, '工作分析法\nTask Analysis',
     ['複雜行為分解',
      '序列步驟教學',
      '系統化訓練'],
     ['A 系列 5～8 步驟拆解',
      '每步驟語音引導',
      '步驟高亮視覺提示',
      '輔助點擊步驟自動推進',
      '完整流程記錄回饋']),
    (AO, AOL, '直接教學法\nDirect Instruction',
     ['明確示範',
      '引導練習',
      '獨立操作'],
     ['即時正確/錯誤回饋',
      '錯誤三次自動提示',
      '語音讚美正向增強',
      '漸進撤除支持',
      'IEP 目標導向評量']),
]
for i,(bc,bl,title,pts1,pts2) in enumerate(theories4):
    bx = Inches(0.35+i*4.32)
    rect(sl, bx, Inches(1.12), Inches(4.12), Inches(5.95),
         fill=W, bc=bc, bpt=2, rad=16000)
    rect(sl, bx, Inches(1.12), Inches(4.12), Inches(0.58), fill=bc, rad=0)
    mtxt(sl, title.split('\n'), bx, Inches(1.15), Inches(4.12), Inches(0.52),
         sz=17, bold=True, col=W, align=PP_ALIGN.CENTER)
    txt(sl, '理論核心', bx+Inches(0.22), Inches(1.78),
        Inches(3.68), Inches(0.35), sz=12, bold=True, col=bc)
    for k,p in enumerate(pts1):
        pill(sl, p, bx+Inches(0.22), Inches(2.15+k*0.42),
             Inches(3.68), Inches(0.35), bl, bc, sz=11, bold=False)
    hline(sl, Inches(3.45), bx+Inches(0.22), bx+Inches(3.9), col=PNL2)
    txt(sl, '軟體落實', bx+Inches(0.22), Inches(3.55),
        Inches(3.68), Inches(0.35), sz=12, bold=True, col=bc)
    for k,p in enumerate(pts2):
        txt(sl, f'▸ {p}', bx+Inches(0.22), Inches(3.92+k*0.4),
            Inches(3.68), Inches(0.38), sz=11.5, col=DK)
page_num(sl, 4)

# ═══════════════════════════════════════════════════════════════
# S05  設計理念 — 多元評量與個別差異
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '02  設計理念 ── 個別差異因應矩陣',
        '三種難度 × 輔助點擊 × 多元評量模式', bg=DK)

rect(sl, Inches(0.35), Inches(1.15), Inches(12.6), Inches(5.95), fill=W, bc=PNL2, bpt=1, rad=12000)
headers = ['', '簡單模式\nEasy', '普通模式\nNormal', '困難模式\nHard']
h_colors = [DK, CG, AM, RD]
for j,hdr in enumerate(headers):
    cx5 = Inches(0.35+j*3.18)
    rect(sl, cx5, Inches(1.15), Inches(3.18), Inches(0.72), fill=h_colors[j], rad=0)
    mtxt(sl, hdr.split('\n'), cx5, Inches(1.18), Inches(3.18), Inches(0.66),
         sz=14, bold=True, col=W, align=PP_ALIGN.CENTER)

rows5 = [
    ('視覺提示', ['✓ 高亮目標\n✓ 步驟動畫', '✓ 步驟提示\n無動畫輔助', '✗ 無提示\n記憶策略']),
    ('語音引導', ['✓ 完整語音\n每步引導', '✓ 提示語音\n答錯才播', '✗ 僅回饋\n無引導語音']),
    ('錯誤容忍', ['無限次數\n自動提示', '三次後\n顯示提示', '三次後\n僅語音提示']),
    ('輔助點擊', ['✓ 全步驟\n自動執行', '✓ 一般模式\n皆可啟用', '✓ 困難模式\n亦可啟用']),
    ('評量目的', ['初學習得\nIEP 基準', '流暢性建立\nIEP 中期', '獨立熟練\nIEP 精熟']),
]
row_colors = [FBL, AML, TEL, PUL, CGL]
for i,(label,cells5) in enumerate(rows5):
    ry5 = Inches(1.92+i*0.92)
    rect(sl, Inches(0.35), ry5, Inches(3.18), Inches(0.85),
         fill=row_colors[i], rad=0)
    txt(sl, label, Inches(0.35), ry5, Inches(3.18), Inches(0.85),
        sz=13, bold=True, col=DK, align=PP_ALIGN.CENTER)
    for j,cell in enumerate(cells5):
        mtxt(sl, cell.split('\n'), Inches(3.53+j*3.18), ry5+Inches(0.06),
             Inches(3.0), Inches(0.78), sz=11.5, col=DK, align=PP_ALIGN.CENTER)
page_num(sl, 5)

# ═══════════════════════════════════════════════════════════════
# S06  三段學習歷程（概覽）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '三段式學習歷程 ── 概覽',
        '數位模擬 → 紙本評量 → 社區類化', bg=TE)

rect(sl, Inches(0.35), Inches(1.12), Inches(8.55), Inches(5.95),
     fill=W, bc=PNL2, bpt=1.2, rad=12000)
txt(sl, '整合設計理念', Inches(0.55), Inches(1.22),
    Inches(8.15), Inches(0.42), sz=15, bold=True, col=DK)
for k,line in enumerate([
    '▸ 數位練習提供安全試誤空間，消除真實情境焦慮',
    '▸ 紙本作業單評估學習是否從數位遷移至靜態情境',
    '▸ 社區類化為 IEP 轉銜目標提供直接評量依據',
    '▸ 三段設計讓教師清楚掌握學生在各層次的表現',
]):
    txt(sl, line, Inches(0.55), Inches(1.72+k*0.48),
        Inches(8.15), Inches(0.45), sz=13, col=DK)

rx6 = Inches(9.15)
stages6 = [
    (FB,'第一段','數位模擬練習','課堂主軸\n即時回饋'),
    (CG,'第二段','紙本形成性評量','課後作業\n評估遷移效果'),
    (AO,'第三段','實地類化應用','社區轉銜\nIEP 成效記錄'),
]
for j,(c6,num6,t6,d6) in enumerate(stages6):
    ry6b = Inches(1.85+j*1.72)
    ellipse(sl, rx6+Inches(0.3), ry6b+Inches(0.12), Inches(0.72), Inches(0.72), fill=c6)
    txt(sl, num6, rx6+Inches(0.3), ry6b+Inches(0.12), Inches(0.72), Inches(0.72),
        sz=11, bold=True, col=W, align=PP_ALIGN.CENTER)
    txt(sl, t6, rx6+Inches(1.1), ry6b+Inches(0.12),
        Inches(1.75), Inches(0.45), sz=14, bold=True, col=c6)
    mtxt(sl, d6.split('\n'), rx6+Inches(1.1), ry6b+Inches(0.58),
         Inches(1.75), Inches(0.72), sz=11, col=GR)
    if j < 2:
        txt(sl, '↓', rx6+Inches(0.54), ry6b+Inches(1.0),
            Inches(0.5), Inches(0.5), sz=18, col=GR, align=PP_ALIGN.CENTER)
page_num(sl, 6)

# ═══════════════════════════════════════════════════════════════
# S07  三層序階課程架構圖（V3 修正版）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '課程架構 ── 三層次螺旋序階',
        'F 數學基礎先備 → C 貨幣認知建立 → A 真實情境類化', bg=DK)

# ── 金字塔：A 在頂（小y＋窄），F 在底（大y＋寬）──
# 左側可用寬度 0.35 ~ 9.0（8.65 吋），右側面板 9.05 起
PYRL = Inches(0.35)   # 左邊距
PYRMAX = Inches(8.65) # 最大寬度

layers7 = [
    # (色, 深色, 主文, 副文, y, 寬)
    (AO, AOD, 'A 系列  ·  真實情境類化', '6 個單元  ·  國中至高中',       Inches(1.15), Inches(4.8)),
    (CG, CGD, 'C 系列  ·  貨幣認知建立', '6 個單元  ·  台灣 9 種面額',    Inches(2.45), Inches(6.8)),
    (FB, FBD, 'F 系列  ·  數學基礎先備', '6 個單元  ·  國小至國中',       Inches(3.75), Inches(8.65)),
]
for i,(bc,bcd,main,sub,y7,w7) in enumerate(layers7):
    cx = PYRL + (PYRMAX - w7) / 2
    rect(sl, cx, y7, w7, Inches(1.18), fill=bc, bc=bcd, bpt=1.5, rad=8000)
    txt(sl, main, cx, y7+Inches(0.1), w7, Inches(0.55),
        sz=19, bold=True, col=W, align=PP_ALIGN.CENTER)
    txt(sl, sub, cx, y7+Inches(0.65), w7, Inches(0.4),
        sz=12, col=RGBColor(0xFF,0xFF,0xE0), align=PP_ALIGN.CENTER, italic=True)
    # 左側序號圓
    ellipse(sl, PYRL+Inches(0.05), y7+Inches(0.34), Inches(0.52), Inches(0.52), fill=bcd)
    txt(sl, str(3-i), PYRL+Inches(0.05), y7+Inches(0.34), Inches(0.52), Inches(0.52),
        sz=14, bold=True, col=W, align=PP_ALIGN.CENTER)

# 連接箭頭（層間）
for ay in [Inches(2.38), Inches(3.68)]:
    txt(sl, '▼', PYRL+Inches(3.5), ay, Inches(2.0), Inches(0.3),
        sz=14, col=GR, align=PP_ALIGN.CENTER)

# ── 右側說明面板 ──
rx7 = Inches(9.08)
rect(sl, rx7, Inches(1.12), Inches(4.0), Inches(5.82),
     fill=W, bc=PNL2, bpt=1.2, rad=12000)
txt(sl, '課程設計核心原則', rx7+Inches(0.2), Inches(1.22),
    Inches(3.6), Inches(0.42), sz=15, bold=True, col=DK)
hline(sl, Inches(1.72), rx7+Inches(0.2), rx7+Inches(3.8), col=PNL2)

prins7 = [
    (FB, '螺旋序階',
     '具體→抽象，簡單→複雜\n每系列建立下一系列先備'),
    (CG, '個別差異',
     '三種難度 × 輔助點擊\n確保不同能力學生成功'),
    (AO, '情境類化',
     '最終目標：真實社區場景\n獨立完成消費交易'),
    (PU, 'IEP 整合',
     '各單元提供評量基準\n直接對應轉銜年度目標'),
]
for j,(c7,t7,d7) in enumerate(prins7):
    ry7 = Inches(1.85+j*1.12)
    rect(sl, rx7+Inches(0.2), ry7, Inches(0.06), Inches(0.78), fill=c7)
    txt(sl, t7, rx7+Inches(0.38), ry7+Inches(0.02),
        Inches(1.4), Inches(0.38), sz=13, bold=True, col=c7)
    mtxt(sl, d7.split('\n'), rx7+Inches(0.38), ry7+Inches(0.42),
         Inches(3.4), Inches(0.62), sz=11.5, col=GR)

# ── 學習路徑（底部，三列分開）──
rect(sl, Inches(0.35), Inches(5.08), Inches(8.55), Inches(2.02),
     fill=PNL, bc=PNL2, bpt=1, rad=10000)
txt(sl, '建議學習路徑', Inches(0.55), Inches(5.14),
    Inches(8.1), Inches(0.35), sz=12, bold=True, col=DK)

path_series = [
    (FB, 'F', ['F1','F2','F3','F4','F5','F6']),
    (CG, 'C', ['C1','C2','C3','C4','C5','C6']),
    (AO, 'A', ['A1','A2','A3','A4','A5','A6']),
]
for j,(c7b,sl_lbl,units) in enumerate(path_series):
    row_y = Inches(5.55 + j*0.45)
    pill(sl, sl_lbl, Inches(0.42), row_y, Inches(0.5), Inches(0.33), c7b, W, sz=11)
    for k,u in enumerate(units):
        ux = Inches(1.02 + k*1.24)
        pill(sl, u, ux, row_y, Inches(1.08), Inches(0.33), c7b, W, sz=10)
    if j < 2:
        txt(sl, '↓', Inches(0.55), Inches(5.89+j*0.45),
            Inches(0.38), Inches(0.28), sz=11, col=GR, align=PP_ALIGN.CENTER)
page_num(sl, 7)

# ═══════════════════════════════════════════════════════════════
# S08  F 系列詳細
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '03  F 系列 ── 數學基礎先備（6 個單元）',
        '適用：國小中高年級至國中  ·  先備能力：能執行觸控動作', bg=FB)

f_units8 = [
    ('F1','一對一對應','拖曳物件配對\n建立數量等量直觀','能執行觸控','拖曳配對'),
    ('F2','唱  數',    '依序點擊物件計數\n語音同步播報','能識別圖片','語音計數'),
    ('F3','數字認讀',  '看圖選數字\n支援自訂主題上傳','能辨識 1–10','選擇配對'),
    ('F4','數字排序',  '拖曳數字由小到大\n或由大到小','能完成 F3','拖曳排序'),
    ('F5','量  比  較','比較兩組物品多寡\n建立大小概念','能完成 F2','比較判斷'),
    ('F6','數的組成',  '合成與分解概念\n理解數的結構','能完成 F4/F5','合成分解'),
]
for i,(code,name,skill,prereq,badge) in enumerate(f_units8):
    col = i%3; row = i//3
    bx = Inches(0.35+col*4.32)
    by = Inches(1.15+row*2.85)
    rect(sl, bx, by, Inches(4.12), Inches(2.68),
         fill=W, bc=FBL, bpt=1.8, rad=14000)
    rect(sl, bx, by, Inches(4.12), Inches(0.56), fill=FB, rad=0)
    txt(sl, code, bx+Inches(0.12), by+Inches(0.07),
        Inches(0.7), Inches(0.42), sz=17, bold=True, col=W)
    txt(sl, name, bx+Inches(0.88), by+Inches(0.1),
        Inches(2.9), Inches(0.4), sz=16, bold=True, col=W)
    pill(sl, badge, bx+Inches(0.12), by+Inches(0.65),
         Inches(1.2), Inches(0.32), FBL, FB, sz=11)
    mtxt(sl, skill.split('\n'), bx+Inches(0.18), by+Inches(1.08),
         Inches(3.74), Inches(0.98), sz=13, col=DK)
    rect(sl, bx+Inches(0.18), by+Inches(2.18), Inches(3.74), Inches(0.38),
         fill=PNL, rad=6000)
    txt(sl, f'先備：{prereq}', bx+Inches(0.3), by+Inches(2.22),
        Inches(3.54), Inches(0.32), sz=11, col=GR)
page_num(sl, 8)

# ═══════════════════════════════════════════════════════════════
# S09  C 系列詳細
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '04  C 系列 ── 貨幣認知建立（6 個單元）',
        '適用：國中至高中  ·  台灣 9 種面額  ·  拖曳操作  ·  三種難度', bg=CG)

c_units9 = [
    ('C1','認識錢幣','辨識台灣 9 種面額\n看圖選字 / 看字選幣 / 聽音選幣','能完成 F3','三種模式'),
    ('C2','數  錢',  '點擊錢幣依序累加\n語音播報累計金額','能完成 C1','累加計數'),
    ('C3','換  錢',  '以不同面額組合等值\n大換小 / 小換大 / 全隨機','能完成 C2','面額兌換'),
    ('C4','付  款',  '拖曳正確金額完成付款\n自動判定超額/不足','能完成 C3','付款操作'),
    ('C5','夠不夠',  '拖曳錢幣至付款區\n判斷金額能否購買','能完成 C4','比較判斷'),
    ('C6','找  零',  '拖曳找零錢幣至對應格\n驗收完整交易閉環','能完成 C5','找零驗收'),
]
for i,(code,name,skill,prereq,badge) in enumerate(c_units9):
    col = i%3; row = i//3
    bx = Inches(0.35+col*4.32)
    by = Inches(1.15+row*2.85)
    rect(sl, bx, by, Inches(4.12), Inches(2.68),
         fill=W, bc=CGL, bpt=1.8, rad=14000)
    rect(sl, bx, by, Inches(4.12), Inches(0.56), fill=CG, rad=0)
    txt(sl, code, bx+Inches(0.12), by+Inches(0.07),
        Inches(0.7), Inches(0.42), sz=17, bold=True, col=W)
    txt(sl, name, bx+Inches(0.88), by+Inches(0.1),
        Inches(2.9), Inches(0.4), sz=16, bold=True, col=W)
    pill(sl, badge, bx+Inches(0.12), by+Inches(0.65),
         Inches(1.2), Inches(0.32), CGL, CG, sz=11)
    mtxt(sl, skill.split('\n'), bx+Inches(0.18), by+Inches(1.08),
         Inches(3.74), Inches(0.98), sz=13, col=DK)
    rect(sl, bx+Inches(0.18), by+Inches(2.18), Inches(3.74), Inches(0.38),
         fill=PNL, rad=6000)
    txt(sl, f'先備：{prereq}', bx+Inches(0.3), by+Inches(2.22),
        Inches(3.54), Inches(0.32), sz=11, col=GR)
page_num(sl, 9)

# ═══════════════════════════════════════════════════════════════
# S10  A 系列 工作分析橫條
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '05  A 系列 ── 真實情境類化（6 個單元）',
        '工作分析法拆解 5～8 個操作步驟  ·  台灣在地化仿真介面', bg=AO)

a_units10 = [
    ('A1','🥤 自動販賣機','5',
     ['選商品','投幣','確認金額','取貨','核對找零'],
     '幣值累計・找零驗收'),
    ('A2','✂️ 理髮廳自助機','6',
     ['選類別','選服務','確認','付款','取號','等候'],
     '服務辨識・等候社會規範'),
    ('A3','🍔 速食店點餐機','7',
     ['瀏覽','加購物車','確認','選付款','付款','等取號','取餐'],
     '分類導覽・多品項加總'),
    ('A4','🛒 超市購物','6',
     ['選商品','確認價格','選付款','付款','找零','完成'],
     '購物清單執行・找零驗收'),
    ('A5','🏧 ATM 操作','7',
     ['插卡','密碼','選功能','輸金額','取款','退卡','取收據'],
     '密碼保護意識・金融交易流程'),
    ('A6','🚂 火車票購票機','8',
     ['選出發站','選目的地','選票種','確認票價','付款','取票','找零','完成'],
     '路線規劃・多步驟序列整合'),
]
for i,(code,name,steps,seq_list,skill) in enumerate(a_units10):
    by10 = Inches(1.15+i*0.98)
    alt = (i%2==1)
    row_bg = AOL if not alt else W
    rect(sl, Inches(0.35), by10, Inches(12.6), Inches(0.88),
         fill=row_bg, bc=RGBColor(0xFF,0xD7,0xB4), bpt=0.8, rad=8000)
    rect(sl, Inches(0.35), by10, Inches(0.9), Inches(0.88), fill=AO, rad=8000)
    txt(sl, code, Inches(0.35), by10, Inches(0.9), Inches(0.88),
        sz=14, bold=True, col=W, align=PP_ALIGN.CENTER)
    txt(sl, name, Inches(1.32), by10+Inches(0.05),
        Inches(2.05), Inches(0.48), sz=13, bold=True, col=AO)
    pill(sl, f'{steps} 步', Inches(1.32), by10+Inches(0.52),
         Inches(0.72), Inches(0.28), AOL, AOD, sz=10)
    txt(sl, skill, Inches(2.12), by10+Inches(0.52),
        Inches(1.75), Inches(0.32), sz=10, col=GR, italic=True)
    for j,step in enumerate(seq_list):
        sx = Inches(4.05+j*1.38)
        ellipse(sl, sx, by10+Inches(0.2), Inches(0.34), Inches(0.34), fill=AO)
        txt(sl, str(j+1), sx, by10+Inches(0.2), Inches(0.34), Inches(0.34),
            sz=11, bold=True, col=W, align=PP_ALIGN.CENTER)
        txt(sl, step, sx-Inches(0.28), by10+Inches(0.54),
            Inches(0.9), Inches(0.28), sz=9.5, col=DK, align=PP_ALIGN.CENTER)
        if j < len(seq_list)-1:
            txt(sl, '→', sx+Inches(0.36), by10+Inches(0.22),
                Inches(0.28), Inches(0.3), sz=11, col=GR, align=PP_ALIGN.CENTER)
page_num(sl, 10)

# ═══════════════════════════════════════════════════════════════
# S11  輔助點擊模式
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '06  輔助點擊模式（Click Mode） ── 特教數位教材首創',
        '精細動作障礙學生：任意位置點擊 → 自動執行下一個正確步驟', bg=PU)

rect(sl, Inches(0.35), Inches(1.12), Inches(5.8), Inches(5.98),
     fill=PUL, bc=PU, bpt=2, rad=16000)
txt(sl, '🖱️', Inches(0.35), Inches(1.18), Inches(5.8), Inches(0.85),
    sz=48, align=PP_ALIGN.CENTER)
txt(sl, '設計理念', Inches(0.55), Inches(2.08),
    Inches(5.4), Inches(0.42), sz=14, bold=True, col=PU)
design_lines = [
    '傳統數位教材要求點擊特定目標元素',
    '→ 精細動作障礙學生難以精確定位',
    '',
    '輔助點擊模式：「位置無關」點擊',
    '→ 只要在畫面任意位置點擊',
    '→ 系統自動識別當前應執行的步驟',
    '→ 並執行該步驟的正確操作',
    '',
    '效果等同「教師手把手協助」',
    '但完全由學生主動觸發',
]
y11 = Inches(2.52)
for line in design_lines:
    if line == '':
        y11 += Inches(0.12)
        continue
    c11 = PU if line.startswith('→') else DK
    txt(sl, line, Inches(0.55), y11, Inches(5.4), Inches(0.42), sz=12.5, col=c11)
    y11 += Inches(0.38)

mx11 = Inches(6.45)
features11 = [
    (PU,  '全 18 單元支援',   '所有教學單元皆可啟用\n設定頁一鍵開關'),
    (FB,  '三種難度均支援',   '簡單/普通/困難模式\n均可搭配輔助點擊'),
    (CG,  '步驟自動推進',    '每次點擊執行下一步\n不需精確定位目標'),
    (AO,  'IEP 評量保留',    '輔助模式下仍記錄操作\n用於形成性評量'),
]
for j,(c11,t11,d11) in enumerate(features11):
    ry11 = Inches(1.15+j*1.48)
    rect(sl, mx11, ry11, Inches(3.25), Inches(1.32),
         fill=W, bc=c11, bpt=1.8, rad=12000)
    ellipse(sl, mx11+Inches(0.2), ry11+Inches(0.28), Inches(0.7), Inches(0.7), fill=c11)
    txt(sl, '✓', mx11+Inches(0.2), ry11+Inches(0.28), Inches(0.7), Inches(0.7),
        sz=18, bold=True, col=W, align=PP_ALIGN.CENTER)
    txt(sl, t11, mx11+Inches(1.0), ry11+Inches(0.08),
        Inches(2.1), Inches(0.42), sz=13, bold=True, col=c11)
    mtxt(sl, d11.split('\n'), mx11+Inches(1.0), ry11+Inches(0.52),
         Inches(2.1), Inches(0.68), sz=11, col=GR)

rx11 = Inches(9.95)
rect(sl, rx11, Inches(1.12), Inches(3.1), Inches(5.98),
     fill=W, bc=PNL2, bpt=1.2, rad=12000)
rect(sl, rx11, Inches(1.12), Inches(3.1), Inches(0.5), fill=PU, rad=0)
txt(sl, '適用學生類型', rx11+Inches(0.2), Inches(1.16),
    Inches(2.7), Inches(0.44), sz=14, bold=True, col=W)
stu_types = [
    (PUL, PU, '腦性麻痺',     '上肢動作控制困難\n難以精確點擊小目標'),
    (FBL, FB, '多重障礙',     '精細動作合併認知\n需要高度結構支持'),
    (CGL, CG, '發展遲緩',     '觸控操作尚在建立\n需要成功經驗鼓勵'),
    (AOL, AO, '視覺動作整合困難','難以手眼協調\n定位特定目標'),
]
for j,(bg11,bc11,t11b,d11b) in enumerate(stu_types):
    ry11b = Inches(1.75+j*1.18)
    rect(sl, rx11+Inches(0.2), ry11b, Inches(2.7), Inches(1.05),
         fill=bg11, bc=bc11, bpt=1, rad=8000)
    txt(sl, t11b, rx11+Inches(0.35), ry11b+Inches(0.05),
        Inches(2.4), Inches(0.38), sz=13, bold=True, col=bc11)
    mtxt(sl, d11b.split('\n'), rx11+Inches(0.35), ry11b+Inches(0.45),
         Inches(2.4), Inches(0.52), sz=11, col=DK)
page_num(sl, 11)

# ═══════════════════════════════════════════════════════════════
# S12  配套系統
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '完整配套系統 ── 四大輔助模組', bg=DK)

modules12 = [
    (AM, AML, '📝', '配套作業單系統',
     ['18 個單元 × 10 種題型',
      '一鍵列印 / PDF 下載',
      '題型包含：',
      '  • 看圖填空（幣值辨識）',
      '  • 找零計算',
      '  • 付款組合選擇',
      '  • 序列排序',
      '  • 問答計算',
      '完整對應 IEP 評量記錄']),
    (PU, PUL, '🏆', '獎勵積分系統',
     ['獨立模組（reward/）',
      '功能：',
      '  • 學生照片管理',
      '  • 個人積分累計',
      '  • 課堂即時頒獎',
      '  • 積分記錄查詢',
      '整合 IEP 行為正增強',
      '教師課堂即時操作']),
    (FB, FBL, '🔊', '全程語音引導系統',
     ['Web Speech API',
      'Microsoft 雅婷優先（Windows）',
      '台灣語境金額唸法',
      '（「兩百元」非「二百元」）',
      '即時讚美語音回饋',
      '錯誤提示語音',
      '步驟引導語音',
      '完全離線語音支援']),
    (TE, TEL, '🌙', '深色模式 / 無障礙',
     ['全 18 單元支援深色模式',
      '適合：',
      '  • 視覺敏感學生',
      '  • 視覺疲勞情況',
      '  • 弱視輔助需求',
      '一鍵切換，無需重載',
      '跨裝置保持設定',
      '零額外安裝費用']),
]
for i,(bc,bl,icon,t12,body) in enumerate(modules12):
    bx = Inches(0.35+i*3.23)
    rect(sl, bx, Inches(1.12), Inches(3.05), Inches(5.98),
         fill=W, bc=bc, bpt=1.8, rad=14000)
    rect(sl, bx, Inches(1.12), Inches(3.05), Inches(0.52), fill=bc, rad=0)
    txt(sl, icon, bx, Inches(1.15), Inches(3.05), Inches(0.46),
        sz=20, align=PP_ALIGN.CENTER, col=W)
    txt(sl, t12, bx, Inches(1.72), Inches(3.05), Inches(0.52),
        sz=13, bold=True, col=bc, align=PP_ALIGN.CENTER)
    hline(sl, Inches(2.28), bx+Inches(0.22), bx+Inches(2.83), col=bc, pt=0.8)
    for k,line in enumerate(body):
        is_sec = line.endswith('：') or line.endswith(':')
        c12 = bc if is_sec else (DK if not line.startswith('  ') else GR)
        txt(sl, line, bx+Inches(0.18), Inches(2.38+k*0.43),
            Inches(2.7), Inches(0.41), sz=11.5 if not is_sec else 12,
            bold=is_sec, col=c12)
page_num(sl, 12)

# ═══════════════════════════════════════════════════════════════
# S13  三段學習歷程（詳細）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '07  三段式學習歷程設計',
        '數位模擬練習 → 紙本形成性評量 → 實地類化應用', bg=TE)

stages13 = [
    (FB, FBD, FBL, '第一段', '數位模擬練習', '課堂主軸',
     ['軟體情境操作（18 個互動單元）',
      '即時觀察學生錯誤類型',
      '三種難度對應個別需求',
      '輔助點擊支援精細動作障礙',
      '提供教師調整教學策略依據',
      '記錄：答對題數 / 嘗試次數'],
     '💻  主要工具：平板電腦 / Windows 電腦'),
    (CG, CGD, CGL, '第二段', '紙本形成性評量', '課後作業',
     ['列印對應作業單（18 單元備有）',
      '每單元 10 種題型可選',
      'PDF 下載或直接列印',
      '評估數位學習是否遷移至紙筆',
      '結果可記入 IEP 形成性評量',
      '教師可依能力調整題數'],
     '📄  主要工具：作業單 / 評量記錄表'),
    (AO, AOD, AOL, '第三段', '實地類化應用', '社區轉銜',
     ['真實社區場景操作演練',
      '行為檢核表記錄達成程度',
      '教師逐步褪除支持策略',
      '最終目標：教師 3 公尺外觀察',
      '達成完全獨立操作',
      '作為 IEP 轉銜評估依據'],
     '🏘️  場景：超市 / 便利商店 / ATM'),
]
for i,(bc,bcd,bl,num,t,sub,body,tool) in enumerate(stages13):
    bx = Inches(0.35+i*4.32)
    rect(sl, bx, Inches(1.12), Inches(4.1), Inches(6.05),
         fill=bl, bc=bc, bpt=2.5, rad=16000)
    ellipse(sl, bx+Inches(1.5), Inches(1.22), Inches(1.1), Inches(1.1), fill=bc)
    txt(sl, num, bx+Inches(1.5), Inches(1.22), Inches(1.1), Inches(1.1),
        sz=14, bold=True, col=W, align=PP_ALIGN.CENTER)
    txt(sl, t, bx, Inches(2.42), Inches(4.1), Inches(0.52),
        sz=18, bold=True, col=bcd, align=PP_ALIGN.CENTER)
    txt(sl, sub, bx, Inches(2.92), Inches(4.1), Inches(0.35),
        sz=12, col=GR, align=PP_ALIGN.CENTER, italic=True)
    hline(sl, Inches(3.35), bx+Inches(0.3), bx+Inches(3.8), col=bc, pt=1.2)
    for k,line in enumerate(body):
        rect(sl, bx+Inches(0.22), Inches(3.48+k*0.48)+Inches(0.14),
             Inches(0.1), Inches(0.2), fill=bc, rad=5000)
        txt(sl, line, bx+Inches(0.42), Inches(3.48+k*0.48),
            Inches(3.5), Inches(0.46), sz=12, col=DK)
    rect(sl, bx+Inches(0.18), Inches(6.45), Inches(3.74), Inches(0.58),
         fill=W, bc=bc, bpt=1, rad=8000)
    txt(sl, tool, bx+Inches(0.3), Inches(6.48),
        Inches(3.54), Inches(0.5), sz=11, col=bc, bold=True)
for ax in [Inches(4.58), Inches(8.9)]:
    txt(sl, '➡', ax, Inches(4.05), Inches(0.6), Inches(0.55),
        sz=28, col=GR, align=PP_ALIGN.CENTER)
page_num(sl, 13)

# ═══════════════════════════════════════════════════════════════
# S14  課程法規依據
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '08  課程法規依據與 IEP 對應', bg=DK)

laws14 = [
    (RD, '🏛️', '特殊教育法 第 28 條',
     '應為身心障礙學生訂定 IEP',
     ['各單元提供 IEP 年度目標範本','提供評量基準與記錄工具','直接對應轉銜目標格式']),
    (AO, '📋', '特殊需求領域課綱\n生活管理・財務管理',
     '生管-Ⅳ-E-1 認識貨幣\n生管-Ⅳ-E-2 選取幣值付款\n生管-Ⅳ-E-3 驗收找零\n生管-Ⅳ-F-1 使用自動化設備',
     ['F/C/A 三系列 18 單元','完整對應財務管理序階','B 系列延伸預算規劃主題']),
    (FB, '🌏', '身心障礙者權利公約\nCRPD 第 9・19 條',
     '無障礙近用設施\n獨立生活於社區',
     ['輔助點擊消除動作障礙','A 系列六種真實社區場景','服務 IEP 轉銜社區獨立目標']),
    (CG, '📚', '十二年國教課綱總綱',
     '自主行動・社會參與\n核心素養',
     ['A 系列社區真實場景應用','F 系列數學概念建構','跨領域融合設計']),
    (PU, '⚖️', '高中以下特教\n教材教法評量辦法 第 6 條',
     '教材應符合學生\n個別學習需求',
     ['三種難度 × 輔助點擊','支援個別化差異教學','完全免費、無授權費用']),
    (AM, '📖', '十二年國教\n特殊需求領域課綱（2019）',
     '財務管理主題\n頁 45–52',
     ['系統性對應課綱架構','提供教學目標與評量規準','可作教師自我效能評估依據']),
]
for i,(bc,icon,law_t,law_d,resp) in enumerate(laws14):
    col=i%2; row=i//2
    bx = Inches(0.35+col*6.48)
    by = Inches(1.15+row*1.95)
    rect(sl, bx, by, Inches(6.2), Inches(1.82),
         fill=W, bc=PNL2, bpt=1.2, rad=10000)
    rect(sl, bx, by, Inches(0.1), Inches(1.82), fill=bc, rad=0)
    txt(sl, icon, bx+Inches(0.18), by+Inches(0.12),
        Inches(0.58), Inches(0.6), sz=22)
    mtxt(sl, law_t.split('\n'), bx+Inches(0.82), by+Inches(0.06),
         Inches(2.4), Inches(0.72), sz=12.5, bold=True, col=bc)
    mtxt(sl, law_d.split('\n'), bx+Inches(0.82), by+Inches(0.78),
         Inches(2.5), Inches(0.9), sz=11, col=GR)
    vline(sl, bx+Inches(3.5), by+Inches(0.15), by+Inches(1.67), col=PNL2)
    txt(sl, '本軟體對應', bx+Inches(3.62), by+Inches(0.08),
        Inches(2.4), Inches(0.32), sz=11, bold=True, col=bc)
    for k,r14 in enumerate(resp):
        txt(sl, f'✓ {r14}', bx+Inches(3.62), by+Inches(0.42+k*0.42),
            Inches(2.4), Inches(0.4), sz=11, col=DK)
page_num(sl, 14)

# ═══════════════════════════════════════════════════════════════
# S15  試教觀察
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '09  試教觀察 ── 類化窗口的關鍵時刻', bg=CG)

rect(sl, Inches(0.35), Inches(1.14), Inches(12.6), Inches(1.18),
     fill=CGL, bc=CG, bpt=1.8, rad=12000)
txt(sl, '🔬  共同轉折點', Inches(0.55), Inches(1.22),
    Inches(4), Inches(0.42), sz=14, bold=True, col=CG)
txt(sl, '學生在第 5～7 次軟體操作後，開始在啟動前主動低聲說出下一步名稱',
    Inches(4.65), Inches(1.22), Inches(8.1), Inches(0.42),
    sz=14, bold=True, col=DK)
txt(sl, '→ 程序性記憶語言化的信號，也是類化窗口即將開啟的前兆',
    Inches(4.65), Inches(1.64), Inches(8.1), Inches(0.38),
    sz=12.5, col=CG, italic=True)

obs_cases = [
    (CG, CGL, 'A2  自閉症學生的突破',
     ['試教背景：曾因等待焦慮情緒失調',
      '         一學期拒絕接觸相關學習',
      '第三節社區實地操作：',
      '→ 面對真實理髮廳觸控機',
      '→ 獨立完成完整付款流程',
      '→ 等待期間主動說：',
      '   「我是 7 號，還沒叫到我，我要繼續等。」',
      '',
      '意義：不只技能習得，更是社交規範的內化'],
     ['⭐ 跨越性成長指標',
      '• 等待行為建立',
      '• 主動語言表達',
      '• 情境焦慮顯著下降']),
    (FB, FBL, 'A1  拒絕接觸設備長達六個月的學生',
     ['試教背景：多次在真實販賣機失敗',
      '         主動繞道迴避所有自助設備',
      '完成軟體 15 次練習後：',
      '→ 採用「先投幣後選商品」模式',
      '→ 第一次在真實販賣機前完整購買',
      '→ 購買完成後主動說：',
      '   「我下次自己來。」',
      '',
      '意義：從「情境迴避」到「主動自信」'],
     ['⭐ 突破性轉化指標',
      '• 情境迴避行為消除',
      '• 主動自信表達建立',
      '• 獨立操作意願確立']),
]
for i,(bc,bl,title,body,highlight) in enumerate(obs_cases):
    bx = Inches(0.35+i*6.48)
    rect(sl, bx, Inches(2.48), Inches(6.2), Inches(4.68),
         fill=W, bc=bc, bpt=2, rad=14000)
    rect(sl, bx, Inches(2.48), Inches(6.2), Inches(0.5), fill=bc, rad=0)
    txt(sl, title, bx+Inches(0.2), Inches(2.52),
        Inches(5.8), Inches(0.42), sz=15, bold=True, col=W)
    for k,line in enumerate(body):
        is_arrow = line.startswith('→')
        is_quote = line.startswith('   「')
        is_blank = line == ''
        if is_blank:
            continue
        c15 = bc if is_arrow else (AM if is_quote else DK)
        txt(sl, line, bx+Inches(0.22), Inches(3.05+k*0.38),
            Inches(3.85), Inches(0.36), sz=12.5, col=c15,
            bold=is_quote, italic=is_arrow)
    for m,h15 in enumerate(highlight):
        hy = Inches(3.22+m*0.62)
        bold15 = h15.startswith('⭐')
        rect(sl, bx+Inches(4.22), hy, Inches(1.78), Inches(0.52),
             fill=bl, bc=bc, bpt=1, rad=8000)
        txt(sl, h15, bx+Inches(4.3), hy+Inches(0.04),
            Inches(1.65), Inches(0.44), sz=11,
            bold=bold15, col=bc)
page_num(sl, 15)

# ═══════════════════════════════════════════════════════════════
# S16  預期成效指標
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '10  預期成效指標 ── 三層面 × 短中長期', bg=CG)

periods = ['短期（1～2 個月）','中期（3～6 個月）','長期（6 個月以上）']
domains = [('🧠','認知'),('🤝','技能'),('❤️','情意')]
domain_colors = [(FB,FBL),(CG,CGL),(RD,RDL)]
period_bc = [AM, CG, PU]

cells = [
    ['9 種面額正確辨識 ≥ 80%\n能說出各 A 系列操作步驟',
     '找零計算正確率 ≥ 70%\n能跨情境類化付款概念',
     '能向他人說明消費流程\n理解金融安全概念'],
    ['A1～A3 普通模式完成\nC4 付款 5 次正確',
     'A4～A6 普通模式完成\nATM 提款獨立操作',
     '真實場景（理髮廳/超市/ATM）\n獨立完成全流程'],
    ['面對自助設備焦慮降低\n主動要求練習',
     '不再迴避社區場景\n主動啟動操作',
     '建立社區消費自主信心\n積極規劃獨立生活'],
]
for j,p in enumerate(periods):
    bx = Inches(3.15+j*3.38)
    rect(sl, bx, Inches(1.12), Inches(3.2), Inches(0.5), fill=period_bc[j], rad=0)
    txt(sl, p, bx, Inches(1.14), Inches(3.2), Inches(0.46),
        sz=13, bold=True, col=W, align=PP_ALIGN.CENTER)

for i,(icon,dom) in enumerate(domains):
    dc,dcl = domain_colors[i]
    by16 = Inches(1.66+i*1.88)
    rect(sl, Inches(0.35), by16, Inches(2.7), Inches(1.78),
         fill=dcl, bc=dc, bpt=1.5, rad=10000)
    txt(sl, icon, Inches(0.35), by16+Inches(0.12),
        Inches(2.7), Inches(0.65), sz=28, align=PP_ALIGN.CENTER)
    txt(sl, dom, Inches(0.35), by16+Inches(0.8),
        Inches(2.7), Inches(0.5), sz=16, bold=True, col=dc, align=PP_ALIGN.CENTER)
    for j,cell in enumerate(cells[i]):
        bx = Inches(3.15+j*3.38)
        rect(sl, bx, by16, Inches(3.2), Inches(1.78),
             fill=W, bc=period_bc[j], bpt=1, rad=8000)
        mtxt(sl, cell.split('\n'), bx+Inches(0.18), by16+Inches(0.22),
             Inches(2.84), Inches(1.3), sz=12.5, col=DK)
page_num(sl, 16)

# ═══════════════════════════════════════════════════════════════
# S17  技術規格
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '14  技術規格與無障礙設計', 'Technical Specifications & Accessibility', bg=DK)

specs17 = [
    ('🖥️','開發技術',
     'HTML5 / CSS3 / JavaScript（ES6+）\n無需安裝任何外掛程式'),
    ('📏','程式規模',
     '18 個主程式，超過 15 萬行程式碼\n研發歷時約 4 個月'),
    ('📱','支援裝置',
     'Windows 電腦（推薦）\niPad / Android 平板（9.7 吋以上）'),
    ('🌐','支援瀏覽器',
     'Google Chrome / Microsoft Edge\n（最新版）'),
    ('📵','離線支援',
     '完整支援完全離線使用\n首次語音啟動後不需網路'),
    ('🔊','語音引擎',
     'Web Speech API\n優先 Microsoft 雅婷（台灣中文）\n台灣語境金額唸法'),
    ('🖱️','輔助點擊',
     '全 18 單元支援\n精細動作障礙學生任意位置點擊\n自動執行正確步驟'),
    ('🌙','深色模式',
     '全 18 單元支援\n適合視覺敏感或視覺疲勞學生\n一鍵切換'),
    ('💰','製作費用',
     '新臺幣 0 元\n教師自主研發\n無任何商業授權費用'),
]
for i,(icon,t,d) in enumerate(specs17):
    col=i%3; row=i//3
    bx = Inches(0.35+col*4.32)
    by = Inches(1.15+row*1.95)
    rect(sl, bx, by, Inches(4.12), Inches(1.82),
         fill=W, bc=PNL2, bpt=1.2, rad=12000)
    rect(sl, bx, by, Inches(4.12), Inches(0.12), fill=DK, rad=0)
    txt(sl, icon, bx+Inches(0.15), by+Inches(0.2),
        Inches(0.65), Inches(0.65), sz=26)
    txt(sl, t, bx+Inches(0.88), by+Inches(0.2),
        Inches(3.0), Inches(0.42), sz=14, bold=True, col=DK)
    mtxt(sl, d.split('\n'), bx+Inches(0.88), by+Inches(0.65),
         Inches(3.0), Inches(1.0), sz=11.5, col=GR)
page_num(sl, 17)

# ═══════════════════════════════════════════════════════════════
# S18  核心創新總結
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BK)
rect(sl, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill=CG)
rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.12), fill=AM)
txt(sl, '15  核心創新總結', Inches(0.45), Inches(0.22),
    Inches(9), Inches(0.55), sz=28, bold=True, col=W)
txt(sl, 'Key Innovations  ── 四大首創特色',
    Inches(0.45), Inches(0.75), Inches(9), Inches(0.38),
    sz=14, col=RGBColor(0x94,0xA3,0xB8), fn='Arial', italic=True)
hline(sl, Inches(1.22), Inches(0.45), Inches(12.8),
      col=RGBColor(0x33,0x4A,0x6E), pt=1)

innovations18 = [
    (FB, '①',
     '全台首套三層次序階\n數位金融素養課程',
     'F→C→A 三系列 18 個互動教學單元\n系統覆蓋「數學先備→貨幣認知→情境類化」完整序列\n每系列 6 單元，每單元三種難度 × 輔助點擊'),
    (CG, '②',
     '輔助點擊模式\n（特教數位教材首創）',
     '精細動作障礙學生「任意位置點擊」即自動執行正確步驟\n全 18 單元支援，消除動作能力門檻\n讓操作能力不等於學習機會'),
    (AO, '③',
     '六種台灣在地化\n真實情境仿真介面',
     '自動販賣機・理髮廳自助機・速食店點餐機\n超市購物・ATM・台鐵購票機\n台灣在地化設備外觀 + 真實費率 + 完整操作流程'),
    (PU, '④',
     '完整配套系統\n對應 IEP 各功能向度',
     '作業單系統（18 × 10 題型）× 獎勵積分模組\n全程語音引導 × 深色模式 × PDF 下載\n教師可直接用於 IEP 評量記錄與轉銜規劃'),
]
for i,(bc,num,title,desc) in enumerate(innovations18):
    col=i%2; row=i//2
    bx = Inches(0.45+col*6.4)
    by = Inches(1.35+row*2.75)
    rect(sl, bx, by, Inches(6.1), Inches(2.6),
         fill=RGBColor(0x1E,0x3A,0x5C), bc=bc, bpt=2, rad=14000)
    rect(sl, bx, by, Inches(0.7), Inches(2.6), fill=bc, rad=0)
    txt(sl, num, bx, by, Inches(0.7), Inches(2.6),
        sz=32, bold=True, col=W, align=PP_ALIGN.CENTER, fn='Arial')
    for j,tl in enumerate(title.split('\n')):
        txt(sl, tl, bx+Inches(0.78), by+Inches(0.08+j*0.52),
            Inches(5.12), Inches(0.5), sz=17, bold=True, col=bc)
    hline(sl, by+Inches(1.18), bx+Inches(0.78), bx+Inches(5.9),
          col=RGBColor(0x33,0x4A,0x6E), pt=1)
    for k,dl in enumerate(desc.split('\n')):
        txt(sl, dl, bx+Inches(0.78), by+Inches(1.3+k*0.42),
            Inches(5.12), Inches(0.4), sz=12.5, col=RGBColor(0xCB,0xD5,0xE1))
rect(sl, Inches(0.45), Inches(6.92), Inches(12.4), Inches(0.45),
     fill=RGBColor(0x1E,0x3A,0x5C), rad=10000)
txt(sl, '身心障礙類・國民中學及高級中等學校組  ｜  製作費用 NT$0  ｜  完全離線・免安裝  ｜  中華民國 115 年 3 月',
    Inches(0.45), Inches(6.94), Inches(12.4), Inches(0.41),
    sz=12, col=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)
page_num(sl, 18)

# ═══════════════════════════════════════════════════════════════
# S19  B 系列延伸 — 預算規劃 6 單元
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '11  B 系列延伸 ── 預算規劃（6 個單元）',
        '適用：已完成 A 系列  ·  延伸：日常預算管理・比價・存錢計畫', bg=PU)

b_units19 = [
    ('B1','今天帶多少錢',
     '依行程計算所需金額\n簡單/普通/困難三難度',
     '能完成 A 系列',
     '行程計算'),
    ('B2','零用錢日記',
     '計算一週收支與餘額\n收入/支出事件序列',
     '能完成 B1',
     '收支記帳'),
    ('B3','存錢計畫',
     '計算每週需存幾元\n月曆模擬實際存錢',
     '能完成 B2',
     '目標儲蓄'),
    ('B4','特賣比一比',
     '比較兩店商品價格\n找出較便宜選項',
     '能完成 B1',
     '價格比較'),
    ('B5','生日派對預算',
     '在預算內規劃派對\n必買+選購商品決策',
     '能完成 B2',
     '預算分配'),
    ('B6','菜市場買菜',
     '按清單各攤採購\n付款並確認找零',
     '能完成 A4/A6',
     '實境採購'),
]
for i,(code,name,skill,prereq,badge) in enumerate(b_units19):
    col = i%3; row = i//3
    bx = Inches(0.35+col*4.32)
    by = Inches(1.12+row*2.85)
    rect(sl, bx, by, Inches(4.12), Inches(2.68),
         fill=W, bc=PUL, bpt=1.8, rad=14000)
    rect(sl, bx, by, Inches(4.12), Inches(0.56), fill=PU, rad=0)
    txt(sl, code, bx+Inches(0.12), by+Inches(0.07),
        Inches(0.7), Inches(0.42), sz=17, bold=True, col=W)
    txt(sl, name, bx+Inches(0.88), by+Inches(0.1),
        Inches(2.9), Inches(0.4), sz=15, bold=True, col=W)
    pill(sl, badge, bx+Inches(0.12), by+Inches(0.65),
         Inches(1.2), Inches(0.32), PUL, PU, sz=11)
    mtxt(sl, skill.split('\n'), bx+Inches(0.18), by+Inches(1.08),
         Inches(3.74), Inches(0.95), sz=13, col=DK)
    rect(sl, bx+Inches(0.18), by+Inches(2.18), Inches(3.74), Inches(0.38),
         fill=PNL, rad=6000)
    txt(sl, f'先備：{prereq}', bx+Inches(0.3), by+Inches(2.22),
        Inches(3.54), Inches(0.32), sz=11, col=GR)
page_num(sl, 19)

# ═══════════════════════════════════════════════════════════════
# S20  IEP 具體教學目標範例
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '12  IEP 具體教學目標 ── 各系列年度目標範本',
        '依 IEP 目標格式：行為主體 + 情境 + 操作行為 + 標準', bg=AM)

# 欄標頭
headers_iep = ['系列', 'IEP 年度目標範本', '評量標準', '轉銜應用場域']
hcols = [AM, AM, AM, AM]
col_w = [Inches(1.0), Inches(5.8), Inches(2.8), Inches(3.3)]
col_x = [Inches(0.35), Inches(1.38), Inches(7.22), Inches(10.05)]

for j,(hw,hx) in enumerate(zip(col_w, col_x)):
    rect(sl, hx, Inches(1.12), hw, Inches(0.48), fill=AM, rad=0)
    txt(sl, headers_iep[j], hx, Inches(1.14), hw, Inches(0.44),
        sz=13, bold=True, col=W, align=PP_ALIGN.CENTER)

iep_rows = [
    (FB, FBL, 'F 系列',
     '學生能在數字排序活動中，依序將 1～10 的數字卡由小到大排列完成',
     '3 次中有 2 次正確完成，教師不提示',
     '數字辨識・日常排序活動'),
    (CG, CGL, 'C 系列',
     '學生能在付款情境中，正確辨識並拿取等值面額完成 NT$50 以內付款',
     '連續 4 次答對，錯誤不超過 1 次',
     '便利商店・超市結帳'),
    (AO, AOL, 'A 系列',
     '學生能在自動販賣機前，獨立完成投幣、選商品、取貨、核對找零之完整流程',
     '教師 3 公尺外觀察，5 步驟均達成',
     '校園販賣機・超商自助結帳'),
    (PU, PUL, 'B 系列',
     '學生能依據當週活動計畫，計算所需帶出門金額，誤差在 NT$10 以內',
     '3 次活動中有 2 次自行計算正確',
     '假日出遊・社區購物規劃'),
]

for i,(bc,bl,series,goal,std,field) in enumerate(iep_rows):
    ry = Inches(1.65+i*1.38)
    row_bg = bl if i%2==0 else W
    rect(sl, Inches(0.35), ry, Inches(12.6), Inches(1.28),
         fill=row_bg, bc=PNL2, bpt=0.8, rad=6000)
    # 系列標籤
    rect(sl, Inches(0.35), ry, Inches(1.0), Inches(1.28), fill=bc, rad=6000)
    txt(sl, series, Inches(0.35), ry, Inches(1.0), Inches(1.28),
        sz=13, bold=True, col=W, align=PP_ALIGN.CENTER)
    # 目標
    mtxt(sl, [goal], Inches(1.42), ry+Inches(0.12),
         Inches(5.72), Inches(1.08), sz=12, col=DK)
    # 標準
    mtxt(sl, [std], Inches(7.26), ry+Inches(0.12),
         Inches(2.7), Inches(1.08), sz=11, col=bc, bold=True)
    # 場域
    mtxt(sl, [field], Inches(10.08), ry+Inches(0.12),
         Inches(2.72), Inches(1.08), sz=11.5, col=GR)

page_num(sl, 20)

# ═══════════════════════════════════════════════════════════════
# S21  三難度 × 輔助點擊 詳細比較
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '13  三種難度 × 輔助點擊 ── 差異化教學設計',
        '相同教材內容，以難度設計實現完全個別化', bg=TE)

# 四欄比較（簡單/普通/困難/輔助點擊）
col4 = [
    (CG,  CGL,  '簡單模式\nEasy',
     ['完整視覺提示','步驟動畫高亮','語音全程引導','答錯無限機會','自動提示介入'],
     '初學・習得期\nIEP 短期目標'),
    (AM,  AML,  '普通模式\nNormal',
     ['部分提示','步驟順序提示','錯誤才播語音','3 次後顯示提示','教師可設定'],
     '流暢性建立\nIEP 中期目標'),
    (RD,  RDL,  '困難模式\nHard',
     ['無視覺提示','無步驟引導','僅結果回饋','強化記憶策略','高度自主操作'],
     '獨立精熟\nIEP 精熟目標'),
    (PU,  PUL,  '輔助點擊\nClick Mode',
     ['任意位置點擊','自動執行步驟','適用三種難度','精細動作障礙','全 18 單元'],
     '動作障礙支持\n不限能力目標'),
]
for i,(bc,bl,title,pts,iep_note) in enumerate(col4):
    bx = Inches(0.35+i*3.23)
    rect(sl, bx, Inches(1.12), Inches(3.05), Inches(6.05),
         fill=W, bc=bc, bpt=2, rad=14000)
    rect(sl, bx, Inches(1.12), Inches(3.05), Inches(0.68), fill=bc, rad=0)
    mtxt(sl, title.split('\n'), bx, Inches(1.14), Inches(3.05), Inches(0.64),
         sz=15, bold=True, col=W, align=PP_ALIGN.CENTER)
    for k,pt in enumerate(pts):
        ry21 = Inches(1.88+k*0.72)
        rect(sl, bx+Inches(0.18), ry21, Inches(2.7), Inches(0.62),
             fill=bl, bc=bc, bpt=0.8, rad=8000)
        txt(sl, pt, bx+Inches(0.28), ry21+Inches(0.1),
            Inches(2.5), Inches(0.42), sz=12, col=bc, bold=True)
    rect(sl, bx+Inches(0.18), Inches(5.62), Inches(2.7), Inches(1.32),
         fill=bc, rad=10000)
    mtxt(sl, iep_note.split('\n'), bx, Inches(5.72), Inches(3.05), Inches(1.12),
         sz=13, bold=True, col=W, align=PP_ALIGN.CENTER)
page_num(sl, 21)

# ═══════════════════════════════════════════════════════════════
# S22  結語・感謝
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BK)
# 裝飾條
rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.15), fill=AO)
rect(sl, Inches(0), SLIDE_H-Inches(0.15), SLIDE_W, Inches(0.15), fill=CG)
# 左側彩色邊條
for j,(c22,y22,h22) in enumerate([(FB,Inches(0.15),Inches(2.45)),(CG,Inches(2.6),Inches(2.45)),(AO,Inches(5.05),Inches(2.45))]):
    rect(sl, Inches(0), y22, Inches(0.18), h22, fill=c22)

# 核心訊息
txt(sl, '謝謝評審委員', Inches(0.45), Inches(0.55), Inches(9.5), Inches(0.88),
    sz=52, bold=True, col=W)
txt(sl, 'Thank you for your consideration',
    Inches(0.45), Inches(1.45), Inches(9.5), Inches(0.42),
    sz=16, col=RGBColor(0x7C,0x8D,0xA8), fn='Arial', italic=True)
hline(sl, Inches(2.0), Inches(0.45), Inches(9.5), col=RGBColor(0x33,0x4A,0x6E), pt=1.5)

closing_points = [
    (FB, '🎯', '教育初衷',
     '每一位特教學生都有權利在社區中自主消費\n這套軟體是讓這個權利可以被練習的起點'),
    (CG, '🌱', '持續精進',
     '本軟體持續更新維護，依教學現場需求優化\n歡迎特教教師提供教學回饋共同改善'),
    (AO, '🔓', '完全開放',
     '製作費用 NT$0・完全離線・免安裝\n歡迎各校特教班自由下載使用'),
]
for i,(bc,icon,t22,d22) in enumerate(closing_points):
    bx = Inches(0.45+i*4.26)
    by = Inches(2.22)
    rect(sl, bx, by, Inches(4.02), Inches(2.42),
         fill=RGBColor(0x1E,0x3A,0x5C), bc=bc, bpt=2, rad=14000)
    txt(sl, icon, bx, by+Inches(0.15), Inches(4.02), Inches(0.62),
        sz=32, align=PP_ALIGN.CENTER)
    txt(sl, t22, bx, by+Inches(0.82), Inches(4.02), Inches(0.45),
        sz=16, bold=True, col=bc, align=PP_ALIGN.CENTER)
    hline(sl, by+Inches(1.35), bx+Inches(0.25), bx+Inches(3.77),
          col=RGBColor(0x33,0x4A,0x6E), pt=1)
    mtxt(sl, d22.split('\n'), bx+Inches(0.22), by+Inches(1.45),
         Inches(3.6), Inches(0.88), sz=12, col=RGBColor(0xCB,0xD5,0xE1),
         align=PP_ALIGN.CENTER)

# 底部摘要橫幅
rect(sl, Inches(0.45), Inches(4.82), Inches(12.4), Inches(0.72),
     fill=RGBColor(0x1E,0x3A,0x5C), rad=10000)
summary_items = [
    '18 個互動教學單元', '24 個（含 B 系列）', '輔助點擊首創', 'NT$0 完全免費', '完全離線'
]
summary_labels = ['F+C+A 系列', '含延伸 B 系列', '精細動作支持', '授權費用', '網路需求']
for i,(val,lbl) in enumerate(zip(summary_items, summary_labels)):
    sx = Inches(0.85+i*2.48)
    txt(sl, val, sx, Inches(4.87), Inches(2.3), Inches(0.35),
        sz=12.5, bold=True, col=AM, align=PP_ALIGN.CENTER)
    txt(sl, lbl, sx, Inches(5.22), Inches(2.3), Inches(0.28),
        sz=10, col=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)
    if i < 4:
        vline(sl, sx+Inches(2.28), Inches(4.88), Inches(5.48),
              col=RGBColor(0x33,0x4A,0x6E), pt=1)

txt(sl, '金錢小達人  ·  第八屆特殊教育教材教具競賽  ·  身心障礙類  ·  中華民國 115 年 3 月',
    Inches(0.45), Inches(5.72), Inches(12.4), Inches(0.38),
    sz=12, col=RGBColor(0x47,0x5A,0x73), align=PP_ALIGN.CENTER)
page_num(sl, 22)

# ═══════════════════════════════════════════════════════════════
#  儲存
# ═══════════════════════════════════════════════════════════════
out = 'doc/金錢小達人_競賽申請簡報_V3.pptx'
prs.save(out)
print(f'OK  {out}  slides={len(prs.slides)}')
