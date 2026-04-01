# -*- coding: utf-8 -*-
"""金錢小達人 競賽申請簡報 V5 — 教科書圖解・教育知識科普風格"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 色票 ────────────────────────────────────────────────────
W    = RGBColor(0xFF,0xFF,0xFF)
BG   = RGBColor(0xFA,0xF8,0xF4)   # 羊皮紙底
BG2  = RGBColor(0xF0,0xED,0xE6)
INK  = RGBColor(0x1A,0x18,0x14)   # 墨色主文
GR   = RGBColor(0x72,0x6B,0x60)   # 灰色說明
LN   = RGBColor(0xD8,0xD0,0xC4)   # 分隔線
NAV  = RGBColor(0x14,0x3D,0x6F)   # 海軍藍（主標題）
NAV2 = RGBColor(0x1C,0x54,0x96)   # 藍色輔助
NAVL = RGBColor(0xD6,0xE6,0xF7)
# 系列色
FC   = RGBColor(0x1A,0x68,0xB8)   # F-藍
FCL  = RGBColor(0xD6,0xE9,0xF8)
CC   = RGBColor(0x12,0x7A,0x52)   # C-綠
CCL  = RGBColor(0xD2,0xF0,0xE4)
AC   = RGBColor(0xC0,0x43,0x10)   # A-橙
ACL  = RGBColor(0xFD,0xE7,0xD5)
BC   = RGBColor(0x6E,0x28,0xB4)   # B-紫
BCL  = RGBColor(0xEE,0xDF,0xFB)
# 重點色
GO   = RGBColor(0xC8,0x80,0x00)   # 金色
GOL  = RGBColor(0xFD,0xF3,0xD6)
RD   = RGBColor(0xB8,0x28,0x28)
RDL  = RGBColor(0xFD,0xE8,0xE8)
TE   = RGBColor(0x0B,0x78,0x6E)
TEL  = RGBColor(0xCC,0xF5,0xF1)

SW = Inches(13.33)
SH = Inches(7.5)
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLK = prs.slide_layouts[6]

# ── 基礎工具 ────────────────────────────────────────────────
def R(sl,x,y,w,h,fill=None,bc=None,bw=0,rad=None):
    sh=sl.shapes.add_shape(1,x,y,w,h)
    sh.line.fill.background()
    if fill: sh.fill.solid();sh.fill.fore_color.rgb=fill
    else: sh.fill.background()
    if bc and bw>0: sh.line.color.rgb=bc;sh.line.width=Pt(bw)
    if rad is not None:
        sp=sh._element;spPr=sp.find(qn('p:spPr'))
        pg=spPr.find(qn('a:prstGeom'))
        if pg is not None:
            pg.set('prst','roundRect')
            av=pg.find(qn('a:avLst'))
            if av is None: av=etree.SubElement(pg,qn('a:avLst'))
            gd=etree.SubElement(av,qn('a:gd'))
            gd.set('name','adj');gd.set('fmla',f'val {rad}')
    return sh

def E(sl,x,y,w,h,fill=None,bc=None,bw=0):
    sh=sl.shapes.add_shape(9,x,y,w,h)
    sh.line.fill.background()
    if fill: sh.fill.solid();sh.fill.fore_color.rgb=fill
    else: sh.fill.background()
    if bc and bw>0: sh.line.color.rgb=bc;sh.line.width=Pt(bw)
    return sh

def T(sl,s,x,y,w,h,sz=12,bold=False,col=None,
      al=PP_ALIGN.LEFT,fn='Microsoft JhengHei',it=False):
    tb=sl.shapes.add_textbox(x,y,w,h)
    tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.alignment=al
    r=p.add_run();r.text=s
    r.font.name=fn;r.font.size=Pt(sz)
    r.font.bold=bold;r.font.italic=it
    if col: r.font.color.rgb=col
    return tb

def MT(sl,lines,x,y,w,h,sz=11,col=None,
       al=PP_ALIGN.LEFT,fn='Microsoft JhengHei',bold=False,sp=None):
    tb=sl.shapes.add_textbox(x,y,w,h)
    tf=tb.text_frame;tf.word_wrap=True
    for i,line in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=al
        if sp:
            from pptx.util import Pt as _P
            p.line_spacing=_P(sp)
        r=p.add_run();r.text=line
        r.font.name=fn;r.font.size=Pt(sz);r.font.bold=bold
        if col: r.font.color.rgb=col
    return tb

def HL(sl,y,x0=Inches(.4),x1=Inches(12.93),col=LN,pt=1.):
    ln=sl.shapes.add_connector(1,x0,y,x1,y)
    ln.line.color.rgb=col;ln.line.width=Pt(pt)

def VL(sl,x,y0,y1,col=LN,pt=1.):
    ln=sl.shapes.add_connector(1,x,y0,x,y1)
    ln.line.color.rgb=col;ln.line.width=Pt(pt)

def pill(sl,s,x,y,w,h,bg,fg,sz=10,bold=True):
    R(sl,x,y,w,h,fill=bg,rad=22000)
    T(sl,s,x,y,w,h,sz=sz,bold=bold,col=fg,al=PP_ALIGN.CENTER)

def num_badge(sl,n,cx,cy,r=Inches(.22),bg=NAV,fg=W):
    E(sl,cx-r,cy-r,r*2,r*2,fill=bg)
    T(sl,str(n),cx-r,cy-r*.8,r*2,r*1.6,sz=11,bold=True,col=fg,al=PP_ALIGN.CENTER)

def pn(sl,n,tot=22):
    T(sl,f'{n} / {tot}',SW-Inches(1.1),SH-Inches(.36),
      Inches(.9),Inches(.3),sz=9,col=GR,al=PP_ALIGN.RIGHT)

# ── 標題欄（教科書風：左色帶＋底線） ────────────────────────
def hbar(sl,title,sub='',accent=NAV,h=Inches(1.02)):
    R(sl,Inches(0),Inches(0),SW,h,fill=W)
    R(sl,Inches(0),Inches(0),Inches(.28),h,fill=accent)
    HL(sl,h,Inches(0),SW,col=accent,pt=2.2)
    T(sl,title,Inches(.48),Inches(.09),SW-Inches(5),Inches(.56),
      sz=22,bold=True,col=NAV)
    if sub:
        T(sl,sub,Inches(.48),Inches(.66),SW-Inches(2),Inches(.3),
          sz=10.5,col=GR,it=True)

# ── 教科書標注框 ─────────────────────────────────────────────
def callout(sl,s,x,y,w,h,bg=GOL,bc=GO):
    R(sl,x,y,w,h,fill=bg,bc=bc,bw=1.3,rad=6000)
    T(sl,s,x+Inches(.14),y+Inches(.1),w-Inches(.28),h-Inches(.2),
      sz=10.5,col=RGBColor(0x6B,0x3F,0x00))

# ── 章節側標 ────────────────────────────────────────────────
def sec(sl,s,x,y,col=NAV):
    R(sl,x,y,Inches(.07),Inches(.32),fill=col)
    T(sl,s,x+Inches(.15),y,Inches(4),Inches(.32),
      sz=11.5,bold=True,col=col)

# ══════════════════════════════════════════════════════════════
# S01  封面（學術期刊＋教科書書脊風格）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
# 頂部色條
R(sl,Inches(0),Inches(0),SW,Inches(.18),fill=NAV)
R(sl,Inches(0),Inches(.18),SW,Inches(.06),fill=GO)
# 左側書脊
R(sl,Inches(0),Inches(.24),Inches(4.1),SH-Inches(.24),fill=NAV)
for i,(c,y) in enumerate([
    (FC,Inches(1.6)),(CC,Inches(3.0)),(AC,Inches(4.4)),(BC,Inches(5.8))]):
    R(sl,Inches(0),y,Inches(4.1),Inches(1.25),fill=c)
    R(sl,Inches(0),y,Inches(4.1),Inches(.04),fill=W)
# 書脊文字
for txt_,y_ in [('F','1.72'),('C','3.12'),('A','4.52'),('B','5.92')]:
    T(sl,txt_,Inches(0),Inches(float(y_)),Inches(4.1),Inches(.9),
      sz=44,bold=True,col=RGBColor(0xFF,0xFF,0xFF,),al=PP_ALIGN.CENTER)
# 書脊說明
T(sl,'數學先備',Inches(.0),Inches(2.08),Inches(4.1),Inches(.38),
  sz=10,col=RGBColor(0xCC,0xE4,0xFF),al=PP_ALIGN.CENTER)
T(sl,'貨幣認知',Inches(.0),Inches(3.48),Inches(4.1),Inches(.38),
  sz=10,col=RGBColor(0xCC,0xF2,0xE2),al=PP_ALIGN.CENTER)
T(sl,'情境類化',Inches(.0),Inches(4.88),Inches(4.1),Inches(.38),
  sz=10,col=RGBColor(0xFF,0xE0,0xCC),al=PP_ALIGN.CENTER)
T(sl,'預算規劃',Inches(.0),Inches(6.28),Inches(4.1),Inches(.38),
  sz=10,col=RGBColor(0xE8,0xD8,0xFF),al=PP_ALIGN.CENTER)
# 右側主體
pill(sl,'第八屆教育部獎助研發特殊教育教材教具競賽　·　身心障礙類　·　國民中學及高級中等學校組',
     Inches(4.35),Inches(.38),Inches(8.7),Inches(.42),
     bg=NAVL,fg=NAV,sz=10.5)
T(sl,'金錢小達人',Inches(4.35),Inches(1.02),Inches(8.7),Inches(1.1),
  sz=58,bold=True,col=NAV)
T(sl,'互動式特殊教育生活數學暨金融素養教學軟體',
  Inches(4.35),Inches(2.1),Inches(8.7),Inches(.52),
  sz=17,bold=True,col=AC)
T(sl,'Money Tutor: Interactive Software for Financial Literacy in Special Education',
  Inches(4.35),Inches(2.62),Inches(8.7),Inches(.34),
  sz=10.5,col=GR,it=True,fn='Arial')
HL(sl,Inches(3.08),Inches(4.35),Inches(12.95),col=LN,pt=1.5)
# 統計徽章
for i,(v,l,c) in enumerate([
    ('24','教學單元',FC),('4 系列','課程序階',CC),
    ('150,000+','程式碼行數',AC),('NT$ 0','授權費用',BC),('完全離線','免安裝',TE)]):
    bx=Inches(4.35+i*1.74)
    T(sl,v,bx,Inches(3.18),Inches(1.7),Inches(.46),
      sz=17,bold=True,col=c,al=PP_ALIGN.CENTER)
    T(sl,l,bx,Inches(3.62),Inches(1.7),Inches(.26),
      sz=9,col=GR,al=PP_ALIGN.CENTER)
HL(sl,Inches(4.0),Inches(4.35),Inches(12.95),col=LN,pt=1.)
T(sl,'🎯 適用對象　國中～高中特教班・資源班',
  Inches(4.35),Inches(4.1),Inches(8.7),Inches(.34),
  sz=11.5,bold=True,col=NAV)
T(sl,'輕度至中度智能障礙・自閉症・學習障礙・多重障礙',
  Inches(4.35),Inches(4.44),Inches(8.7),Inches(.3),sz=11,col=INK)
callout(sl,'✦ 科技輔助教學軟體（HTML5）・免安裝・完全離線・授權費 0 元・原始碼完整開放',
        Inches(4.35),Inches(4.9),Inches(8.7),Inches(.42))
T(sl,'中華民國 115 年 3 月',Inches(4.35),Inches(5.55),Inches(4),Inches(.3),
  sz=10,col=GR)
pn(sl,1)

# ══════════════════════════════════════════════════════════════
# S02  目錄
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'簡報綱要','Table of Contents',accent=NAV)
toc=[
    (RD,'01','設計動機','特教現場三重結構性困境'),
    (AC,'02','設計理念','UDL × 工作分析 × 直接教學 × 融合三層'),
    (FC,'03','F 系列','數學基礎先備 6 單元'),
    (CC,'04','C 系列','貨幣認知建立 6 單元'),
    (AC,'05','A 系列','真實情境類化 6 單元'),
    (BC,'06','B 系列','預算規劃延伸 6 單元'),
    (TE,'07','輔助點擊','精細動作障礙首創解方'),
    (GO,'08','三難度系統','簡單 / 普通 / 困難完整比較'),
    (NAV,'09','工作分析法','步驟化教學流程圖解'),
    (CC,'10','融合教育','三層支持架構圖解'),
    (NAV,'11','IEP 對應','課綱法規完整對應表'),
    (RD,'12','試教觀察','三個案・類化窗口記錄'),
    (GO,'13','預期成效','學生・教師・學校三層面'),
    (FC,'14','技術規格','無障礙・作業單・獎勵系統'),
    (BC,'15','核心創新','四大首創特色・結語'),
]
for idx,(c,num,t,s) in enumerate(toc):
    col=idx%3; row=idx//3
    bx=Inches(.32+col*4.35)
    by=Inches(1.12+row*1.23)
    R(sl,bx,by,Inches(4.12),Inches(1.1),fill=W,bc=LN,bw=1.1,rad=7000)
    R(sl,bx,by,Inches(.09),Inches(1.1),fill=c)
    T(sl,num,bx+Inches(.18),by+Inches(.1),Inches(.65),Inches(.44),
      sz=20,bold=True,col=c,fn='Arial')
    T(sl,t,bx+Inches(.9),by+Inches(.1),Inches(3.1),Inches(.38),
      sz=13,bold=True,col=INK)
    T(sl,s,bx+Inches(.9),by+Inches(.52),Inches(3.1),Inches(.44),
      sz=10,col=GR)
pn(sl,2)

# ══════════════════════════════════════════════════════════════
# S03  設計動機——三重困境（時間軸個案卡）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'設計動機——特教現場三重結構性困境','壹、來自十五年特教現場的教材研發回應',accent=RD)
# 橫軸時間線
HL(sl,Inches(3.8),Inches(.6),Inches(12.73),col=RD,pt=2.5)
for i,xp in enumerate([Inches(1.45),Inches(5.2),Inches(9.0)]):
    E(sl,xp-Inches(.18),Inches(3.62),Inches(.36),Inches(.36),fill=RD)

cases=[
    (RD,RDL,'困境 1','交易序列\n意義理解障礙',
     '小涵（14歲，輕度智能障礙）',
     '能算 100−35=65\n卻無法連結「我的錢」\n與「這個多少錢」',
     '靜態計算≠動態交易\n序列意義理解障礙'),
    (AC,ACL,'困境 2','安全練習\n場域長期缺位',
     '小俊（輕度自閉症）',
     '速食店選擇迴路→情緒失調\n此後一學期拒絕\n接觸相關學習',
     '真實環境試誤風險過高\n需要「失敗零後果」\n的安全練習場'),
    (BC,BCL,'困境 3','系統性課程\n工具十年缺口',
     '課綱 vs 市場現況',
     '生管-Ⅳ-E-1/F-1/E-4\n已明列財務管理主題\n市場無對應互動軟體',
     '課綱要求明確\n系統性數位工具\n缺口長達十年'),
]
for i,(hc,lc,tag,title,case,prob,key) in enumerate(cases):
    bx=Inches(.32+i*4.35)
    # 上半（問題）
    R(sl,bx,Inches(1.18),Inches(4.12),Inches(2.44),
      fill=W,bc=hc,bw=1.6,rad=10000)
    R(sl,bx,Inches(1.18),Inches(4.12),Inches(.48),fill=hc,rad=0)
    T(sl,tag,bx+Inches(.15),Inches(1.24),Inches(1.2),Inches(.38),
      sz=12,bold=True,col=W)
    T(sl,title,bx+Inches(.15),Inches(1.74),Inches(3.82),Inches(.66),
      sz=13,bold=True,col=hc)
    pill(sl,case,bx+Inches(.15),Inches(2.46),Inches(3.82),Inches(.28),
         bg=lc,fg=hc,sz=9.5,bold=False)
    MT(sl,prob.split('\n'),bx+Inches(.2),Inches(2.82),
       Inches(3.72),Inches(.72),sz=10.5,col=INK,sp=15)
    # 下半（關鍵詞）
    R(sl,bx,Inches(3.98),Inches(4.12),Inches(1.0),
      fill=lc,bc=hc,bw=1.2,rad=8000)
    T(sl,'▶ '+key.replace('\n',' '),bx+Inches(.2),Inches(4.1),
      Inches(3.72),Inches(.76),sz=10.5,col=hc)
# 底部統整
callout(sl,'✦ 三重困境的共同本質：特殊需求學生需要「步驟化引導」「安全試誤」「系統化工具」三位一體的解方',
        Inches(.32),Inches(5.18),Inches(12.7),Inches(.42))
# 解方標示
R(sl,Inches(.32),Inches(5.78),Inches(12.7),Inches(.9),
  fill=NAVL,bc=NAV,bw=1.5,rad=8000)
T(sl,'→ 解方：金錢小達人',Inches(.55),Inches(5.85),Inches(3.5),Inches(.38),
  sz=13,bold=True,col=NAV)
for i,(ic,t) in enumerate([
    ('🔢','步驟化工作分析'),('🛡️','安全數位練習環境'),
    ('📚','24單元系統課程'),('🖱️','輔助點擊無障礙')]):
    T(sl,ic+' '+t,Inches(4.0+i*2.2),Inches(5.88),Inches(2.1),Inches(.32),
      sz=11,col=NAV)
pn(sl,3)

# ══════════════════════════════════════════════════════════════
# S04  設計理念——四大教育理論（巢狀結構＋四柱）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'設計理念——四大教育理論支柱','貳、UDL × 工作分析法 × 直接教學法 × 融合教育三層支持',accent=AC)

pillars=[
    (FC,FCL,'🎯','UDL\n通用設計學習',
     ['多元表徵（What）','多元行動與表達（How）','多元參與方式（Why）'],
     'CAST 2018'),
    (CC,CCL,'🔧','工作分析法\n步驟化教學',
     ['拆解設備操作序列','每步可單獨教學評量','對應 IEP 行為目標'],
     'Gold 1980 / Snell & Brown 2011'),
    (AC,ACL,'📈','直接教學法\n三階段難度',
     ['I 示範（簡單模式）','II 引導（普通模式）','III 獨立（困難模式）'],
     'Engelmann & Carnine 1982'),
    (BC,BCL,'🤝','融合教育\n三層支持',
     ['L1 普通班同步參與','L2 資源班強化練習','L3 自足班精熟學習'],
     '特殊教育法 §28・LRE 原則'),
]
for i,(hc,lc,ic,title,pts,ref) in enumerate(pillars):
    bx=Inches(.42+i*3.25)
    cw=Inches(3.0)
    # 頂部icon圓
    E(sl,bx+Inches(1.0),Inches(1.15),Inches(1.0),Inches(1.0),fill=hc)
    T(sl,ic,bx+Inches(1.0),Inches(1.18),Inches(1.0),Inches(.9),
      sz=26,al=PP_ALIGN.CENTER)
    # 柱體
    R(sl,bx,Inches(2.2),cw,Inches(4.7),fill=lc,bc=hc,bw=1.8,rad=10000)
    R(sl,bx,Inches(2.2),cw,Inches(.5),fill=hc,rad=0)
    T(sl,title,bx,Inches(2.24),cw,Inches(.44),
      sz=12,bold=True,col=W,al=PP_ALIGN.CENTER)
    HL(sl,Inches(2.78),bx+Inches(.2),bx+cw-Inches(.2),col=hc,pt=.9)
    MT(sl,['▸ '+p for p in pts],bx+Inches(.18),Inches(2.86),
       cw-Inches(.3),Inches(2.5),sz=11,col=INK,sp=17)
    callout(sl,ref,bx+Inches(.14),Inches(5.5),cw-Inches(.28),Inches(.48))
# 底部整合說明
R(sl,Inches(.42),Inches(6.1),Inches(12.5),Inches(.52),
  fill=NAVL,bc=NAV,bw=1.3,rad=7000)
T(sl,'🔗 整合原則：相同教材・相同目標方向，差異在「支持程度（Level of Support）」，實現最少限制環境（LRE）',
  Inches(.65),Inches(6.18),Inches(12.1),Inches(.38),sz=11,col=NAV)
pn(sl,4)

# ══════════════════════════════════════════════════════════════
# S05  課程架構全覽（斜角階梯圖）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'課程架構全覽——四系列 24 單元完整序階',
     '三層遞進：數學基礎 → 貨幣認知 → 真實情境 → 預算規劃延伸',accent=FC)

# 四層階梯（由下往上，左縮排遞增）
stairs=[
    (BC,BCL,'B','預算規劃','B1–B6',
     '今日帶多少錢・零用錢日記・存錢計畫・特賣比一比・生日派對預算・菜市場買菜',
     Inches(3.0),Inches(1.18)),
    (AC,ACL,'A','情境類化','A1–A6',
     '自動販賣機・理髮廳・麥當勞・超市購物・ATM・火車票',
     Inches(2.0),Inches(2.38)),
    (CC,CCL,'C','貨幣認知','C1–C6',
     '認識錢幣・數錢・換錢・付款・夠不夠・找零',
     Inches(1.0),Inches(3.58)),
    (FC,FCL,'F','數學先備','F1–F6',
     '一對一對應・唱數・數字認讀・數字排序・量比較・數的組成',
     Inches(0),Inches(4.78)),
]
for hc,lc,ser,name,units,detail,indent,sy in stairs:
    sw_=SW-Inches(.27)-indent
    R(sl,Inches(.27)+indent,sy,sw_,Inches(1.0),fill=lc,bc=hc,bw=1.8,rad=7000)
    R(sl,Inches(.27)+indent,sy,Inches(1.5),Inches(1.0),fill=hc,rad=0)
    T(sl,ser,Inches(.27)+indent,sy+Inches(.1),Inches(1.5),Inches(.44),
      sz=18,bold=True,col=W,al=PP_ALIGN.CENTER)
    T(sl,name,Inches(.27)+indent+Inches(1.65),sy+Inches(.08),
      Inches(2.2),Inches(.38),sz=12,bold=True,col=hc)
    pill(sl,units,Inches(.27)+indent+Inches(4.0),sy+Inches(.2),
         Inches(1.0),Inches(.3),bg=hc,fg=W,sz=9.5)
    T(sl,detail,Inches(.27)+indent+Inches(5.2),sy+Inches(.12),
      sw_-Inches(5.4),Inches(.36),sz=10,col=INK)
# 學習方向箭頭
T(sl,'學習方向\n↑',Inches(0),Inches(1.5),Inches(.28),Inches(3.7),
  sz=9,bold=True,col=GR,al=PP_ALIGN.CENTER)
# 適用年段
R(sl,Inches(.27),Inches(5.98),Inches(12.68),Inches(.5),
  fill=W,bc=LN,bw=1.,rad=5000)
for lx,lt,lc_ in [
    (Inches(.5),'國小中高年級',FC),(Inches(3.5),'國中特教班',CC),
    (Inches(6.5),'高中特教班',AC),(Inches(9.5),'延伸應用',BC)]:
    R(sl,lx,Inches(6.06),Inches(.14),Inches(.26),fill=lc_,rad=3000)
    T(sl,lt,lx+Inches(.22),Inches(6.06),Inches(2.5),Inches(.26),
      sz=11,col=lc_,bold=True)
callout(sl,'✦ 每系列 6 個互動單元，共 24 個教學單元；各系列可獨立使用，亦可完整序階進行',
        Inches(.27),Inches(6.62),Inches(12.68),Inches(.38))
pn(sl,5)

# ══════════════════════════════════════════════════════════════
# S06  F 系列——六宮格＋技能進程
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'F 系列——數學基礎先備技能模組',
     '適用：國小中高年級至國中　目標：建立貨幣操作所需數學先備技能',accent=FC)

units_f=[
    ('F1','一對一對應','🔢',
     '透過拖曳物件配對建立數量等量直觀理解\n物件與容器配對・數量等量概念建立'),
    ('F2','唱數','🔊',
     '將數量概念轉化為有序語言序列\n視覺化數列・跟讀模式・獨立唱數'),
    ('F3','數字認讀','👁️',
     '辨識阿拉伯數字及其對應量\n10種自訂主題圖示・正反辨識'),
    ('F4','數字排序','🔀',
     '理解數字大小關係並進行排列操作\n由小到大・由大到小雙向練習'),
    ('F5','量比較','⚖️',
     '判斷兩組數量的多少大小關係\n視覺化比較・語言表達・符號理解'),
    ('F6','數的組成','🧩',
     '理解數字可分解與組合的內部結構\n分解遊戲・加法橋接・組成圖示'),
]
for i,(code,name,ic,desc) in enumerate(units_f):
    col=i%3; row=i//3
    bx=Inches(.35+col*4.35)
    by=Inches(1.22+row*2.88)
    R(sl,bx,by,Inches(4.12),Inches(2.62),fill=W,bc=FC,bw=1.6,rad=10000)
    R(sl,bx,by,Inches(4.12),Inches(.46),fill=FC,rad=0)
    T(sl,f'{code}  {name}',bx+Inches(.15),by+Inches(.06),
      Inches(2.9),Inches(.36),sz=14,bold=True,col=W)
    T(sl,ic,bx+Inches(3.5),by+Inches(.02),Inches(.55),Inches(.42),
      sz=20,al=PP_ALIGN.CENTER)
    MT(sl,desc.split('\n'),bx+Inches(.18),by+Inches(.62),
       Inches(3.76),Inches(1.8),sz=11,col=INK,sp=16)
# 技能進程箭頭
T(sl,'F1 → F2 → F3 → F4 → F5 → F6　（建議學習順序）',
  Inches(.35),Inches(7.08),Inches(10),Inches(.32),sz=10.5,col=FC,bold=True)
pn(sl,6)

# ══════════════════════════════════════════════════════════════
# S07  C 系列
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'C 系列——貨幣認知建立模組',
     '適用：國中至高中　目標：建立完整的金錢認知與交易技能',accent=CC)

units_c=[
    ('C1','認識錢幣','🪙',
     '台灣9種面額辨識\n硬幣正反面・紙鈔防偽標誌・聽音辨幣'),
    ('C2','數錢','🔢',
     '多面額混合計算・逐枚累加程序\n視覺化計數・語音同步・面額分組'),
    ('C3','換錢','🔄',
     '面額兌換雙向練習\n小換大・大換小・隨機組合・全隨機模式'),
    ('C4','付款','💳',
     '湊出指定金額操作\n面額最佳組合・拖曳付款・提示鈕輔助'),
    ('C5','夠不夠','⚖️',
     '判斷錢包金額是否足夠購買\n比較思維・差額計算・視覺化呈現'),
    ('C6','找零','💰',
     '計算找零並驗收的完整交易閉環\n三選一題型・算式提示・精準付款'),
]
for i,(code,name,ic,desc) in enumerate(units_c):
    col=i%3; row=i//3
    bx=Inches(.35+col*4.35); by=Inches(1.22+row*2.88)
    R(sl,bx,by,Inches(4.12),Inches(2.62),fill=W,bc=CC,bw=1.6,rad=10000)
    R(sl,bx,by,Inches(4.12),Inches(.46),fill=CC,rad=0)
    T(sl,f'{code}  {name}',bx+Inches(.15),by+Inches(.06),
      Inches(2.9),Inches(.36),sz=14,bold=True,col=W)
    T(sl,ic,bx+Inches(3.5),by+Inches(.02),Inches(.55),Inches(.42),
      sz=20,al=PP_ALIGN.CENTER)
    MT(sl,desc.split('\n'),bx+Inches(.18),by+Inches(.62),
       Inches(3.76),Inches(1.8),sz=11,col=INK,sp=16)
callout(sl,'📌 完成 C 系列建議進入 A 系列：將靜態貨幣技能類化至真實社區情境設備操作',
        Inches(.35),Inches(7.08),Inches(12.6),Inches(.38))
pn(sl,7)

# ══════════════════════════════════════════════════════════════
# S08  A 系列
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'A 系列——真實情境類化模組',
     '適用：國中至高中　目標：社區設備獨立操作・類化至真實生活',accent=AC)

units_a=[
    ('A1','自動販賣機','🥤',
     '7步驟完整操作序列\n投幣・選飲料・確認・取出・找零驗收'),
    ('A2','理髮廳','✂️',
     '觸控繳費機服務選擇付款\n先投幣再選服務・先選服務再付款雙模式'),
    ('A3','麥當勞','🍔',
     '自助點餐機完整流程\n餐點分類選擇・組合・付款・取餐'),
    ('A4','超市購物','🛒',
     '收銀台購物結帳全流程\n12種商店類型・商品掃描・付款・找零'),
    ('A5','ATM 操作','🏧',
     '提款・存款・轉帳三大功能\n密碼輸入・金額選擇・收據取卡'),
    ('A6','火車票','🚃',
     '自動售票機購票全流程\n車站選擇・張數・付款・取票・找零'),
]
for i,(code,name,ic,desc) in enumerate(units_a):
    col=i%3; row=i//3
    bx=Inches(.35+col*4.35); by=Inches(1.22+row*2.88)
    R(sl,bx,by,Inches(4.12),Inches(2.62),fill=W,bc=AC,bw=1.6,rad=10000)
    R(sl,bx,by,Inches(4.12),Inches(.46),fill=AC,rad=0)
    T(sl,f'{code}  {name}',bx+Inches(.15),by+Inches(.06),
      Inches(2.9),Inches(.36),sz=14,bold=True,col=W)
    T(sl,ic,bx+Inches(3.5),by+Inches(.02),Inches(.55),Inches(.42),
      sz=20,al=PP_ALIGN.CENTER)
    MT(sl,desc.split('\n'),bx+Inches(.18),by+Inches(.62),
       Inches(3.76),Inches(1.8),sz=11,col=INK,sp=16)
callout(sl,'🎯 終極目標：教師 3 公尺外靜態觀察，學生在真實社區環境完全獨立操作',
        Inches(.35),Inches(7.08),Inches(12.6),Inches(.38))
pn(sl,8)

# ══════════════════════════════════════════════════════════════
# S09  B 系列
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'B 系列——預算規劃延伸模組',
     '適用：國中至高中　目標：日常預算管理與消費決策能力',accent=BC)

units_b=[
    ('B1','今天帶多少錢','👜',
     '依行程規劃所需金額\n場景類別篩選・費用明細・最少張數提示'),
    ('B2','零用錢日記','📒',
     '記帳概念建立與收支追蹤\n收入/支出事件・餘額計算・日記主題篩選'),
    ('B3','存錢計畫','🐷',
     '月曆存錢模擬（拖曳存幣）\n目標設定・週存金額規劃・里程碑慶祝'),
    ('B4','特賣比一比','🏷️',
     '商品比價與消費決策\n兩商店・三商店・單位比價三種模式'),
    ('B5','生日派對預算','🎂',
     '必買/選購 項目預算分配\n三關制・派對主題篩選・超支智慧提示'),
    ('B6','菜市場買菜','🥦',
     '多攤位採購清單執行\n三種市場類型・付款・找零三選一'),
]
for i,(code,name,ic,desc) in enumerate(units_b):
    col=i%3; row=i//3
    bx=Inches(.35+col*4.35); by=Inches(1.22+row*2.88)
    R(sl,bx,by,Inches(4.12),Inches(2.62),fill=W,bc=BC,bw=1.6,rad=10000)
    R(sl,bx,by,Inches(4.12),Inches(.46),fill=BC,rad=0)
    T(sl,f'{code}  {name}',bx+Inches(.15),by+Inches(.06),
      Inches(2.9),Inches(.36),sz=14,bold=True,col=W)
    T(sl,ic,bx+Inches(3.5),by+Inches(.02),Inches(.55),Inches(.42),
      sz=20,al=PP_ALIGN.CENTER)
    MT(sl,desc.split('\n'),bx+Inches(.18),by+Inches(.62),
       Inches(3.76),Inches(1.8),sz=11,col=INK,sp=16)
pn(sl,9)

# ══════════════════════════════════════════════════════════════
# S10  輔助點擊模式（T型對比＋流程圖）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'輔助點擊模式——精細動作障礙首創解方',
     '全 24 單元支援・僅需「單一點擊」即可完成所有操作',accent=TE)

# 左：傳統障礙
sec(sl,'❌ 傳統操作障礙',Inches(.4),Inches(1.22),col=RD)
R(sl,Inches(.4),Inches(1.6),Inches(4.2),Inches(2.5),
  fill=RDL,bc=RD,bw=1.5,rad=8000)
MT(sl,['✗  拖曳操作需精確手部控制',
       '✗  觸控滑動對肢障學生困難',
       '✗  多步驟序列記憶負擔沉重',
       '✗  連續失敗誘發情緒挫折',
       '✗  教師須全程逐步介入協助'],
   Inches(.6),Inches(1.75),Inches(3.8),Inches(2.1),
   sz=11.5,col=RD,sp=17)

# 中間箭頭
T(sl,'解方\n→',Inches(4.75),Inches(2.25),Inches(.9),Inches(.9),
  sz=20,bold=True,col=TE,al=PP_ALIGN.CENTER)

# 右：解方
sec(sl,'✅ 輔助點擊解方',Inches(5.75),Inches(1.22),col=TE)
R(sl,Inches(5.75),Inches(1.6),Inches(4.2),Inches(2.5),
  fill=TEL,bc=TE,bw=1.5,rad=8000)
MT(sl,['✓  系統自動偵測點擊位置',
       '✓  黃色邊框視覺提示下一步',
       '✓  自動引導至正確操作目標',
       '✓  全程只需「點擊」不需拖曳',
       '✓  可疊加於任何難度層級'],
   Inches(5.95),Inches(1.75),Inches(3.8),Inches(2.1),
   sz=11.5,col=TE,sp=17)

# 適用族群
sec(sl,'適用對象',Inches(10.1),Inches(1.22),col=BC)
R(sl,Inches(10.1),Inches(1.6),Inches(3.0),Inches(2.5),
  fill=BCL,bc=BC,bw=1.5,rad=8000)
MT(sl,['• 腦性麻痺','• 脊髓損傷',
       '• 肌肉萎縮症','• 注意力缺失過動',
       '• 其他精細動作困難'],
   Inches(10.3),Inches(1.75),Inches(2.6),Inches(2.1),
   sz=11.5,col=BC,sp=17)

# 操作流程
sec(sl,'操作流程圖',Inches(.4),Inches(4.32),col=TE)
steps_=['①啟用\n輔助點擊','②出現\n黃色提示框','③點擊\n任意處',
        '④自動執行\n下一步','⑤完成\n學習任務']
for i,s in enumerate(steps_):
    bx=Inches(.4+i*2.5)
    E(sl,bx,Inches(4.72),Inches(.5),Inches(.5),fill=TE)
    T(sl,str(i+1),bx,Inches(4.72),Inches(.5),Inches(.5),
      sz=14,bold=True,col=W,al=PP_ALIGN.CENTER)
    R(sl,bx-Inches(.32),Inches(5.3),Inches(1.14),Inches(.78),
      fill=TEL,bc=TE,bw=1.2,rad=7000)
    T(sl,s,bx-Inches(.32),Inches(5.34),Inches(1.14),Inches(.7),
      sz=9.5,col=TE,al=PP_ALIGN.CENTER)
    if i<4:
        T(sl,'→',Inches(1.28+i*2.5),Inches(4.82),Inches(.4),Inches(.3),
          sz=14,col=GR,al=PP_ALIGN.CENTER)

callout(sl,'✦ 首創設計：全台同類型特教軟體中，首個支援「全程單一點擊」輔助模式',
        Inches(.4),Inches(6.28),Inches(12.53),Inches(.42))
pn(sl,10)

# ══════════════════════════════════════════════════════════════
# S11  三難度系統（對齊矩陣比較表）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'三段難度設計——直接教學法系統落實',
     '每個單元均提供三種難度，確保從示範到獨立操作的完整進展',accent=GO)

hdrs=['','😊 簡單模式\nEasy','😐 普通模式\nNormal','😤 困難模式\nHard']
hcols=[None,CC,GO,RD]
cws=[Inches(2.1),Inches(3.4),Inches(3.4),Inches(3.4)]
bx0=Inches(.32); by0=Inches(1.18); rh=Inches(.42)

# 表頭
cx=bx0
for h,hc,cw in zip(hdrs,hcols,cws):
    if hc:
        R(sl,cx,by0,cw,Inches(.72),fill=hc,rad=0)
        T(sl,h,cx,by0+Inches(.06),cw,Inches(.62),
          sz=13,bold=True,col=W,al=PP_ALIGN.CENTER)
    cx+=cw

rows_data=[
    ('教學策略',
     '完整步驟提示\n視覺指引箭頭\n自動輔助確認',
     '提示鈕按需提示\n3次錯誤自動提示\n減少自動引導',
     '僅語音引導一次\n無視覺提示框\n完全靠記憶操作'),
    ('適用學生',
     '初學者\n需要高度支持\n建立程序記憶',
     '中級學習者\n部分提示需求\n強化記憶階段',
     '進階學習者\n準備類化現實\n類化前最終確認'),
    ('提示類型',
     '黃框高亮＋語音\n步驟自動播報\n輔助點擊可用',
     '錯誤時紅色提示\n按鈕觸發提示\n語音選擇性',
     '語音僅提示一次\n無自動輔助\n錯誤需重嘗試'),
    ('IEP 對應',
     'FP 全程協助\nPP 部分肢體\nVP 口語提示',
     'VP 口語提示\nM 示範提示\n（按需觸發）',
     'I 獨立完成\n類化準備指標\n社區轉銜基準'),
]
row_bgs=[W,BG2,W,BG2]
for ri,(label,*cells) in enumerate(rows_data):
    ry=by0+Inches(.72)+ri*Inches(1.26)
    cx=bx0
    # label
    R(sl,cx,ry,cws[0],Inches(1.26),fill=NAVL,bc=NAV,bw=.8)
    T(sl,label,cx+Inches(.12),ry+Inches(.1),cws[0]-Inches(.24),Inches(.5),
      sz=12,bold=True,col=NAV,al=PP_ALIGN.CENTER)
    cx+=cws[0]
    for ci,(cell,hc) in enumerate(zip(cells,hcols[1:])):
        R(sl,cx,ry,cws[ci+1],Inches(1.26),fill=row_bgs[ri],bc=LN,bw=.8)
        MT(sl,cell.split('\n'),cx+Inches(.12),ry+Inches(.1),
           cws[ci+1]-Inches(.2),Inches(1.04),sz=11,col=INK,sp=15)
        cx+=cws[ci+1]

callout(sl,'＋ 輔助點擊模式可疊加於任何難度，IEP對應：簡單→VP/PP/FP　普通→VP/M　困難→I（獨立）',
        Inches(.32),Inches(6.46),Inches(12.7),Inches(.42))
pn(sl,11)

# ══════════════════════════════════════════════════════════════
# S12  工作分析法（水平步驟流程圖）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'工作分析法——步驟化教學流程圖解',
     '以 A1 自動販賣機為例　Gold 1980 / Snell & Brown 2011',accent=CC)

steps7=[
    (FC,'Step 1','確認攜帶金額','開啟錢包\n核對面額'),
    (FC,'Step 2','選擇目標飲料','瀏覽選單\n確認售價'),
    (CC,'Step 3','投入硬幣/紙鈔','面額選擇\n投入確認'),
    (CC,'Step 4','確認金額足夠','螢幕顯示\n比對售價'),
    (AC,'Step 5','按下購買鍵','確認按鈕\n等待出貨'),
    (AC,'Step 6','取出飲料','出貨口取出\n確認品項'),
    (BC,'Step 7','確認找零','找零口取出\n核對金額'),
]
# 上排 4步
for i,(hc,st,nm,desc) in enumerate(steps7[:4]):
    bx=Inches(.32+i*3.26)
    R(sl,bx,Inches(1.22),Inches(3.0),Inches(1.9),
      fill=W,bc=hc,bw=1.6,rad=8000)
    E(sl,bx+Inches(.08),Inches(1.06),Inches(.44),Inches(.44),fill=hc)
    T(sl,st,bx+Inches(.08),Inches(1.06),Inches(.44),Inches(.4),
      sz=10,bold=True,col=W,al=PP_ALIGN.CENTER)
    T(sl,nm,bx+Inches(.14),Inches(1.3),Inches(2.72),Inches(.36),
      sz=12,bold=True,col=hc)
    MT(sl,desc.split('\n'),bx+Inches(.14),Inches(1.72),
       Inches(2.72),Inches(1.2),sz=11,col=INK,sp=15)
    if i<3:
        T(sl,'→',Inches(3.45+i*3.26),Inches(1.85),Inches(.4),Inches(.36),
          sz=16,col=GR,al=PP_ALIGN.CENTER)

# 下排 3步（從右往左掉頭）
for i,(hc,st,nm,desc) in enumerate(steps7[4:]):
    bx=Inches(.32+(2-i)*3.26)
    R(sl,bx,Inches(3.28),Inches(3.0),Inches(1.9),
      fill=W,bc=hc,bw=1.6,rad=8000)
    E(sl,bx+Inches(.08),Inches(3.12),Inches(.44),Inches(.44),fill=hc)
    T(sl,st,bx+Inches(.08),Inches(3.12),Inches(.44),Inches(.4),
      sz=10,bold=True,col=W,al=PP_ALIGN.CENTER)
    T(sl,nm,bx+Inches(.14),Inches(3.36),Inches(2.72),Inches(.36),
      sz=12,bold=True,col=hc)
    MT(sl,desc.split('\n'),bx+Inches(.14),Inches(3.78),
       Inches(2.72),Inches(1.2),sz=11,col=INK,sp=15)
    if i<2:
        T(sl,'←',Inches(3.45+(1-i)*3.26),Inches(3.9),Inches(.4),Inches(.36),
          sz=16,col=GR,al=PP_ALIGN.CENTER)

# 評量說明欄
R(sl,Inches(.32),Inches(5.38),Inches(12.7),Inches(.42),fill=NAVL,bc=NAV,bw=1.)
T(sl,'評量標準：  I 獨立完成　｜　VP 口語提示（≤2次）　｜　PP 部分肢體提示　｜　FP 全程協助',
  Inches(.5),Inches(5.44),Inches(12.3),Inches(.32),sz=11.5,col=NAV)
callout(sl,'✦ 每步驟均可獨立評量、單獨教學，完整對應 IEP 短期目標行為敘述格式（行為+條件+標準）',
        Inches(.32),Inches(5.96),Inches(12.7),Inches(.42))
pn(sl,12)

# ══════════════════════════════════════════════════════════════
# S13  融合教育三層支持（倒三角梯形）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'融合教育三層支持架構',
     '相同教材・相同目標方向，差異在「支持程度（Level of Support）」',accent=CC)

# 三層（由上到下寬度遞增）
tiers=[
    (FC,FCL,'Level 1  普通班同步參與路徑',
     Inches(.5),Inches(1.18),Inches(12.33),Inches(1.1),
     ['• 以輔助點擊模式在普通班數學課中參與相同任務',
      '• 軟體自動提供支持，教師不須全程個別介入，實現無縫融合',
      '• 普通班同儕理解如何給予適當協助，雙向融合目標同步實現']),
    (CC,CCL,'Level 2  資源班強化練習路徑',
     Inches(1.6),Inches(2.46),Inches(10.13),Inches(1.1),
     ['• 資源班抽離課使用，針對個別化弱點進行系統性強化練習',
      '• 可在結構化環境中重複操作，不受普通班進度影響',
      '• 教師即時調整難度，精確對應每位學生的 IEP 短期目標']),
    (AC,ACL,'Level 3  自足班精熟學習路徑',
     Inches(2.7),Inches(3.74),Inches(7.93),Inches(1.1),
     ['• 自足式特教班完整課程使用，依學生程度客製化學習路徑',
      '• 作業單輔助評量，系統化記錄進步，提供 IEP 具體數據',
      '• 類化至社區環境前的最後精熟確認階段']),
]
for hc,lc,title,lx,ly,lw,lh,pts in tiers:
    R(sl,lx,ly,lw,lh,fill=lc,bc=hc,bw=1.8,rad=7000)
    R(sl,lx,ly,Inches(.12),lh,fill=hc)
    T(sl,title,lx+Inches(.24),ly+Inches(.08),lw-Inches(.4),Inches(.34),
      sz=12,bold=True,col=hc)
    MT(sl,pts,lx+Inches(.24),ly+Inches(.44),lw-Inches(.4),Inches(.6),
       sz=10.5,col=INK,sp=14)

# 右側說明
sec(sl,'法規依據',Inches(.5),Inches(5.06),col=NAV)
items_=[
    (NAV,'特殊教育法 §28','依個別需求訂定 IEP・相同教材差異化支持'),
    (CC,'最少限制環境（LRE）','從結構化到自然情境逐步褪除支持'),
    (FC,'CRPD 第 24 條','融合教育・平等受教機會保障'),
]
for i,(c,t,d) in enumerate(items_):
    bx=Inches(.5+i*4.25)
    R(sl,bx,Inches(5.44),Inches(4.0),Inches(.8),fill=W,bc=c,bw=1.2,rad=6000)
    pill(sl,t,bx+Inches(.1),Inches(5.52),Inches(3.8),Inches(.24),bg=c,fg=W,sz=9.5)
    T(sl,d,bx+Inches(.14),Inches(5.8),Inches(3.72),Inches(.28),sz=10,col=INK)
callout(sl,'✦ 體現《特殊教育法》第 28 條精神：相同教材・相同目標方向，差異在支持程度，而非課程內容的全然切割',
        Inches(.5),Inches(6.44),Inches(12.33),Inches(.42))
pn(sl,13)

# ══════════════════════════════════════════════════════════════
# S14  IEP / 課綱 / 法規三欄對應
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'IEP 目標・課綱・法規三對應','課程設計完整對應國家課程綱要與特殊教育法規',accent=NAV)

cols14=[
    (FC,FCL,'📋 十二年國教課綱對應',
     [('F系列','生管-Ⅳ-E-1','認識台灣常用貨幣面額與使用'),
      ('C1-C3','生管-Ⅳ-E-2','使用金錢完成購物交易'),
      ('C5-C6','生管-Ⅳ-E-3','了解找零概念與驗收'),
      ('A系列','生管-Ⅳ-E-4','社區消費技能與自助設備'),
      ('A5',   'Ⅳ-F-1',    '使用自動化設備（ATM）'),
      ('B系列','生管-Ⅳ-G-1','財務管理基礎規劃')]),
    (CC,CCL,'🎯 IEP 年度目標範本',
     [('F1-F2','數學先備','在口語提示下完成10以內一對一對應操作'),
      ('C1-C3','貨幣認知','獨立辨識台灣六種硬幣並說出面額'),
      ('C4-C6','金錢操作','在部分協助下完成付款並核對找零'),
      ('A1',   '情境類化','獨立完成自動販賣機購買飲料完整流程'),
      ('A5',   '情境類化','在口語提示下完成 ATM 提款操作'),
      ('B1-B2','財務管理','依行程計算所需金額並準備正確零錢')]),
    (AC,ACL,'⚖️ 法規依據',
     [('特殊教育法','§28','依個別需求訂定個別化教育計畫（IEP）'),
      ('特殊教育法','§24','提供必要之輔助科技及支援服務'),
      ('課綱','2019','特殊需求領域生活管理：財務管理核心主題'),
      ('CRPD','Art.24','融合教育・平等受教機會保障'),
      ('UDL','CAST 2018','通用設計學習三原則系統落實'),
      ('工作分析','Gold 1980','步驟化教學完整對應行為目標格式')]),
]
for ci,(hc,lc,title,items) in enumerate(cols14):
    bx=Inches(.32+ci*4.35); cw=Inches(4.12)
    R(sl,bx,Inches(1.18),cw,Inches(5.68),fill=W,bc=hc,bw=1.6,rad=10000)
    R(sl,bx,Inches(1.18),cw,Inches(.46),fill=hc,rad=0)
    T(sl,title,bx+Inches(.15),Inches(1.24),cw-Inches(.3),Inches(.38),
      sz=12,bold=True,col=W)
    for ri,(a,b,c) in enumerate(items):
        ry=Inches(1.76+ri*.82)
        R(sl,bx+Inches(.14),ry,cw-Inches(.28),Inches(.74),
          fill=lc,bc=hc,bw=.8,rad=5000)
        pill(sl,a,bx+Inches(.2),ry+Inches(.1),Inches(.82),Inches(.22),
             bg=hc,fg=W,sz=8.5)
        pill(sl,b,bx+Inches(1.09),ry+Inches(.1),Inches(.92),Inches(.22),
             bg=W,fg=hc,sz=8.5)
        T(sl,c,bx+Inches(.2),ry+Inches(.38),cw-Inches(.5),Inches(.3),
          sz=10,col=INK)
pn(sl,14)

# ══════════════════════════════════════════════════════════════
# S15  試教觀察——三個案（前後對比卡）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'試教觀察——三個案・類化窗口的關鍵時刻',
     '學生在第 5–7 次操作後，出現程序性記憶語言化信號',accent=RD)

cases15=[
    (RD,RDL,'個案一  小涵','14歲・輕度智能障礙',
     ['✗ 能算 100−35=65','✗ 模擬購物反覆困惑','✗ 無法連結「我的錢」','  與「售價」的關係'],
     ['✓ 主動拿零錢逐枚核對','✓ 說出「這個35元，我有」','✓ 程序性記憶語言化','✓ 類化窗口開啟'],
     '突破點：步驟化流程\n讓靜態計算轉為\n理解交易序列意義'),
    (CC,CCL,'個案二  小俊','輕度自閉症',
     ['✗ 速食店選擇迴路','✗ 後方催促→情緒失調','✗ 一學期迴避相關學習','✗ 聽到「麥當勞」即迴避'],
     ['✓ 面對真實觸控機完成全程','✓ 主動說「我是7號還沒叫」','✓ 等待焦慮顯著改善','✓ 情境迴避行為消失'],
     '突破點：安全練習建立信心\n類化窗口在第5次\n後自然開啟'),
    (AC,ACL,'個案三  小雯','販賣機迴避行為',
     ['✗ 硬幣選錯多次失敗','✗ 繞道迴避所有自助設備','✗ 告訴老師「不會用」','✗ 情境焦慮長期積累'],
     ['✓ 主動走向真實販賣機','✓ 獨立完成選飲料投幣','✓ 說「我下次自己來」','✓ 情境迴避完全消失'],
     '突破點：無限次安全嘗試\n消除「失敗的社會\n後果」焦慮'),
]
for i,(hc,lc,name,tag,bef,aft,note) in enumerate(cases15):
    bx=Inches(.32+i*4.35); cw=Inches(4.12)
    # 標頭
    R(sl,bx,Inches(1.18),cw,Inches(.52),fill=hc,rad=0)
    T(sl,name,bx+Inches(.15),Inches(1.24),cw-Inches(.3),Inches(.38),
      sz=14,bold=True,col=W)
    pill(sl,tag,bx+Inches(.15),Inches(1.78),cw-Inches(.3),Inches(.28),
         bg=lc,fg=hc,sz=9.5,bold=False)
    # 試教前
    T(sl,'試教前',bx+Inches(.15),Inches(2.14),Inches(1.5),Inches(.28),
      sz=10.5,bold=True,col=RD)
    R(sl,bx+Inches(.15),Inches(2.4),cw-Inches(.3),Inches(1.72),
      fill=RDL,bc=RD,bw=.8,rad=5000)
    MT(sl,bef,bx+Inches(.28),Inches(2.5),cw-Inches(.5),Inches(1.5),
       sz=10.5,col=INK,sp=15)
    # 試教後
    T(sl,'試教後',bx+Inches(.15),Inches(4.2),Inches(1.5),Inches(.28),
      sz=10.5,bold=True,col=CC)
    R(sl,bx+Inches(.15),Inches(4.46),cw-Inches(.3),Inches(1.72),
      fill=CCL,bc=CC,bw=.8,rad=5000)
    MT(sl,aft,bx+Inches(.28),Inches(4.56),cw-Inches(.5),Inches(1.5),
       sz=10.5,col=INK,sp=15)
    # 突破標注
    callout(sl,note,bx+Inches(.15),Inches(6.3),cw-Inches(.3),Inches(.62))
pn(sl,15)

# ══════════════════════════════════════════════════════════════
# S16  預期成效（三層面×三時期矩陣）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'預期成效——三層面 × 短中長期指標',
     '學生・教師・學校三層面，建立完整評估框架',accent=CC)

# 橫向欄頭（三時期）
periods=['短期（1學期）','中期（1學年）','長期（跨學年）']
per_cols=[FC,CC,AC]
per_x=[Inches(3.0),Inches(6.5),Inches(10.0)]
for px,pt,pc in zip(per_x,periods,per_cols):
    R(sl,px,Inches(1.18),Inches(3.2),Inches(.44),fill=pc,rad=0)
    T(sl,pt,px,Inches(1.2),Inches(3.2),Inches(.38),
      sz=12,bold=True,col=W,al=PP_ALIGN.CENTER)

# 縱向列頭（三層面）
layers16=[
    (FC,FCL,'🎓 學生層面'),
    (CC,CCL,'👩‍🏫 教師層面'),
    (AC,ACL,'🏫 學校層面'),
]
cells=[
    ['完成 F/C 系列 80% 以上\n建立貨幣辨識付款基礎',
     '完成 A 系列 4 個以上\n部分提示下完成社區操作',
     '真實社區環境獨立完成\nATM・超市・自助點餐機'],
    ['熟悉軟體操作\nIEP 目標有效對應',
     '建立系統化評量記錄\n追蹤步驟別進步',
     '累積班級學習資料\n形成校本課程材料'],
    ['提供零成本數位教材\n補充紙本互動不足',
     '建立特教班數位學習環境\n支援多元評量實施',
     '成為區域特教分享典範\n連結融合教育推廣政策'],
]
for ri,(hc,lc,label) in enumerate(layers16):
    ry=Inches(1.62+ri*1.72)
    R(sl,Inches(.32),ry,Inches(2.5),Inches(1.6),fill=lc,bc=hc,bw=1.4,rad=7000)
    T(sl,label,Inches(.32),ry+Inches(.6),Inches(2.5),Inches(.4),
      sz=12,bold=True,col=hc,al=PP_ALIGN.CENTER)
    for ci,(px,pc) in enumerate(zip(per_x,per_cols)):
        R(sl,px,ry,Inches(3.2),Inches(1.6),fill=W,bc=LN,bw=.8)
        R(sl,px,ry,Inches(3.2),Inches(.06),fill=pc)
        MT(sl,cells[ri][ci].split('\n'),px+Inches(.14),ry+Inches(.14),
           Inches(2.92),Inches(1.3),sz=11,col=INK,sp=16)
pn(sl,16)

# ══════════════════════════════════════════════════════════════
# S17  三段學習旅程（橫向旅程圖）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'三段完整學習旅程',
     '數位互動練習 → 紙本作業評量 → 社區真實類化',accent=BC)

stages17=[
    (FC,FCL,'🖥️','第一段','數位互動練習',
     ['課堂使用軟體反覆操作',
      '三難度漸進個別適應',
      '即時語音回饋強化記憶',
      '允許無限次安全嘗試',
      '輔助點擊支援精細動作'],
     '重點：建立程序性記憶\n失敗不帶社會後果'),
    (CC,CCL,'📝','第二段','紙本作業評量',
     ['一鍵 PDF 作業單生成',
      '看圖付款/找零多元題型',
      '學生版/答案版一鍵切換',
      '步驟別評量對應 IEP',
      '形成性評量持續記錄'],
     '重點：形成性評量\n確認技能遷移程度'),
    (AC,ACL,'🏘️','第三段','社區真實類化',
     ['教師 3 公尺外靜態觀察',
      '逐步褪除支持提示',
      '真實設備實地操作',
      '社區技能類化確認',
      '終極目標：完全獨立'],
     '重點：類化窗口開啟\n終極成效確認'),
]
for i,(hc,lc,ic,stage,title,pts,note) in enumerate(stages17):
    bx=Inches(.35+i*4.35); cw=Inches(4.12)
    # 頂部圓形badge
    E(sl,bx+Inches(1.31),Inches(1.15),Inches(1.5),Inches(1.5),fill=hc)
    T(sl,ic,bx+Inches(1.31),Inches(1.18),Inches(1.5),Inches(1.4),
      sz=32,al=PP_ALIGN.CENTER)
    # 卡體
    R(sl,bx,Inches(2.72),cw,Inches(4.6),fill=lc,bc=hc,bw=1.8,rad=10000)
    T(sl,stage,bx,Inches(2.8),cw,Inches(.3),
      sz=10,col=GR,al=PP_ALIGN.CENTER,bold=True)
    T(sl,title,bx,Inches(3.1),cw,Inches(.38),
      sz=13,bold=True,col=hc,al=PP_ALIGN.CENTER)
    HL(sl,Inches(3.56),bx+Inches(.2),bx+cw-Inches(.2),col=hc,pt=.9)
    MT(sl,['• '+p for p in pts],bx+Inches(.2),Inches(3.64),
       cw-Inches(.3),Inches(2.2),sz=11,col=INK,sp=16)
    callout(sl,note,bx+Inches(.15),Inches(5.95),cw-Inches(.3),Inches(.56))
    if i<2:
        T(sl,'⟹',Inches(4.6+i*4.35),Inches(3.2),
          Inches(.55),Inches(.5),sz=22,col=GR,al=PP_ALIGN.CENTER)

callout(sl,'✦ 三段設計呼應課綱「社區功能性技能」目標・《特殊教育法》§24 輔助科技條款',
        Inches(.35),Inches(6.68),Inches(12.62),Inches(.38))
pn(sl,17)

# ══════════════════════════════════════════════════════════════
# S18  技術規格與配套系統
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'技術規格與配套系統','無障礙設計・作業單系統・獎勵系統',accent=FC)

# 左：技術規格
sec(sl,'技術規格',Inches(.4),Inches(1.22),col=FC)
specs18=[
    ('🌐','平台技術','HTML5 + 純 JavaScript（無框架依賴・純靜態頁面）'),
    ('📴','部署方式','完全離線・免安裝・免伺服器・USB 隨身碟即可使用'),
    ('💰','授權費用','NT$ 0・原始碼完整開放・可自由改作用於教育目的'),
    ('📏','程式規模','150,000+ 行程式碼・24 個教學單元・4 個配套系統'),
    ('📱','設備相容','平板・桌機・手機全相容・觸控與滑鼠均支援'),
    ('🔊','語音系統','Web Speech API 繁體中文・雅婷語音優先・速率可調'),
    ('♿','無障礙','視覺2倍縮放・高對比・語音全程引導・輔助點擊'),
    ('📊','評量整合','作業單 PDF 下載・獎勵系統・步驟評量記錄'),
]
for i,(ic,k,v) in enumerate(specs18):
    ry=Inches(1.62+i*.6)
    bg=FCL if i%2==0 else W
    R(sl,Inches(.4),ry,Inches(6.0),Inches(.52),fill=bg,bc=LN,bw=.8,rad=4000)
    T(sl,ic,Inches(.55),ry+Inches(.09),Inches(.4),Inches(.34),sz=14)
    T(sl,k,Inches(1.0),ry+Inches(.1),Inches(1.2),Inches(.3),
      sz=11,bold=True,col=FC)
    T(sl,v,Inches(2.32),ry+Inches(.1),Inches(3.9),Inches(.3),sz=11,col=INK)

# 右：作業單系統
sec(sl,'作業單系統',Inches(6.75),Inches(1.22),col=AC)
ws18=[
    ('🖨️','一鍵 PDF 生成','即時產生對應單元個人化作業單\n題目・答案版雙版本切換'),
    ('📋','智慧題目生成','自動依難度調整題型與數量\n看圖填空・付款找零多元題型'),
    ('🏪','商店客製','A4 依 12 種商店類型對應題目\nB 系列依主題生成場景題'),
    ('🏆','獎勵系統整合','完成任務解鎖貼紙徽章\n照片上傳・成就累積'),
]
for i,(ic,t,d) in enumerate(ws18):
    col=i%2; row=i//2
    bx=Inches(6.75+col*3.25)
    ry=Inches(1.62+row*1.72)
    R(sl,bx,ry,Inches(3.1),Inches(1.56),fill=ACL,bc=AC,bw=1.3,rad=8000)
    T(sl,ic,bx+Inches(.15),ry+Inches(.15),Inches(.52),Inches(.5),sz=20)
    T(sl,t,bx+Inches(.72),ry+Inches(.18),Inches(2.24),Inches(.3),
      sz=11.5,bold=True,col=AC)
    MT(sl,d.split('\n'),bx+Inches(.2),ry+Inches(.6),
       Inches(2.76),Inches(.86),sz=10.5,col=INK,sp=14)
pn(sl,18)

# ══════════════════════════════════════════════════════════════
# S19  教案摘要（F+C 系列）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'教案摘要（一）——F 系列・C 系列',
     '教案一：F1/F2 數學先備　教案二：C4/C5/C6 貨幣認知整合',accent=FC)

# F 系列教案
R(sl,Inches(.32),Inches(1.18),Inches(6.1),Inches(5.68),
  fill=FCL,bc=FC,bw=1.8,rad=10000)
T(sl,'📚 教案一：F 系列',Inches(.5),Inches(1.28),Inches(5.7),Inches(.38),
  sz=14,bold=True,col=FC)
callout(sl,'適用：國中特教班・資源班　3節×45分鐘',Inches(.5),Inches(1.72),Inches(5.7),Inches(.36))
fnodes=[
    ('第一節','F1 一對一對應',
     ['導入：實物配對引發生活連結',
      '示範：教師操作 F1 展示步驟',
      '引導：輔助點擊模式學生嘗試',
      '評量：步驟完成記錄 I/VP/PP/FP']),
    ('第二節','F2 唱數',
     ['複習：F1 連結「幾個就說幾」',
      '示範：逐步唱數視覺化操作',
      '引導：語音輔助模式跟讀練習',
      '評量：錄音樣本建立基線']),
    ('第三節','整合評量',
     ['整合：F1+F2 混合任務',
      '紙本：PDF 作業單形成性評量',
      '融合：普通班同儕互相示範',
      '更新：IEP 短期目標記錄']),
]
for i,(title,sub,pts) in enumerate(fnodes):
    bx=Inches(.5+i*1.96)
    R(sl,bx,Inches(2.18),Inches(1.84),Inches(4.38),
      fill=W,bc=FC,bw=1.2,rad=7000)
    R(sl,bx,Inches(2.18),Inches(1.84),Inches(.4),fill=FC,rad=0)
    T(sl,title,bx,Inches(2.2),Inches(1.84),Inches(.36),
      sz=11,bold=True,col=W,al=PP_ALIGN.CENTER)
    T(sl,sub,bx+Inches(.1),Inches(2.66),Inches(1.64),Inches(.3),
      sz=10.5,bold=True,col=FC)
    MT(sl,pts,bx+Inches(.1),Inches(3.02),Inches(1.64),Inches(3.38),
       sz=10,col=INK,sp=14)

# C 系列教案
R(sl,Inches(6.6),Inches(1.18),Inches(6.42),Inches(5.68),
  fill=CCL,bc=CC,bw=1.8,rad=10000)
T(sl,'🪙 教案二：C 系列',Inches(6.78),Inches(1.28),Inches(6.0),Inches(.38),
  sz=14,bold=True,col=CC)
callout(sl,'C4 付款・C5 夠不夠・C6 找零　核心融合節次',Inches(6.78),Inches(1.72),Inches(5.9),Inches(.36))
cnodes=[
    ('第一節','C4 付款技能建立',
     ['付款面額組合策略訓練',
      '拖曳操作程序記憶建立',
      '三難度漸進挑戰']),
    ('第二節','C5+C6 整合',
     ['C5：足額判斷比較思維',
      'C6：找零計算驗收閉環',
      '多步驟序列程序記憶']),
    ('第三節','融合情境評量🏆',
     ['普通班同儕擔任「收銀員」',
      '特殊需求學生擔任「顧客」',
      '雙向融合目標同步實現']),
]
for i,(title,sub,pts) in enumerate(cnodes):
    bx=Inches(6.78+i*2.06)
    R(sl,bx,Inches(2.18),Inches(1.9),Inches(4.38),
      fill=W,bc=CC,bw=1.2,rad=7000)
    R(sl,bx,Inches(2.18),Inches(1.9),Inches(.4),fill=CC,rad=0)
    T(sl,title,bx,Inches(2.2),Inches(1.9),Inches(.36),
      sz=11,bold=True,col=W,al=PP_ALIGN.CENTER)
    T(sl,sub,bx+Inches(.1),Inches(2.66),Inches(1.7),Inches(.3),
      sz=10.5,bold=True,col=CC)
    MT(sl,pts,bx+Inches(.1),Inches(3.02),Inches(1.7),Inches(3.38),
       sz=10.5,col=INK,sp=16)
pn(sl,19)

# ══════════════════════════════════════════════════════════════
# S20  教案三：A 系列
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'教案摘要（二）——A 系列真實情境類化模組',
     '教案三：A1 自動販賣機（核心融合節次）・A5 ATM 操作精要',accent=AC)

# A1 教案（左側大區塊）
R(sl,Inches(.32),Inches(1.18),Inches(7.9),Inches(5.7),
  fill=ACL,bc=AC,bw=1.8,rad=10000)
T(sl,'🥤 A1 自動販賣機教案（3 節）',Inches(.52),Inches(1.28),
  Inches(7.5),Inches(.38),sz=14,bold=True,col=AC)
callout(sl,'核心任務：建立 7 步驟完整操作程序記憶・達成社區真實類化',
        Inches(.52),Inches(1.72),Inches(7.5),Inches(.36))

a1nodes=[
    (FC,'第一節','流程建立',
     ['工作分析 7 步驟逐一教導',
      '輔助點擊全程引導',
      '語音提示每步操作要點',
      '簡單模式→普通模式過渡']),
    (CC,'第二節','投幣找零強化',
     ['面額選擇策略訓練',
      '找零驗收程序建立',
      '三難度循序挑戰',
      'coinFirst 先投幣模式']),
    (AC,'第三節','融合教育頂點 🏆',
     ['邀請普通班同儕為夥伴',
      '在校園真實販賣機操作',
      '教師 3 公尺外靜態觀察',
      '→ 類化窗口開啟關鍵節次']),
]
for i,(hc,title,sub,pts) in enumerate(a1nodes):
    bx=Inches(.52+i*2.6)
    R(sl,bx,Inches(2.18),Inches(2.44),Inches(4.4),
      fill=W,bc=hc,bw=1.3,rad=7000)
    R(sl,bx,Inches(2.18),Inches(2.44),Inches(.42),fill=hc,rad=0)
    T(sl,title,bx,Inches(2.2),Inches(2.44),Inches(.36),
      sz=12,bold=True,col=W,al=PP_ALIGN.CENTER)
    T(sl,sub,bx+Inches(.12),Inches(2.68),Inches(2.2),Inches(.3),
      sz=11,bold=True,col=hc)
    MT(sl,pts,bx+Inches(.12),Inches(3.04),Inches(2.2),Inches(3.36),
       sz=10.5,col=INK,sp=15)

# A5 教案（右側）
R(sl,Inches(8.4),Inches(1.18),Inches(4.6),Inches(5.7),
  fill=BCL,bc=BC,bw=1.8,rad=10000)
T(sl,'🏧 A5 ATM 操作精要',Inches(8.6),Inches(1.28),Inches(4.2),Inches(.38),
  sz=13,bold=True,col=BC)
a5pts=[
    ('教學時間','3節×45分鐘'),
    ('學習階段','高中特教班'),
    ('學習目標','獨立完成提款（I）\n部分提示下存款/轉帳'),
    ('第一節','介面認識・提款流程'),
    ('第二節','存款・轉帳流程'),
    ('第三節','真實 ATM 社區演練'),
    ('融合設計','普通班同學擔任\n「等候示範者」'),
    ('評量工具','8步驟評量表\nIEP 轉銜目標'),
]
ry=Inches(1.76)
for k,v in a5pts:
    R(sl,Inches(8.6),ry,Inches(4.1),Inches(.58),
      fill=W if a5pts.index((k,v))%2==0 else BCL,bc=LN,bw=.6,rad=3000)
    T(sl,k+'：',Inches(8.72),ry+Inches(.09),Inches(1.1),Inches(.38),
      sz=10.5,bold=True,col=BC)
    MT(sl,v.split('\n'),Inches(9.82),ry+Inches(.09),Inches(2.7),Inches(.38),
       sz=10.5,col=INK,sp=13)
    ry+=Inches(.62)
pn(sl,20)

# ══════════════════════════════════════════════════════════════
# S21  核心創新——四大首創
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
hbar(sl,'核心創新——四大首創特色',
     '填補台灣特殊教育金融素養教材十年空白',accent=BC)

innov=[
    (BC,BCL,'🖱️','首創\n輔助點擊模式','精細動作障礙解方',
     ['全 24 單元均支援單一點擊','不需拖曳・不需精確觸控',
      '可疊加任何難度層級','全台同類型特教軟體首例'],
     '腦性麻痺・脊損・肌萎・ADHD'),
    (AC,ACL,'🏪','首創\n六場景仿真','社區設備互動模擬',
     ['自動販賣機・理髮廳・麥當勞','ATM・超市收銀台・火車票',
      '像素到流程忠實仿真','降低真實場域焦慮的過渡橋'],
     '類化轉移・社區獨立操作'),
    (CC,CCL,'📚','首創\n完整序階課程','24單元×4系列系統',
     ['F→C→A→B 完整學習序階','從數學先備到預算規劃',
      '全台首個完整金融素養套件','課綱全部學習主題完整覆蓋'],
     '國小至高中完整學習路徑'),
    (FC,FCL,'🎯','首創\nIEP 整合設計','評量-作業單-獎勵三合一',
     ['工作分析法步驟評量格式','PDF 作業單一鍵生成',
      '獎勵系統強化學習動機','教師行政負擔降低 50%+'],
     'IEP目標・形成性評量・激勵'),
]
for i,(hc,lc,ic,title,sub,pts,tag) in enumerate(innov):
    bx=Inches(.35+i*3.25); cw=Inches(3.0)
    # 頂部圓形icon
    E(sl,bx+Inches(1.0),Inches(1.15),Inches(1.0),Inches(1.0),fill=hc)
    T(sl,ic,bx+Inches(1.0),Inches(1.18),Inches(1.0),Inches(.9),
      sz=26,al=PP_ALIGN.CENTER)
    # 卡體
    R(sl,bx,Inches(2.22),cw,Inches(4.88),fill=lc,bc=hc,bw=1.8,rad=10000)
    R(sl,bx,Inches(2.22),cw,Inches(.52),fill=hc,rad=0)
    T(sl,title,bx,Inches(2.26),cw,Inches(.44),
      sz=13,bold=True,col=W,al=PP_ALIGN.CENTER)
    T(sl,sub,bx,Inches(2.78),cw,Inches(.28),
      sz=10,col=GR,al=PP_ALIGN.CENTER)
    HL(sl,Inches(3.14),bx+Inches(.2),bx+cw-Inches(.2),col=hc,pt=.8)
    MT(sl,['▸ '+p for p in pts],bx+Inches(.18),Inches(3.22),
       cw-Inches(.3),Inches(2.28),sz=11,col=INK,sp=17)
    pill(sl,tag,bx+Inches(.15),Inches(5.58),cw-Inches(.3),Inches(.26),
         bg=hc,fg=W,sz=9,bold=False)
pn(sl,21)

# ══════════════════════════════════════════════════════════════
# S22  結語（統計儀表板＋核心訴求）
# ══════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(BLK)
R(sl,Inches(0),Inches(0),SW,SH,fill=BG)
R(sl,Inches(0),Inches(0),SW,Inches(.18),fill=NAV)
R(sl,Inches(0),Inches(.18),SW,Inches(.06),fill=GO)
# 左側深藍欄
R(sl,Inches(0),Inches(.24),Inches(4.5),SH-Inches(.24),fill=NAV)
for i,(c,fy) in enumerate([(FC,1.5),(CC,3.0),(AC,4.5),(BC,6.0)]):
    R(sl,Inches(0),Inches(fy),Inches(4.5),Inches(.08),fill=c)
T(sl,'感謝聆聽',Inches(.2),Inches(1.0),Inches(4.1),Inches(.8),
  sz=30,bold=True,col=W,al=PP_ALIGN.CENTER)
T(sl,'金錢小達人',Inches(.2),Inches(1.9),Inches(4.1),Inches(.7),
  sz=22,bold=True,col=RGBColor(0xFD,0xC8,0x70),al=PP_ALIGN.CENTER)
T(sl,'互動式特殊教育\n生活數學暨金融素養教學軟體',
  Inches(.2),Inches(2.65),Inches(4.1),Inches(.72),
  sz=12,col=RGBColor(0x9B,0xBB,0xDD),al=PP_ALIGN.CENTER)
for i,(v,l,c) in enumerate([
    ('24','教學單元',FC),('4','課程系列',CC),
    ('150K+','程式碼行數',AC),('$0','授權費用',BC)]):
    ey=Inches(3.62+i*.8)
    R(sl,Inches(.2),ey,Inches(4.1),Inches(.7),fill=RGBColor(0x1E,0x45,0x7A),rad=6000)
    T(sl,v,Inches(.2),ey+Inches(.04),Inches(1.8),Inches(.44),
      sz=22,bold=True,col=c,al=PP_ALIGN.CENTER)
    T(sl,l,Inches(2.0),ey+Inches(.2),Inches(2.1),Inches(.28),
      sz=11,col=RGBColor(0xCC,0xDD,0xEE))
T(sl,'中華民國 115 年 3 月',Inches(.2),Inches(7.12),Inches(4.1),Inches(.28),
  sz=10,col=RGBColor(0x5A,0x7A,0x9A),al=PP_ALIGN.CENTER)

# 右側核心訴求
T(sl,'核心訴求',Inches(4.9),Inches(.44),Inches(4),Inches(.38),
  sz=14,bold=True,col=NAV)
msgs=[
    (BC,BCL,'🎯','填補十年空白',
     '全台首個完整系統化特殊教育\n金融素養互動教學軟體套件（24單元）'),
    (AC,ACL,'🤝','零門檻推廣',
     'NT$0 授權・完全離線・免安裝\n讓每位特教教師立即可用'),
    (CC,CCL,'🌱','真實成效',
     '三個個案的類化轉折紀錄\n從情境迴避到社區獨立操作的突破'),
    (FC,FCL,'📋','完整支持',
     'IEP整合・作業單・獎勵・融合教育\n提供教師全方位教學支援'),
]
for i,(hc,lc,ic,title,desc) in enumerate(msgs):
    col=i%2; row=i//2
    bx=Inches(4.9+col*4.15)
    ry=Inches(.9+row*3.12)
    R(sl,bx,ry,Inches(3.9),Inches(2.88),fill=lc,bc=hc,bw=1.8,rad=10000)
    T(sl,ic,bx+Inches(.2),ry+Inches(.28),Inches(.7),Inches(.6),sz=26)
    T(sl,title,bx+Inches(1.0),ry+Inches(.32),Inches(2.72),Inches(.4),
      sz=14,bold=True,col=hc)
    HL(sl,ry+Inches(.82),bx+Inches(.18),bx+Inches(3.72),col=hc,pt=.8)
    MT(sl,desc.split('\n'),bx+Inches(.18),ry+Inches(.96),
       Inches(3.54),Inches(1.7),sz=11,col=INK,sp=16)

callout(sl,'✦ 程式碼逾 15 萬行，完全由教師個人研發，零外部資金，零商業目的，只為讓特殊需求學生平等享有數位學習機會',
        Inches(4.9),Inches(7.1),Inches(8.15),Inches(.34))
pn(sl,22)

# ── 儲存 ──────────────────────────────────────────────────────
out='金錢小達人_競賽申請簡報_V5.pptx'
prs.save(out)
print(f'✓ 已儲存：{out}')
