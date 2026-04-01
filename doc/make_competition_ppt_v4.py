# -*- coding: utf-8 -*-
"""
金錢小達人 競賽申請 PPT V4
視覺風格：教科書圖解・教育知識科普
共 22 張投影片
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 教科書配色 ──────────────────────────────────────────────
W    = RGBColor(0xFF,0xFF,0xFF)
BG   = RGBColor(0xF9,0xF7,0xF2)   # 米白底
BG2  = RGBColor(0xF0,0xEC,0xE3)   # 淺暖灰
LN   = RGBColor(0xD6,0xCE,0xBD)   # 分隔線
GR   = RGBColor(0x6B,0x65,0x5A)   # 說明文字
DK   = RGBColor(0x2C,0x27,0x1E)   # 深色主文
HDB  = RGBColor(0x1A,0x3A,0x5C)   # 標題欄深藍
# 系列色
F_H  = RGBColor(0x1D,0x6F,0xB8)   # F藍
F_L  = RGBColor(0xD6,0xEA,0xF8)
C_H  = RGBColor(0x14,0x7A,0x4E)   # C綠
C_L  = RGBColor(0xD5,0xF0,0xE3)
A_H  = RGBColor(0xC0,0x44,0x0C)   # A橙
A_L  = RGBColor(0xFD,0xE8,0xD5)
B_H  = RGBColor(0x7B,0x2F,0xBE)   # B紫
B_L  = RGBColor(0xEE,0xE0,0xFB)
RD   = RGBColor(0xB9,0x2B,0x27)
RDL  = RGBColor(0xFD,0xE8,0xE7)
AM   = RGBColor(0xB4,0x5D,0x07)
AML  = RGBColor(0xFE,0xF3,0xD7)
TE   = RGBColor(0x0C,0x7C,0x72)
TEL  = RGBColor(0xCE,0xF5,0xF1)
IND  = RGBColor(0x3A,0x2F,0xB0)
INDL = RGBColor(0xE4,0xE1,0xFB)
# 教科書標注色
NOTE_BG = RGBColor(0xFF,0xF8,0xDC)   # 淡黃標注
NOTE_BD = RGBColor(0xE0,0xC8,0x30)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── 基礎繪圖函數 ────────────────────────────────────────────
def rect(sl, x, y, w, h, fill=None, bc=None, bpt=0, rad=None):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    if bc and bpt > 0:
        sh.line.color.rgb = bc; sh.line.width = Pt(bpt)
    if rad is not None:
        sp = sh._element; spPr = sp.find(qn('p:spPr'))
        pg = spPr.find(qn('a:prstGeom'))
        if pg is not None:
            pg.set('prst','roundRect')
            av = pg.find(qn('a:avLst'))
            if av is None: av = etree.SubElement(pg, qn('a:avLst'))
            gd = etree.SubElement(av, qn('a:gd'))
            gd.set('name','adj'); gd.set('fmla',f'val {rad}')
    return sh

def ellipse(sl, x, y, w, h, fill=None, bc=None, bpt=0):
    sh = sl.shapes.add_shape(9, x, y, w, h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    if bc and bpt>0: sh.line.color.rgb=bc; sh.line.width=Pt(bpt)
    return sh

def txt(sl, s, x, y, w, h, sz=13, bold=False, col=None,
        align=PP_ALIGN.LEFT, fn='Microsoft JhengHei', italic=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.name = fn; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    if col: r.font.color.rgb = col
    return tb

def mtxt(sl, lines, x, y, w, h, sz=12, col=None,
         align=PP_ALIGN.LEFT, fn='Microsoft JhengHei', bold=False, lh=None):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if lh:
            from pptx.util import Pt as _Pt
            p.line_spacing = _Pt(lh)
        r = p.add_run(); r.text = line
        r.font.name = fn; r.font.size = Pt(sz)
        r.font.bold = bold
        if col: r.font.color.rgb = col
    return tb

def hline(sl, y, x0=Inches(0.5), x1=Inches(12.83), col=LN, pt=1.0):
    ln = sl.shapes.add_connector(1, x0, y, x1, y)
    ln.line.color.rgb = col; ln.line.width = Pt(pt)

def vline(sl, x, y0, y1, col=LN, pt=1.0):
    ln = sl.shapes.add_connector(1, x, y0, x, y1)
    ln.line.color.rgb = col; ln.line.width = Pt(pt)

def pill(sl, s, x, y, w, h, bg, fg, sz=11, bold=True):
    rect(sl, x, y, w, h, fill=bg, rad=20000)
    txt(sl, s, x, y, w, h, sz=sz, bold=bold, col=fg, align=PP_ALIGN.CENTER)

def page_num(sl, n, total=22):
    txt(sl, f'{n} / {total}',
        SLIDE_W-Inches(1.1), SLIDE_H-Inches(0.38),
        Inches(0.9), Inches(0.32),
        sz=9, col=GR, align=PP_ALIGN.RIGHT)

def top_bar(sl, title, subtitle='', accent=F_H, h=Inches(1.0)):
    """教科書風格標題欄：左側色帶 + 白底橫線"""
    rect(sl, Inches(0), Inches(0), SLIDE_W, h, fill=W)
    rect(sl, Inches(0), Inches(0), Inches(0.22), h, fill=accent)
    hline(sl, h, Inches(0), SLIDE_W, col=accent, pt=2.0)
    txt(sl, title, Inches(0.45), Inches(0.1), SLIDE_W-Inches(5.0), Inches(0.58),
        sz=22, bold=True, col=HDB)
    if subtitle:
        txt(sl, subtitle, Inches(0.45), Inches(0.63), SLIDE_W-Inches(2.0), Inches(0.3),
            sz=10.5, col=GR, italic=True)

def section_label(sl, s, x, y, col=F_H):
    """教科書章節標籤"""
    rect(sl, x, y, Inches(0.06), Inches(0.32), fill=col)
    txt(sl, s, x+Inches(0.14), y, Inches(3.5), Inches(0.32),
        sz=12, bold=True, col=col)

def note_box(sl, s, x, y, w, h):
    """淡黃色標注框（教科書重點框）"""
    rect(sl, x, y, w, h, fill=NOTE_BG, bc=RGBColor(0xD4,0xAC,0x0D), bpt=1.2, rad=8000)
    txt(sl, s, x+Inches(0.12), y+Inches(0.08), w-Inches(0.24), h-Inches(0.16),
        sz=11, col=RGBColor(0x78,0x49,0x04))

def icon_card(sl, icon, title, lines, x, y, w, h,
              bg=F_L, bd=F_H, tc=F_H):
    rect(sl, x, y, w, h, fill=bg, bc=bd, bpt=1.4, rad=10000)
    txt(sl, icon, x, y+Inches(0.12), w, Inches(0.52),
        sz=24, align=PP_ALIGN.CENTER)
    txt(sl, title, x, y+Inches(0.62), w, Inches(0.35),
        sz=11.5, bold=True, col=tc, align=PP_ALIGN.CENTER)
    mtxt(sl, lines, x+Inches(0.1), y+Inches(0.95),
         w-Inches(0.2), h-Inches(1.05), sz=10, col=DK)

def arrow_right(sl, x, y, length=Inches(0.35), col=GR):
    from pptx.util import Emu as _E
    ln = sl.shapes.add_connector(1, x, y, x+length, y)
    ln.line.color.rgb = col; ln.line.width = Pt(1.8)

# ═══════════════════════════════════════════════════════════════
# S01  封面（教科書風格）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
# 左側深色側欄
rect(sl, Inches(0), Inches(0), Inches(3.2), SLIDE_H, fill=HDB)
# 側欄裝飾條
for i,(c,frac) in enumerate([(F_H,0.25),(C_H,0.5),(A_H,0.75),(B_H,1.0)]):
    rect(sl, Inches(0), Inches(7.5*frac-0.05), Inches(3.2), Inches(0.08), fill=c)
# 側欄文字
txt(sl, '金', Inches(0.3), Inches(0.8), Inches(2.6), Inches(0.9),
    sz=52, bold=True, col=W, align=PP_ALIGN.CENTER)
txt(sl, '錢', Inches(0.3), Inches(1.65), Inches(2.6), Inches(0.9),
    sz=52, bold=True, col=W, align=PP_ALIGN.CENTER)
txt(sl, '小', Inches(0.3), Inches(2.5), Inches(2.6), Inches(0.9),
    sz=52, bold=True, col=W, align=PP_ALIGN.CENTER)
txt(sl, '達', Inches(0.3), Inches(3.35), Inches(2.6), Inches(0.9),
    sz=52, bold=True, col=W, align=PP_ALIGN.CENTER)
txt(sl, '人', Inches(0.3), Inches(4.2), Inches(2.6), Inches(0.9),
    sz=52, bold=True, col=W, align=PP_ALIGN.CENTER)
txt(sl, 'MONEY', Inches(0.3), Inches(5.4), Inches(2.6), Inches(0.38),
    sz=13, col=RGBColor(0x8A,0xA8,0xCC), align=PP_ALIGN.CENTER, fn='Arial', bold=True)
txt(sl, 'TUTOR', Inches(0.3), Inches(5.75), Inches(2.6), Inches(0.38),
    sz=13, col=RGBColor(0x8A,0xA8,0xCC), align=PP_ALIGN.CENTER, fn='Arial', bold=True)
# 右側主體
pill(sl, '第八屆教育部獎助研發特殊教育教材教具競賽　·　身心障礙類　·　國民中學及高級中等學校組',
     Inches(3.5), Inches(0.42), Inches(9.5), Inches(0.4),
     bg=RGBColor(0xE8,0xF0,0xFB), fg=HDB, sz=10.5)
txt(sl, '互動式特殊教育生活數學暨金融素養教學軟體',
    Inches(3.5), Inches(1.05), Inches(9.5), Inches(0.55),
    sz=19, bold=True, col=HDB)
txt(sl, 'Interactive Software for Financial Literacy in Special Education',
    Inches(3.5), Inches(1.6), Inches(9.5), Inches(0.36),
    sz=11, col=GR, italic=True, fn='Arial')
hline(sl, Inches(2.1), Inches(3.5), Inches(12.8), col=LN, pt=1.5)
# 四大特色卡片
for i,(ic,t,d,bg,bd) in enumerate([
    ('📚','F 系列','數學基礎先備\n6 個互動單元',F_L,F_H),
    ('🪙','C 系列','貨幣認知建立\n6 個互動單元',C_L,C_H),
    ('🏪','A 系列','真實情境類化\n6 個互動單元',A_L,A_H),
    ('💰','B 系列','預算規劃延伸\n6 個互動單元',B_L,B_H),
]):
    bx = Inches(3.5 + i*2.4)
    rect(sl, bx, Inches(2.22), Inches(2.2), Inches(1.22),
         fill=bg, bc=bd, bpt=1.5, rad=10000)
    txt(sl, ic, bx, Inches(2.28), Inches(2.2), Inches(0.45),
        sz=20, align=PP_ALIGN.CENTER)
    txt(sl, t, bx, Inches(2.7), Inches(2.2), Inches(0.3),
        sz=11, bold=True, col=bd, align=PP_ALIGN.CENTER)
    txt(sl, d, bx, Inches(2.98), Inches(2.2), Inches(0.44),
        sz=9.5, col=DK, align=PP_ALIGN.CENTER)
# 統計數字列
hline(sl, Inches(3.58), Inches(3.5), Inches(12.8), col=LN, pt=1.0)
for i,(v,l,c) in enumerate([
    ('18','互動教學單元',F_H),
    ('4 系列','課程序階設計',C_H),
    ('150,000+','程式碼行數',A_H),
    ('NT$ 0','授權費用',B_H),
    ('完全離線','無需網路',TE),
]):
    bx = Inches(3.5 + i*1.88)
    txt(sl, v, bx, Inches(3.68), Inches(1.85), Inches(0.46),
        sz=18, bold=True, col=c, align=PP_ALIGN.CENTER)
    txt(sl, l, bx, Inches(4.12), Inches(1.85), Inches(0.28),
        sz=9, col=GR, align=PP_ALIGN.CENTER)
hline(sl, Inches(4.52), Inches(3.5), Inches(12.8), col=LN, pt=1.0)
# 適用對象
txt(sl, '🎯 適用對象', Inches(3.5), Inches(4.62), Inches(3.5), Inches(0.35),
    sz=12, bold=True, col=HDB)
txt(sl, '國中～高中特教班・資源班', Inches(3.5), Inches(4.95), Inches(3.5), Inches(0.3),
    sz=11, col=DK)
txt(sl, '輕度至中度智能障礙・自閉症・學習障礙・多重障礙',
    Inches(3.5), Inches(5.25), Inches(6.5), Inches(0.28),
    sz=10.5, col=GR)
note_box(sl, '✦ 科技輔助教學軟體（HTML5）・免安裝・完全離線・授權費 0 元',
         Inches(3.5), Inches(5.7), Inches(9.0), Inches(0.42))
txt(sl, '中華民國 115 年 3 月',
    Inches(3.5), Inches(6.5), Inches(4), Inches(0.3),
    sz=10, col=GR)
page_num(sl, 1)

# ═══════════════════════════════════════════════════════════════
# S02  目錄
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '簡報綱要', 'Table of Contents', accent=HDB)
toc = [
    (RD,  '01','設計動機',    '特教現場三重結構性困境'),
    (A_H, '02','設計理念',    'UDL × 工作分析 × 直接教學 × 融合三層'),
    (F_H, '03','F 系列',      '數學基礎先備 6 單元序階'),
    (C_H, '04','C 系列',      '貨幣認知建立 6 單元序階'),
    (A_H, '05','A 系列',      '真實情境類化 6 單元序階'),
    (B_H, '06','B 系列',      '預算規劃延伸 6 單元介紹'),
    (TE,  '07','輔助點擊',    '精細動作障礙首創解方'),
    (AM,  '08','三難度系統',  '簡單 / 普通 / 困難完整比較'),
    (IND, '09','工作分析法',  '步驟化教學流程圖解'),
    (C_H, '10','融合教育',    '三層支持架構圖解'),
    (HDB, '11','課程法規',    'IEP / 特教法 / 課綱對應'),
    (RD,  '12','試教觀察',    '三個個案・類化轉折記錄'),
    (AM,  '13','預期成效',    '三層面 × 短中長期指標'),
    (F_H, '14','技術規格',    '無障礙設計・作業單系統'),
    (DK,  '15','核心創新',    '四大首創特色・結語'),
]
cols = 3
for idx,(c,num,t,s) in enumerate(toc):
    col = idx % cols
    row = idx // cols
    bx = Inches(0.32 + col*4.35)
    by = Inches(1.12 + row*1.22)
    rect(sl, bx, by, Inches(4.12), Inches(1.1),
         fill=W, bc=LN, bpt=1.2, rad=8000)
    rect(sl, bx, by, Inches(0.08), Inches(1.1), fill=c, rad=0)
    txt(sl, num, bx+Inches(0.18), by+Inches(0.12),
        Inches(0.6), Inches(0.42), sz=20, bold=True, col=c, fn='Arial')
    txt(sl, t,   bx+Inches(0.85), by+Inches(0.12),
        Inches(3.15), Inches(0.36), sz=13, bold=True, col=DK)
    txt(sl, s,   bx+Inches(0.85), by+Inches(0.5),
        Inches(3.15), Inches(0.46), sz=9.5, col=GR)
page_num(sl, 2)

# ═══════════════════════════════════════════════════════════════
# S03  設計動機 — 三重困境（教科書案例框）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '設計動機——特教現場三重結構性困境', '壹、來自十五年特教現場的教材研發回應', accent=RD)
section_label(sl, '三個結構性問題  →  一個系統性解方', Inches(0.4), Inches(1.2), col=RD)

# 三欄困境卡
cases = [
    (RD, RDL, '困境一', '交易序列的意義理解障礙',
     '個案・小涵（14歲，輕度智能障礙）',
     ['✗ 能計算 100−35=65', '✗ 卻反覆困惑「我有多少錢？」', '  與「這個東西多少錢？」的連結',
      '→ 問題不在數字，在交易序列的邏輯鏈'],
     '交易序列 = 付出 → 接受 → 驗收\n三步驟意義理解是核心障礙點'),
    (A_H, A_L, '困境二', '安全練習場域的長期缺位',
     '個案・小俊（輕度自閉症）',
     ['✗ 速食店自助點餐機選擇迴路', '✗ 顧客催促誘發情緒失調', '✗ 此後一學期迴避相關學習',
      '→ 真實環境試誤風險過高'],
     '特教現場迫切需要：\n允許無限次安全嘗試・失敗零社會後果'),
    (B_H, B_L, '困境三', '系統性課程工具的十年缺口',
     '課綱依據：生管-Ⅳ-E-1／F-1／E-4',
     ['✗ 課綱明確列出財務管理學習主題', '✗ 市場無對應互動軟體達十年', '✗ 教師只能自製紙本，缺乏動態回饋',
      '→ 系統性工具缺口'],
     '《十二年國教特殊需求領域課程綱要》\n2019 頁 45–52 明確要求，市場空白十年'),
]
for i,(hc,lc,tag,title,case,pts,note) in enumerate(cases):
    bx = Inches(0.32 + i*4.35)
    by = Inches(1.58)
    h  = Inches(5.58)
    rect(sl, bx, by, Inches(4.12), h, fill=W, bc=hc, bpt=1.8, rad=10000)
    # 標籤頂欄
    rect(sl, bx, by, Inches(4.12), Inches(0.46), fill=hc, rad=0)
    txt(sl, tag, bx+Inches(0.15), by+Inches(0.04), Inches(1.2), Inches(0.38),
        sz=13, bold=True, col=W)
    txt(sl, title, bx+Inches(0.15), by+Inches(0.55), Inches(3.8), Inches(0.4),
        sz=12, bold=True, col=hc)
    # 個案標籤
    pill(sl, case, bx+Inches(0.15), by+Inches(0.98),
         Inches(3.82), Inches(0.3), bg=lc, fg=hc, sz=10, bold=False)
    # 要點
    mtxt(sl, pts, bx+Inches(0.18), by+Inches(1.4),
         Inches(3.76), Inches(2.0), sz=10.5, col=DK, lh=16)
    # 底部標注
    note_box(sl, note, bx+Inches(0.15), by+Inches(3.52),
             Inches(3.82), Inches(0.72))
page_num(sl, 3)

# ═══════════════════════════════════════════════════════════════
# S04  設計理念 — UDL × 工作分析 × 直接教學 × 融合教育
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '設計理念——四大教育理論支柱', '貳、UDL × 工作分析法 × 直接教學法 × 融合教育三層支持', accent=A_H)

# 四柱圖（教科書圖解風格：底座＋柱子）
base_y = Inches(6.6)
cols_data = [
    (F_H, F_L, '🎯', 'UDL 通用設計學習',
     ['多元表徵', '多元行動與表達', '多元參與方式'],
     'CAST 2018 框架\n系統落實三大原則'),
    (C_H, C_L, '🔧', '工作分析法',
     ['拆解操作步驟', '每步可單獨評量', 'IEP 行為目標格式'],
     'Gold 1980\nSnell & Brown 2011'),
    (A_H, A_L, '📈', '直接教學法三階段',
     ['I 示範階段', 'II 引導練習', 'III 獨立操作'],
     'Engelmann & Carnine 1982\n難度漸進設計'),
    (B_H, B_L, '🤝', '融合教育三層支持',
     ['L1 普通班同步參與', 'L2 資源班強化練習', 'L3 自足班精熟學習'],
     '特殊教育法 §28\n最少限制環境原則'),
]
for i,(hc,lc,ic,title,pts,ref) in enumerate(cols_data):
    bx = Inches(0.42 + i*3.25)
    col_w = Inches(3.0)
    col_h = Inches(1.15 + i*0.0)  # equal height
    col_h = Inches(4.6)
    # 柱子底色
    rect(sl, bx, Inches(1.75), col_w, col_h, fill=lc, bc=hc, bpt=1.6, rad=10000)
    # 頂部icon區
    rect(sl, bx, Inches(1.75), col_w, Inches(0.5), fill=hc, rad=0)
    txt(sl, ic, bx, Inches(1.8), col_w, Inches(0.4),
        sz=16, align=PP_ALIGN.CENTER)
    txt(sl, title, bx, Inches(2.32), col_w, Inches(0.38),
        sz=11.5, bold=True, col=hc, align=PP_ALIGN.CENTER)
    hline(sl, Inches(2.72), bx+Inches(0.2), bx+col_w-Inches(0.2), col=hc, pt=1.0)
    mtxt(sl, ['• '+p for p in pts], bx+Inches(0.15), Inches(2.8),
         col_w-Inches(0.3), Inches(1.65), sz=11, col=DK, lh=18)
    note_box(sl, ref, bx+Inches(0.12), Inches(4.58),
             col_w-Inches(0.24), Inches(0.55))
# 底部統整說明
rect(sl, Inches(0.42), Inches(5.32), Inches(12.5), Inches(0.52),
     fill=RGBColor(0xE8,0xF0,0xFB), bc=HDB, bpt=1.2, rad=8000)
txt(sl, '🔗 四大理論整合應用：相同教材・相同目標方向，差異在「支持程度（Level of Support）」，而非課程內容的全然切割',
    Inches(0.65), Inches(5.36), Inches(12.1), Inches(0.44),
    sz=11, col=HDB)
page_num(sl, 4)

# ═══════════════════════════════════════════════════════════════
# S05  課程架構全覽（樓梯圖解）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '課程架構全覽——四系列 24 單元完整序階', '三層遞進設計：數學基礎 → 貨幣認知 → 真實情境 → 預算規劃延伸', accent=F_H)

# 樓梯圖
stairs = [
    (B_H, B_L, 'B 系列', '預算規劃延伸', 'B1–B6', '今日帶多少錢・零用錢日記・存錢計畫・特賣比一比・生日派對預算・菜市場買菜', Inches(3.0), Inches(1.18)),
    (A_H, A_L, 'A 系列', '真實情境類化', 'A1–A6', '自動販賣機・理髮廳・麥當勞・超市購物・ATM・火車票', Inches(2.0), Inches(2.32)),
    (C_H, C_L, 'C 系列', '貨幣認知建立', 'C1–C6', '認識錢幣・數錢・換錢・付款・夠不夠・找零', Inches(1.0), Inches(3.46)),
    (F_H, F_L, 'F 系列', '數學基礎先備', 'F1–F6', '一對一對應・唱數・數字認讀・數字排序・量比較・數的組成', Inches(0.0), Inches(4.6)),
]
for hc,lc,ser,name,units,detail,indent,sy in stairs:
    sw = Inches(12.8) - indent
    rect(sl, Inches(0.27)+indent, sy, sw, Inches(0.95),
         fill=lc, bc=hc, bpt=1.8, rad=8000)
    rect(sl, Inches(0.27)+indent, sy, Inches(1.4), Inches(0.95), fill=hc, rad=0)
    txt(sl, ser, Inches(0.27)+indent, sy+Inches(0.08), Inches(1.4), Inches(0.42),
        sz=14, bold=True, col=W, align=PP_ALIGN.CENTER)
    txt(sl, name, Inches(0.27)+indent+Inches(1.5), sy+Inches(0.08),
        Inches(2.2), Inches(0.38), sz=12, bold=True, col=hc)
    pill(sl, units, Inches(0.27)+indent+Inches(3.8), sy+Inches(0.2),
         Inches(1.0), Inches(0.3), bg=hc, fg=W, sz=10)
    txt(sl, detail, Inches(0.27)+indent+Inches(5.1), sy+Inches(0.12),
        sw-Inches(5.3), Inches(0.36), sz=10, col=DK)

# 左側箭頭標示「由下往上學習」
txt(sl, '學\n習\n方\n向\n↑', Inches(0.0), Inches(1.6), Inches(0.3), Inches(3.8),
    sz=10, bold=True, col=GR, align=PP_ALIGN.CENTER)

# 適用年段說明
rect(sl, Inches(0.3), Inches(5.82), Inches(12.7), Inches(0.5),
     fill=W, bc=LN, bpt=1.0, rad=6000)
labels = [
    (Inches(0.5), '國小中高年級', F_H),
    (Inches(3.5), '國中特教班', C_H),
    (Inches(6.5), '高中特教班', A_H),
    (Inches(9.5), '延伸應用', B_H),
]
for lx,lt,lc in labels:
    rect(sl, lx, Inches(5.9), Inches(0.12), Inches(0.28), fill=lc, rad=4000)
    txt(sl, lt, lx+Inches(0.2), Inches(5.9), Inches(2.2), Inches(0.28),
        sz=11, col=lc, bold=True)

note_box(sl, '✦ 每系列包含 6 個互動單元，共 24 個教學單元；各系列可獨立使用，亦可完整序階進行',
         Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.38))
page_num(sl, 5)

# ═══════════════════════════════════════════════════════════════
# S06  F 系列——數學基礎先備 6 單元
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, 'F 系列——數學基礎先備技能模組', '適用：國小中高年級至國中　目標：建立貨幣操作所需數學先備技能', accent=F_H)

units_f = [
    ('F1','一對一對應','🔢','拖曳物件配對\n建立數量等量直觀理解'),
    ('F2','唱數','🔊','視覺化數列\n將數量轉化為語言序列'),
    ('F3','數字認讀','👁️','數字辨識練習\n10種自訂主題圖示'),
    ('F4','數字排序','🔀','大小排列拖曳\n建立數字序列概念'),
    ('F5','量比較','⚖️','視覺化比較\n多少大小量的判斷'),
    ('F6','數的組成','🧩','分解與組合\n理解數字內部結構'),
]
for i,(code,name,ic,desc) in enumerate(units_f):
    col = i % 3
    row = i // 3
    bx = Inches(0.35 + col*4.35)
    by = Inches(1.22 + row*2.72)
    rect(sl, bx, by, Inches(4.12), Inches(2.5), fill=F_L, bc=F_H, bpt=1.5, rad=10000)
    rect(sl, bx, by, Inches(4.12), Inches(0.44), fill=F_H, rad=0)
    txt(sl, f'{code} {name}', bx+Inches(0.15), by+Inches(0.05),
        Inches(2.8), Inches(0.35), sz=14, bold=True, col=W)
    txt(sl, ic, bx+Inches(3.5), by+Inches(0.02),
        Inches(0.55), Inches(0.4), sz=20, align=PP_ALIGN.CENTER)
    txt(sl, desc, bx+Inches(0.18), by+Inches(0.6),
        Inches(3.76), Inches(1.6), sz=11.5, col=DK)

page_num(sl, 6)

# ═══════════════════════════════════════════════════════════════
# S07  C 系列——貨幣認知建立 6 單元
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, 'C 系列——貨幣認知建立模組', '適用：國中至高中　目標：建立完整的金錢認知與交易技能', accent=C_H)

units_c = [
    ('C1','認識錢幣','🪙','台灣9種面額辨識\n硬幣/紙鈔正反面記憶'),
    ('C2','數錢','🔢','多幣種混合計算\n逐枚累加的程序性技能'),
    ('C3','換錢','🔄','面額兌換練習\n小換大/大換小雙向'),
    ('C4','付款','💳','湊出指定金額\n面額組合最佳化'),
    ('C5','夠不夠','⚖️','比較金額與售價\n判斷足額與否'),
    ('C6','找零','💰','計算找零金額\n完整交易閉環驗收'),
]
for i,(code,name,ic,desc) in enumerate(units_c):
    col = i % 3
    row = i // 3
    bx = Inches(0.35 + col*4.35)
    by = Inches(1.22 + row*2.72)
    rect(sl, bx, by, Inches(4.12), Inches(2.5), fill=C_L, bc=C_H, bpt=1.5, rad=10000)
    rect(sl, bx, by, Inches(4.12), Inches(0.44), fill=C_H, rad=0)
    txt(sl, f'{code} {name}', bx+Inches(0.15), by+Inches(0.05),
        Inches(2.8), Inches(0.35), sz=14, bold=True, col=W)
    txt(sl, ic, bx+Inches(3.5), by+Inches(0.02),
        Inches(0.55), Inches(0.4), sz=20, align=PP_ALIGN.CENTER)
    txt(sl, desc, bx+Inches(0.18), by+Inches(0.6),
        Inches(3.76), Inches(1.6), sz=11.5, col=DK)

# C→A 連結說明
note_box(sl, '📌 C 系列學習完成後建議進入 A 系列——將靜態貨幣技能類化至真實情境設備操作',
         Inches(0.35), Inches(6.5), Inches(12.6), Inches(0.38))
page_num(sl, 7)

# ═══════════════════════════════════════════════════════════════
# S08  A 系列——真實情境類化 6 單元
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, 'A 系列——真實情境類化模組', '適用：國中至高中　目標：社區設備獨立操作・類化至真實生活', accent=A_H)

units_a = [
    ('A1','自動販賣機','🥤','投幣選飲料・找零取貨\n7步驟完整操作序列'),
    ('A2','理髮廳','✂️','觸控繳費機操作\n預付服務選擇付款'),
    ('A3','麥當勞','🍔','自助點餐機流程\n餐點選擇・付款・取餐'),
    ('A4','超市購物','🛒','收銀台購物結帳\n商品掃描・付款・找零'),
    ('A5','ATM 操作','🏧','提款・存款・轉帳\n金融自助設備完整流程'),
    ('A6','火車票','🚃','自動售票機購票\n車站選擇・付款・取票'),
]
for i,(code,name,ic,desc) in enumerate(units_a):
    col = i % 3
    row = i // 3
    bx = Inches(0.35 + col*4.35)
    by = Inches(1.22 + row*2.72)
    rect(sl, bx, by, Inches(4.12), Inches(2.5), fill=A_L, bc=A_H, bpt=1.5, rad=10000)
    rect(sl, bx, by, Inches(4.12), Inches(0.44), fill=A_H, rad=0)
    txt(sl, f'{code} {name}', bx+Inches(0.15), by+Inches(0.05),
        Inches(2.8), Inches(0.35), sz=14, bold=True, col=W)
    txt(sl, ic, bx+Inches(3.5), by+Inches(0.02),
        Inches(0.55), Inches(0.4), sz=20, align=PP_ALIGN.CENTER)
    txt(sl, desc, bx+Inches(0.18), by+Inches(0.6),
        Inches(3.76), Inches(1.6), sz=11.5, col=DK)

note_box(sl, '🎯 A 系列終極目標：教師 3 公尺外觀察，學生在真實社區環境完全獨立操作',
         Inches(0.35), Inches(6.5), Inches(12.6), Inches(0.38))
page_num(sl, 8)

# ═══════════════════════════════════════════════════════════════
# S09  B 系列——預算規劃延伸 6 單元
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, 'B 系列——預算規劃延伸模組', '適用：國中至高中　目標：日常預算管理與消費決策能力', accent=B_H)

units_b = [
    ('B1','今天帶多少錢','👜','依行程規劃所需金額\n場景類別・費用明細計算'),
    ('B2','零用錢日記','📒','記帳概念建立\n收支追蹤・餘額計算'),
    ('B3','存錢計畫','🐷','月曆存錢模擬\n目標設定・週存金額規劃'),
    ('B4','特賣比一比','🏷️','商品比價能力\n兩商店/三商店/單位比價'),
    ('B5','生日派對預算','🎂','預算分配決策\n必買/選購 項目規劃'),
    ('B6','菜市場買菜','🥦','採購清單執行\n多攤位・付款・找零'),
]
for i,(code,name,ic,desc) in enumerate(units_b):
    col = i % 3
    row = i // 3
    bx = Inches(0.35 + col*4.35)
    by = Inches(1.22 + row*2.72)
    rect(sl, bx, by, Inches(4.12), Inches(2.5), fill=B_L, bc=B_H, bpt=1.5, rad=10000)
    rect(sl, bx, by, Inches(4.12), Inches(0.44), fill=B_H, rad=0)
    txt(sl, f'{code} {name}', bx+Inches(0.15), by+Inches(0.05),
        Inches(2.8), Inches(0.35), sz=14, bold=True, col=W)
    txt(sl, ic, bx+Inches(3.5), by+Inches(0.02),
        Inches(0.55), Inches(0.4), sz=20, align=PP_ALIGN.CENTER)
    txt(sl, desc, bx+Inches(0.18), by+Inches(0.6),
        Inches(3.76), Inches(1.6), sz=11.5, col=DK)

page_num(sl, 9)

# ═══════════════════════════════════════════════════════════════
# S10  輔助點擊模式——首創精細動作障礙解方
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '輔助點擊模式——精細動作障礙首創解方', '全 24 單元支援，僅需「單一點擊」即可完成所有操作', accent=TE)

# 左側說明
section_label(sl, '傳統障礙', Inches(0.4), Inches(1.22), col=RD)
rect(sl, Inches(0.4), Inches(1.6), Inches(4.0), Inches(2.2),
     fill=RDL, bc=RD, bpt=1.4, rad=8000)
mtxt(sl, [
    '✗ 拖曳操作需精確手部控制',
    '✗ 觸控滑動對肢障學生困難',
    '✗ 多步驟操作序列記憶負擔重',
    '✗ 失誤容易導致情緒挫折',
], Inches(0.58), Inches(1.72), Inches(3.64), Inches(1.85),
   sz=11.5, col=RD, lh=18)

# 中間箭頭
txt(sl, '→\n輔助\n點擊\n解方', Inches(4.55), Inches(1.68), Inches(1.0), Inches(2.0),
    sz=13, bold=True, col=TE, align=PP_ALIGN.CENTER)

# 右側解方
section_label(sl, '輔助點擊解方', Inches(5.7), Inches(1.22), col=TE)
rect(sl, Inches(5.7), Inches(1.6), Inches(4.0), Inches(2.2),
     fill=TEL, bc=TE, bpt=1.4, rad=8000)
mtxt(sl, [
    '✓ 系統自動偵測點擊位置',
    '✓ 黃色邊框視覺提示下一步',
    '✓ 自動引導至正確目標',
    '✓ 全程只需「點擊」不需拖曳',
], Inches(5.88), Inches(1.72), Inches(3.64), Inches(1.85),
   sz=11.5, col=TE, lh=18)

# 適用對象
section_label(sl, '適用對象說明', Inches(9.85), Inches(1.22), col=B_H)
rect(sl, Inches(9.85), Inches(1.6), Inches(3.1), Inches(2.2),
     fill=B_L, bc=B_H, bpt=1.4, rad=8000)
mtxt(sl, [
    '• 腦性麻痺',
    '• 脊髓損傷',
    '• 肌肉萎縮症',
    '• 其他精細動作困難',
    '• 注意力不足過動症',
], Inches(10.05), Inches(1.72), Inches(2.72), Inches(1.85),
   sz=11, col=B_H, lh=18)

# 操作流程圖
section_label(sl, '輔助點擊操作流程', Inches(0.4), Inches(4.05), col=TE)
steps = ['啟用輔助點擊模式', '出現黃色提示框', '點擊任意處', '自動執行下一步', '完成任務']
for i,s in enumerate(steps):
    bx = Inches(0.4 + i*2.56)
    ellipse(sl, bx, Inches(4.45), Inches(0.45), Inches(0.45), fill=TE)
    txt(sl, str(i+1), bx, Inches(4.45), Inches(0.45), Inches(0.45),
        sz=14, bold=True, col=W, align=PP_ALIGN.CENTER)
    rect(sl, bx-Inches(0.3), Inches(5.0), Inches(1.05), Inches(0.56),
         fill=TEL, bc=TE, bpt=1.2, rad=8000)
    txt(sl, s, bx-Inches(0.3), Inches(5.04), Inches(1.05), Inches(0.48),
        sz=9.5, col=TE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, '→', Inches(1.45+i*2.56), Inches(4.5), Inches(0.5), Inches(0.35),
            sz=14, col=GR, align=PP_ALIGN.CENTER)

note_box(sl, '✦ 首創設計：市場上同類型特教軟體中，首個支援「全程單一點擊」輔助模式的教學軟體',
         Inches(0.4), Inches(5.78), Inches(12.5), Inches(0.42))
page_num(sl, 10)

# ═══════════════════════════════════════════════════════════════
# S11  三難度系統比較表
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '三段難度設計——直接教學法的系統落實', '每個單元均提供三種難度，確保從示範到獨立操作的完整進展', accent=AM)

# 表格標題行
headers = ['難度設定', '教學策略', '適用學生', '提示類型', '評量標準']
col_ws  = [Inches(1.8), Inches(2.8), Inches(2.8), Inches(2.8), Inches(2.8)]
row_h = Inches(0.42)
base_x = Inches(0.28)
base_y = Inches(1.2)
# 表頭
rect(sl, base_x, base_y, Inches(13.0), row_h, fill=HDB, rad=0)
cx = base_x
for h, cw in zip(headers, col_ws):
    txt(sl, h, cx+Inches(0.08), base_y+Inches(0.06), cw-Inches(0.12), row_h-Inches(0.1),
        sz=12, bold=True, col=W)
    cx += cw

# 三行內容
rows = [
    (C_H, C_L, ['😊 簡單模式', '（Easy）'],
     ['• 完整步驟提示\n• 視覺指引箭頭\n• 自動輔助確認'],
     ['• 初學者\n• 需要高度支持\n• 建立程序記憶'],
     ['• 黃色邊框高亮\n• 語音提示每步\n• 輔助點擊可用'],
     ['• 完成即通過\n• 無時間限制\n• IEP 最低指標']),
    (AM, AML, ['😐 普通模式', '（Normal）'],
     ['• 提示鈕按需提示\n• 3次錯誤自動提示\n• 減少自動引導'],
     ['• 中級學習者\n• 部分提示需求\n• 強化記憶階段'],
     ['• 錯誤時紅色提示\n• 按鈕觸發提示\n• 語音選擇性播放'],
     ['• 正確率計算\n• 計時評量\n• IEP 達成目標']),
    (RD, RDL, ['😤 困難模式', '（Hard）'],
     ['• 僅語音引導\n• 無視覺提示框\n• 完全靠記憶操作'],
     ['• 進階學習者\n• 準備類化至現實\n• 類化前最終確認'],
     ['• 語音提示僅一次\n• 無自動輔助\n• 錯誤需重新嘗試'],
     ['• 完整正確率\n• 最短時間紀錄\n• 類化準備指標']),
]
for ri,(hc,lc,r0,r1,r2,r3,r4) in enumerate(rows):
    ry = base_y + row_h*(ri+1) + Inches(0.05)*ri
    row_data = [r0, r1, r2, r3, r4]
    bg = W if ri%2==0 else BG2
    rect(sl, base_x, ry, Inches(13.0), Inches(1.3), fill=bg, bc=LN, bpt=1.0)
    rect(sl, base_x, ry, Inches(0.12), Inches(1.3), fill=hc)
    cx = base_x + Inches(0.12)
    # 第一欄特殊
    pill(sl, row_data[0][0], cx+Inches(0.05), ry+Inches(0.12),
         Inches(1.58), Inches(0.3), bg=hc, fg=W, sz=10.5)
    txt(sl, row_data[0][1], cx+Inches(0.05), ry+Inches(0.5),
        Inches(1.58), Inches(0.25), sz=9.5, col=hc, align=PP_ALIGN.CENTER)
    cx += col_ws[0]
    for j,d in enumerate(row_data[1:]):
        d_str = '\n'.join(d) if isinstance(d, list) else d
        mtxt(sl, d_str.split('\n') if '\n' in d_str else [d_str],
             cx+Inches(0.08), ry+Inches(0.08),
             col_ws[j+1]-Inches(0.14), Inches(1.2), sz=10, col=DK, lh=15)
        cx += col_ws[j+1]

note_box(sl, '+ 輔助點擊模式：可疊加於任何難度層級，解決精細動作障礙學生的操作問題',
         Inches(0.28), Inches(5.68), Inches(12.72), Inches(0.4))
# IEP對應
rect(sl, Inches(0.28), Inches(6.18), Inches(12.72), Inches(0.68),
     fill=RGBColor(0xE8,0xF0,0xFB), bc=HDB, bpt=1.0, rad=6000)
mtxt(sl, ['IEP 目標對應：  簡單模式 → 「在口語提示下完成」(VP)　｜　普通模式 → 「在示範提示下完成」(M)　｜　困難模式 → 「獨立完成」(I)'],
     Inches(0.5), Inches(6.28), Inches(12.2), Inches(0.48), sz=11, col=HDB)
page_num(sl, 11)

# ═══════════════════════════════════════════════════════════════
# S12  工作分析法——步驟化教學流程圖解
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '工作分析法——步驟化教學流程圖解', 'Gold 1980 / Snell & Brown 2011 系統應用　以 A1 自動販賣機為例', accent=C_H)

# A1 流程圖（7步驟）
steps_a1 = [
    (F_H, 'Step 1', '確認攜帶金額', '開啟錢包・核對面額'),
    (F_H, 'Step 2', '選擇目標飲料', '瀏覽選單・確認售價'),
    (C_H, 'Step 3', '投入硬幣/紙鈔', '面額選擇・投入確認'),
    (C_H, 'Step 4', '確認金額足夠', '螢幕顯示・比對售價'),
    (A_H, 'Step 5', '按下購買鍵', '確認按鈕・等待出貨'),
    (A_H, 'Step 6', '取出飲料', '出貨口取出・確認品項'),
    (B_H, 'Step 7', '確認找零', '找零口取出・核對金額'),
]
row1 = steps_a1[:4]
row2 = steps_a1[4:]
for row, ry in [(row1, Inches(1.32)), (row2, Inches(3.62))]:
    for i,(hc,st,nm,desc) in enumerate(row):
        bx = Inches(0.35 + i*3.24)
        rect(sl, bx, ry, Inches(3.0), Inches(1.88),
             fill=W, bc=hc, bpt=1.6, rad=8000)
        ellipse(sl, bx+Inches(0.1), ry-Inches(0.22),
                Inches(0.5), Inches(0.5), fill=hc)
        txt(sl, st, bx+Inches(0.1), ry-Inches(0.2),
            Inches(0.5), Inches(0.46), sz=10, bold=True, col=W, align=PP_ALIGN.CENTER)
        txt(sl, nm, bx+Inches(0.15), ry+Inches(0.12),
            Inches(2.7), Inches(0.38), sz=12, bold=True, col=hc)
        txt(sl, desc, bx+Inches(0.15), ry+Inches(0.56),
            Inches(2.7), Inches(1.12), sz=11, col=DK)
        if i < len(row)-1:
            txt(sl, '→', Inches(3.48+i*3.24), ry+Inches(0.72),
                Inches(0.35), Inches(0.38), sz=16, col=GR, align=PP_ALIGN.CENTER)

# 連接兩行
txt(sl, '↓', Inches(0.35+0*3.24)+Inches(1.2), Inches(3.28),
    Inches(0.5), Inches(0.3), sz=16, col=GR, align=PP_ALIGN.CENTER)

# 評量記錄說明
section_label(sl, '步驟評量記錄格式', Inches(0.4), Inches(5.68), col=C_H)
headers2 = ['步驟', '行為描述', 'I 獨立', 'VP 口語提示', 'PP 部分肢體', 'FP 全程協助', '備注']
col_ws2  = [Inches(0.65), Inches(3.2), Inches(0.85), Inches(1.35), Inches(1.35), Inches(1.35), Inches(1.8)]
rect(sl, Inches(0.4), Inches(6.05), sum(col_ws2), Inches(0.36), fill=C_H)
cx2 = Inches(0.4)
for h2,cw2 in zip(headers2, col_ws2):
    txt(sl, h2, cx2+Inches(0.04), Inches(6.08), cw2, Inches(0.28),
        sz=10, bold=True, col=W, align=PP_ALIGN.CENTER)
    cx2 += cw2
note_box(sl, '✦ 每步驟均可獨立評量、單獨教學，完整對應 IEP 短期目標行為敘述格式',
         Inches(0.4), Inches(6.52), Inches(12.5), Inches(0.38))
page_num(sl, 12)

# ═══════════════════════════════════════════════════════════════
# S13  融合教育三層支持架構
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '融合教育三層支持架構', '相同教材・相同目標方向，差異在「支持程度」，實現最少限制環境原則', accent=C_H)

# 三層金字塔（從下到上）
layers = [
    (F_H, F_L, 'Level 1  普通班同步參與路徑',
     Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.05),
     ['• 特殊需求學生在普通班數學課中以輔助點擊模式參與相同任務',
      '• 軟體自動提供支持，教師不須全程個別介入，實現無縫融合',
      '• 普通班同儕同步了解如何給予適當協助，雙向融合目標同步實現']),
    (C_H, C_L, 'Level 2  資源班強化練習路徑',
     Inches(1.8), Inches(4.45), Inches(9.7), Inches(1.05),
     ['• 資源班抽離課使用，針對個別化弱點進行系統性強化',
      '• 可在結構化環境中重複練習，不受普通班進度影響',
      '• 教師可即時調整難度，精確對應 IEP 目標']),
    (A_H, A_L, 'Level 3  自足班精熟學習路徑',
     Inches(3.2), Inches(3.1), Inches(6.9), Inches(1.05),
     ['• 自足式特殊教育班完整課程使用',
      '• 依學生程度客製化難度與學習路徑',
      '• 作業單輔助評量，系統化記錄進步']),
]
for hc,lc,title,lx,ly,lw,lh,pts in layers:
    rect(sl, lx, ly, lw, lh, fill=lc, bc=hc, bpt=1.8, rad=6000)
    txt(sl, title, lx+Inches(0.2), ly+Inches(0.05), lw-Inches(0.4), Inches(0.35),
        sz=12, bold=True, col=hc)

# 右側展開說明
for i,(hc,lc,title,lx,ly,lw,lh,pts) in enumerate(layers):
    ey = Inches(1.2 + i*1.95)
    rect(sl, Inches(0.4), ey, Inches(12.5), Inches(1.72),
         fill=lc, bc=hc, bpt=1.4, rad=8000)
    txt(sl, title, Inches(0.6), ey+Inches(0.1), Inches(12.1), Inches(0.35),
        sz=12, bold=True, col=hc)
    mtxt(sl, pts, Inches(0.6), ey+Inches(0.44),
         Inches(12.1), Inches(1.1), sz=10.5, col=DK, lh=16)

note_box(sl, '法規依據：《特殊教育法》第 28 條「依個別需求訂定 IEP」・最少限制環境（LRE）原則',
         Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.38))
page_num(sl, 13)

# ═══════════════════════════════════════════════════════════════
# S14  IEP 課綱法規對應
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, 'IEP 目標・課綱・法規三對應', '課程設計完整對應國家課程綱要與特殊教育法規', accent=HDB)

# 三欄對應表
cols3 = [
    (F_H, F_L, '十二年國教課綱對應',
     [('F系列', '生管-Ⅳ-E-1', '認識台灣常用貨幣'),
      ('C系列', '生管-Ⅳ-E-2', '使用金錢購物'),
      ('C5/C6', '生管-Ⅳ-E-3', '了解找零概念'),
      ('A系列', '生管-Ⅳ-E-4', '社區消費技能'),
      ('A5',    '生管-Ⅳ-F-1', '使用自動化設備'),
      ('B系列', '生管-Ⅳ-G-1', '財務管理基礎')]),
    (C_H, C_L, 'IEP 年度目標範本',
     [('F1-F2', '數學先備', '在口語提示下完成10以內的一對一對應操作'),
      ('C1-C3', '貨幣認知', '獨立辨識台灣六種硬幣並說出面額'),
      ('C4-C6', '金錢操作', '在部分協助下完成付款並核對找零'),
      ('A1',    '類化應用', '獨立完成自動販賣機購買飲料的完整流程'),
      ('A5',    '類化應用', '在口語提示下完成 ATM 提款操作'),
      ('B1-B2', '財務管理', '能依行程計算所需金額並準備正確零錢')]),
    (A_H, A_L, '法規依據',
     [('特殊教育法', '§28',    '依個別需求訂定個別化教育計畫（IEP）'),
      ('特殊教育法', '§24',    '提供必要之輔助科技及支援服務'),
      ('課綱', '2019',         '特殊需求領域生活管理：財務管理核心主題'),
      ('CRPD',  'Art.24',      '融合教育・平等受教機會保障'),
      ('UDL',   'CAST 2018',   '通用設計學習三原則系統落實'),
      ('工作分析', 'Gold 1980','步驟化教學完整對應行為目標格式')]),
]
for ci,(hc,lc,title,items) in enumerate(cols3):
    bx = Inches(0.32 + ci*4.35)
    cw = Inches(4.12)
    rect(sl, bx, Inches(1.22), cw, Inches(5.62),
         fill=W, bc=hc, bpt=1.6, rad=8000)
    rect(sl, bx, Inches(1.22), cw, Inches(0.44), fill=hc, rad=0)
    txt(sl, title, bx+Inches(0.15), Inches(1.28), cw-Inches(0.3), Inches(0.36),
        sz=12, bold=True, col=W)
    for ri,(a,b,c) in enumerate(items):
        ry = Inches(1.78 + ri*0.8)
        rect(sl, bx+Inches(0.15), ry, cw-Inches(0.3), Inches(0.72),
             fill=lc, bc=hc, bpt=0.8, rad=6000)
        pill(sl, a, bx+Inches(0.2), ry+Inches(0.1),
             Inches(0.78), Inches(0.22), bg=hc, fg=W, sz=8.5)
        pill(sl, b, bx+Inches(1.05), ry+Inches(0.1),
             Inches(0.78), Inches(0.22), bg=W, fg=hc, sz=8.5)
        txt(sl, c, bx+Inches(0.2), ry+Inches(0.37),
            cw-Inches(0.55), Inches(0.3), sz=9.5, col=DK)
page_num(sl, 14)

# ═══════════════════════════════════════════════════════════════
# S15  試教觀察——三個個案類化轉折
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '試教觀察——三個個案・類化窗口的關鍵時刻', '肆、試教觀察：學生在第 5–7 次操作後出現程序性記憶語言化信號', accent=RD)

cases2 = [
    (RD, RDL, '個案一  小涵', '14歲 ・輕度智能障礙',
     ['試教前：', '✗ 計算 100-35=65 正確', '✗ 實際付款卻反覆困惑', '✗ 無法連結「我的錢」與「售價」'],
     ['試教後（第7次操作）：', '✓ 主動拿出零錢逐枚核對', '✓ 說出「這個要35元，我有」', '✓ 程序性記憶語言化出現'],
     '突破點：步驟化流程讓「知道計算」\n轉化為「理解交易序列意義」'),
    (C_H, C_L, '個案二  小俊', '輕度自閉症',
     ['試教前：', '✗ 速食店選擇迴路情緒失調', '✗ 一學期迴避相關學習', '✗ 聽到「麥當勞」即迴避行為'],
     ['A3 試教→社區實地（第3節）：', '✓ 面對真實觸控機完成完整流程', '✓ 主動說：「我是7號，還沒叫到我」', '✓ 等待焦慮顯著改善'],
     '突破點：安全練習環境建立信心\n類化窗口在第5次後自然開啟'),
    (A_H, A_L, '個案三  小雯', '販賣機迴避行為',
     ['試教前：', '✗ 多次因硬幣選錯而失敗', '✗ 繞道迴避所有自助設備', '✗ 告訴教師「那個機器我不會用」'],
     ['A1 完成 10 次練習後：', '✓ 主動走向走廊販賣機', '✓ 獨立完成選飲料和投幣', '✓ 情境迴避行為完全消失'],
     '突破點：無限次安全嘗試消除\n「失敗的社會後果」焦慮'),
]
for i,(hc,lc,name,tag,before,after,note) in enumerate(cases2):
    bx = Inches(0.32 + i*4.35)
    cw = Inches(4.12)
    rect(sl, bx, Inches(1.22), cw, Inches(5.58),
         fill=W, bc=hc, bpt=1.8, rad=10000)
    rect(sl, bx, Inches(1.22), cw, Inches(0.5), fill=hc, rad=0)
    txt(sl, name, bx+Inches(0.15), Inches(1.28), cw-Inches(0.3), Inches(0.38),
        sz=14, bold=True, col=W)
    pill(sl, tag, bx+Inches(0.15), Inches(1.82),
         cw-Inches(0.3), Inches(0.28), bg=lc, fg=hc, sz=9.5, bold=False)
    # 試教前後對比
    rect(sl, bx+Inches(0.12), Inches(2.2), cw-Inches(0.24), Inches(1.62),
         fill=RGBColor(0xFF,0xF5,0xF5), bc=RD, bpt=0.8, rad=6000)
    mtxt(sl, before, bx+Inches(0.22), Inches(2.28),
         cw-Inches(0.4), Inches(1.42), sz=10, col=DK, lh=15)
    rect(sl, bx+Inches(0.12), Inches(3.9), cw-Inches(0.24), Inches(1.62),
         fill=RGBColor(0xF0,0xFB,0xF4), bc=C_H, bpt=0.8, rad=6000)
    mtxt(sl, after, bx+Inches(0.22), Inches(3.98),
         cw-Inches(0.4), Inches(1.42), sz=10, col=DK, lh=15)
    note_box(sl, note, bx+Inches(0.12), Inches(5.62),
             cw-Inches(0.24), Inches(0.6))
page_num(sl, 15)

# ═══════════════════════════════════════════════════════════════
# S16  預期成效——三層面×短中長期
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '預期成效——三層面×短中長期指標', '肆、學生・教師・學校三層面成效，建立完整評估框架', accent=C_H)

layers2 = [
    (F_H, F_L, '🎓 學生層面', [
        ('短期（1學期）', '完成 F/C 系列 80% 以上單元，建立貨幣辨識與付款基礎技能'),
        ('中期（1學年）', '完成 A 系列 4 個以上單元，達到「部分提示下完成」社區操作'),
        ('長期（跨學年）', '在真實社區環境獨立完成 ATM 操作・超市購物・自助點餐等日常生活技能'),
    ]),
    (C_H, C_L, '👩‍🏫 教師層面', [
        ('短期（1學期）', '熟悉軟體操作，將 IEP 目標有效對應各系列單元，縮短備課時間 50% 以上'),
        ('中期（1學年）', '建立系統化評量記錄，追蹤學生步驟別進步，提供具體 IEP 報告數據'),
        ('長期（跨學年）', '累積班級學習資料，形成校本課程材料，分享予特教社群推廣應用'),
    ]),
    (A_H, A_L, '🏫 學校層面', [
        ('短期（1學期）', '提供零成本數位教材，補充現有紙本教材的互動回饋不足'),
        ('中期（1學年）', '建立特教班數位學習環境，支援多元評量實施'),
        ('長期（跨學年）', '成為區域特教教材分享典範，連結融合教育推廣政策目標'),
    ]),
]
for ri,(hc,lc,title,items) in enumerate(layers2):
    ry = Inches(1.18 + ri*1.98)
    rect(sl, Inches(0.4), ry, Inches(12.5), Inches(1.82),
         fill=lc, bc=hc, bpt=1.6, rad=8000)
    rect(sl, Inches(0.4), ry, Inches(1.6), Inches(1.82), fill=hc, rad=0)
    txt(sl, title, Inches(0.4), ry+Inches(0.7), Inches(1.6), Inches(0.38),
        sz=12, bold=True, col=W, align=PP_ALIGN.CENTER)
    for ti,(period,content) in enumerate(items):
        tx = Inches(2.2 + ti*3.55)
        rect(sl, tx, ry+Inches(0.1), Inches(3.3), Inches(1.6),
             fill=W, bc=hc, bpt=0.8, rad=6000)
        pill(sl, period, tx+Inches(0.1), ry+Inches(0.18),
             Inches(3.1), Inches(0.26), bg=hc, fg=W, sz=9.5)
        txt(sl, content, tx+Inches(0.1), ry+Inches(0.54),
            Inches(3.1), Inches(1.0), sz=10, col=DK)
page_num(sl, 16)

# ═══════════════════════════════════════════════════════════════
# S17  技術規格與作業單系統
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '技術規格與配套系統', '二、無障礙設計・作業單系統・獎勵系統', accent=F_H)

# 左側：技術規格
section_label(sl, '技術規格', Inches(0.4), Inches(1.22), col=F_H)
specs = [
    ('平台','HTML5 + 純 JavaScript（無框架依賴）'),
    ('部署','完全離線・免安裝・免伺服器'),
    ('授權','NT$ 0（授權費用）・完整原始碼開放'),
    ('規模','150,000+ 行程式碼・24個教學單元'),
    ('設備','平板・電腦・手機全相容'),
    ('音訊','Web Speech API 繁體中文語音合成'),
    ('視覺','雅婷語音優先・2倍縮放（視覺障礙）'),
    ('評量','作業單 PDF 下載・獎勵系統整合'),
]
for i,(k,v) in enumerate(specs):
    ry = Inches(1.62 + i*0.58)
    bg = F_L if i%2==0 else W
    rect(sl, Inches(0.4), ry, Inches(5.8), Inches(0.5),
         fill=bg, bc=LN, bpt=0.8, rad=4000)
    txt(sl, k, Inches(0.55), ry+Inches(0.1), Inches(1.0), Inches(0.3),
        sz=11, bold=True, col=F_H)
    txt(sl, v, Inches(1.7), ry+Inches(0.1), Inches(4.3), Inches(0.3),
        sz=11, col=DK)

# 右側：作業單系統
section_label(sl, '作業單系統', Inches(6.55), Inches(1.22), col=A_H)
ws_features = [
    ('🖨️', '一鍵 PDF 生成', '即時產生對應單元的個人化作業單'),
    ('📋', '智慧題目生成', '自動依難度調整題型與數量'),
    ('🖼️', '圖文並茂版面', '錢幣圖示 + 看圖填空的多元題型'),
    ('🎯', '答案版/練習版', '教師/學生雙版本一鍵切換'),
    ('🏪', '商店客製', 'A4 超市依商店類型生成對應題目'),
    ('🏆', '獎勵系統整合', '完成後解鎖貼紙・累積學習成就'),
]
for i,(ic,t,d) in enumerate(ws_features):
    col = i % 2
    row = i // 2
    bx = Inches(6.55 + col*3.3)
    ry = Inches(1.62 + row*1.58)
    rect(sl, bx, ry, Inches(3.1), Inches(1.38),
         fill=A_L, bc=A_H, bpt=1.2, rad=8000)
    txt(sl, ic, bx+Inches(0.1), ry+Inches(0.1), Inches(0.5), Inches(0.45),
        sz=18, align=PP_ALIGN.CENTER)
    txt(sl, t, bx+Inches(0.65), ry+Inches(0.12), Inches(2.3), Inches(0.3),
        sz=11, bold=True, col=A_H)
    txt(sl, d, bx+Inches(0.18), ry+Inches(0.58), Inches(2.82), Inches(0.68),
        sz=10, col=DK)
page_num(sl, 17)

# ═══════════════════════════════════════════════════════════════
# S18  數位→紙本評量→社區類化 三段學習旅程
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '三段完整學習旅程', '數位互動練習 → 紙本作業評量 → 社區真實類化', accent=B_H)

stages = [
    (F_H, F_L, '🖥️ 第一段', '數位互動練習',
     ['課堂使用軟體反覆練習', '三難度漸進適應個別程度', '即時語音回饋強化記憶',
      '作業單系統記錄進步', '輔助點擊支援精細動作障礙'],
     '重點：建立程序性記憶\n允許無限次安全嘗試'),
    (C_H, C_L, '📝 第二段', '紙本作業評量',
     ['一鍵生成 PDF 作業單', '看圖付款/找零多元題型', '學生版/教師答案版切換',
      '形成性評量紀錄', '步驟別評量對應 IEP'],
     '重點：形成性評量\n確認技能遷移程度'),
    (A_H, A_L, '🏘️ 第三段', '社區真實類化',
     ['教師遠距觀察（3公尺外）', '逐步褪除支持提示', '真實設備實地操作演練',
      '社區技能類化確認', '最終目標：完全獨立操作'],
     '重點：類化窗口開啟\n終極成效確認'),
]
for i,(hc,lc,stage,title,pts,note) in enumerate(stages):
    bx = Inches(0.35 + i*4.35)
    cw = Inches(4.12)
    # 上方階段標示
    rect(sl, bx, Inches(1.2), cw, Inches(0.5), fill=hc, rad=0)
    txt(sl, stage, bx+Inches(0.15), Inches(1.26), cw-Inches(0.3), Inches(0.38),
        sz=13, bold=True, col=W)
    # 卡片主體
    rect(sl, bx, Inches(1.7), cw, Inches(3.8), fill=lc, bc=hc, bpt=1.6, rad=8000)
    txt(sl, title, bx+Inches(0.15), Inches(1.82), cw-Inches(0.3), Inches(0.38),
        sz=14, bold=True, col=hc)
    hline(sl, Inches(2.28), bx+Inches(0.15), bx+cw-Inches(0.15), col=hc, pt=1.0)
    mtxt(sl, ['• '+p for p in pts], bx+Inches(0.18), Inches(2.36),
         cw-Inches(0.3), Inches(2.0), sz=11, col=DK, lh=17)
    note_box(sl, note, bx+Inches(0.15), Inches(5.62), cw-Inches(0.3), Inches(0.56))
    # 箭頭
    if i < 2:
        txt(sl, '⟹', Inches(4.6+i*4.35), Inches(2.8),
            Inches(0.55), Inches(0.5), sz=22, col=GR, align=PP_ALIGN.CENTER)

note_box(sl, '✦ 三段設計呼應《特殊教育法》第 24 條「提供輔助科技」及課綱「社區功能性技能」學習目標',
         Inches(0.35), Inches(6.45), Inches(12.6), Inches(0.38))
page_num(sl, 18)

# ═══════════════════════════════════════════════════════════════
# S19  F系列詳細教案摘要
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, 'F 系列教案摘要——數學基礎先備技能模組', '教案一：F1 一對一對應・F2 唱數（共 3 節）', accent=F_H)

# 教案基本資訊
rect(sl, Inches(0.4), Inches(1.18), Inches(12.5), Inches(1.05),
     fill=F_L, bc=F_H, bpt=1.6, rad=8000)
info_cols = [('教學對象','國中特教班・資源班（輕中度智障・自閉症）'),
             ('學習階段','國中（7-9年級）'),
             ('教學時間','3節×45分鐘'),
             ('先備能力','基本物品辨識')]
for i,(k,v) in enumerate(info_cols):
    bx = Inches(0.6 + i*3.1)
    txt(sl, k+'：', bx, Inches(1.28), Inches(1.2), Inches(0.3),
        sz=10, bold=True, col=F_H)
    txt(sl, v, bx, Inches(1.58), Inches(2.8), Inches(0.5),
        sz=10, col=DK)

# 三節課程說明
lessons = [
    (F_H, '第一節（45分鐘）', 'F1 一對一對應——建立數量等量直觀',
     ['導入：呈現實物（蘋果配碗、筆配本）引發生活連結',
      '示範：教師操作 F1，展示拖曳配對步驟',
      '引導：輔助點擊模式，學生逐步完成配對任務',
      '獨立：切換普通/困難模式自行嘗試',
      '評量：記錄步驟完成情況（I / VP / PP / FP）']),
    (C_H, '第二節（45分鐘）', 'F2 唱數——數量概念轉化為語言序列',
     ['複習：回顧 F1 操作，連結「幾個就說幾」',
      '示範：操作 F2，展示逐步唱數視覺化',
      '引導：語音輔助模式，跟著說出數字',
      '獨立：關閉提示音，自行唱數完成',
      '評量：錄製語音唱數樣本作為基線']),
    (A_H, '第三節（45分鐘）', '整合評量與形成性評量',
     ['整合：F1+F2 混合任務，連結貨幣數量概念',
      '評量：使用作業單（PDF）進行紙本評量',
      '融合：普通班同儕參與，互相示範操作',
      '記錄：步驟別評量表，更新 IEP 短期目標',
      '回饋：語音鼓勵+獎勵系統貼紙']),
]
for i,(hc,title,sub,pts) in enumerate(lessons):
    bx = Inches(0.4 + i*4.35)
    cw = Inches(4.12)
    rect(sl, bx, Inches(2.42), cw, Inches(4.42),
         fill=W, bc=hc, bpt=1.6, rad=8000)
    rect(sl, bx, Inches(2.42), cw, Inches(0.5), fill=hc, rad=0)
    txt(sl, title, bx+Inches(0.12), Inches(2.48), cw-Inches(0.24), Inches(0.38),
        sz=12, bold=True, col=W)
    txt(sl, sub, bx+Inches(0.12), Inches(3.0), cw-Inches(0.24), Inches(0.34),
        sz=11, bold=True, col=hc)
    mtxt(sl, pts, bx+Inches(0.18), Inches(3.4), cw-Inches(0.3), Inches(3.2),
         sz=10.5, col=DK, lh=16)
page_num(sl, 19)

# ═══════════════════════════════════════════════════════════════
# S20  A 系列教案摘要
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, 'A 系列教案摘要——真實情境類化模組', '教案三：A1 自動販賣機（核心融合節次）・A5 ATM 操作', accent=A_H)

# A1 教案
rect(sl, Inches(0.4), Inches(1.22), Inches(7.8), Inches(5.56),
     fill=A_L, bc=A_H, bpt=1.8, rad=10000)
txt(sl, 'A1 自動販賣機教案（3 節）', Inches(0.6), Inches(1.32),
    Inches(7.4), Inches(0.38), sz=14, bold=True, col=A_H)

a1_nodes = [
    (Inches(0.6), Inches(1.82), '第一節', '流程建立',
     ['工作分析 7 步驟逐一教導',
      '輔助點擊模式全程引導',
      '語音提示每步驟操作要點']),
    (Inches(4.1), Inches(1.82), '第二節', '投幣找零強化',
     ['面額選擇策略訓練',
      '找零驗收程序建立',
      '三難度循序挑戰']),
    (Inches(0.6), Inches(3.82), '第三節', '融合教育頂點🏆',
     ['邀請普通班同儕擔任「夥伴」',
      '在校園真實販賣機操作',
      '教師 3 公尺外靜態觀察',
      '→ 類化窗口開啟的關鍵節次']),
]
for bx,by,title,sub,pts in a1_nodes:
    rect(sl, bx, by, Inches(3.25), Inches(1.68),
         fill=W, bc=A_H, bpt=1.0, rad=6000)
    txt(sl, title, bx+Inches(0.12), by+Inches(0.08), Inches(1.2), Inches(0.28),
        sz=11, bold=True, col=A_H)
    txt(sl, sub, bx+Inches(0.12), by+Inches(0.35), Inches(3.0), Inches(0.26),
        sz=10.5, col=DK, bold=True)
    mtxt(sl, pts, bx+Inches(0.15), by+Inches(0.68),
         Inches(3.0), Inches(0.9), sz=10, col=DK, lh=15)

# A5 教案摘要
rect(sl, Inches(8.45), Inches(1.22), Inches(4.5), Inches(5.56),
     fill=B_L, bc=B_H, bpt=1.8, rad=10000)
txt(sl, 'A5 ATM 操作精要教案', Inches(8.65), Inches(1.32),
    Inches(4.1), Inches(0.36), sz=13, bold=True, col=B_H)
a5_pts = [
    ('教學時間', '3節×45分鐘・高中特教班'),
    ('學習目標', '獨立完成提款操作（I 獨立）'),
    ('', '部分提示下完成存款/轉帳'),
    ('第一節', '介面認識・提款流程建立'),
    ('第二節', '存款・轉帳流程'),
    ('第三節', '融合情境：真實 ATM 操作演練'),
    ('融合設計', '普通班同學擔任「等候示範者」'),
    ('', '特殊需求學生自主完成操作'),
    ('評量', '5步驟評量表・IEP 轉銜目標'),
]
for i,(k,v) in enumerate(a5_pts):
    ry = Inches(1.78 + i*0.58)
    if k:
        txt(sl, k+'：', Inches(8.65), ry, Inches(1.2), Inches(0.3),
            sz=10.5, bold=True, col=B_H)
    txt(sl, v, Inches(9.75) if k else Inches(8.88), ry,
        Inches(3.0), Inches(0.3), sz=10, col=DK)

note_box(sl, '🔗 A 系列社區轉銜設計遵循「最少限制環境」原則：軟體練習→褪除支持→真實社區完全獨立',
         Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.38))
page_num(sl, 20)

# ═══════════════════════════════════════════════════════════════
# S21  核心創新——四大首創特色
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
top_bar(sl, '核心創新——四大首創特色', '填補台灣特殊教育金融素養教材的十年空白', accent=B_H)

innovations = [
    (B_H, B_L, '🖱️', '首創輔助點擊模式',
     '精細動作障礙解方',
     ['全 24 個教學單元均支援', '單一點擊完成全部操作', '不需拖曳・不需精確觸控',
      '市場同類型特教軟體首例', ''],
     '適用：腦性麻痺・脊損・肌萎・ADHD'),
    (A_H, A_L, '🏪', '首創全場景仿真',
     '6種真實社區設備模擬',
     ['自動販賣機・理髮廳・麥當勞', 'ATM・超市收銀台・火車售票機',
      '從像素到流程忠實仿真', '降低真實場域焦慮的過渡橋梁', ''],
     '目標：類化轉移・社區獨立操作'),
    (C_H, C_L, '📚', '首創完整序階課程',
     '24單元 × 4系列系統設計',
     ['F（數學先備）→ C（貨幣認知）', '→ A（情境類化）→ B（預算規劃）',
      '全台首個完整金融素養課程套件', '完整對應課綱所有學習主題', ''],
     '覆蓋：國小至高中完整學習序階'),
    (F_H, F_L, '🎯', '首創IEP整合設計',
     '步驟評量→作業單→獎勵',
     ['工作分析法步驟化評量格式', 'PDF 作業單一鍵生成', '獎勵系統強化學習動機',
      '教師行政負擔降低 50%+', ''],
     '整合：IEP目標・形成性評量・激勵'),
]
for i,(hc,lc,ic,title,sub,pts,note) in enumerate(innovations):
    bx = Inches(0.35 + i*3.25)
    cw = Inches(3.0)
    # 頂部icon徽章
    ellipse(sl, bx+Inches(1.0), Inches(1.15), Inches(1.0), Inches(1.0), fill=hc)
    txt(sl, ic, bx+Inches(1.0), Inches(1.18), Inches(1.0), Inches(0.94),
        sz=28, align=PP_ALIGN.CENTER)
    rect(sl, bx, Inches(2.22), cw, Inches(4.78),
         fill=lc, bc=hc, bpt=1.8, rad=10000)
    txt(sl, title, bx+Inches(0.12), Inches(2.32), cw-Inches(0.24), Inches(0.38),
        sz=12.5, bold=True, col=hc, align=PP_ALIGN.CENTER)
    txt(sl, sub, bx+Inches(0.12), Inches(2.7), cw-Inches(0.24), Inches(0.28),
        sz=10, col=GR, align=PP_ALIGN.CENTER)
    hline(sl, Inches(3.06), bx+Inches(0.2), bx+cw-Inches(0.2), col=hc, pt=0.8)
    mtxt(sl, ['▶ '+p for p in pts if p], bx+Inches(0.15), Inches(3.12),
         cw-Inches(0.3), Inches(2.0), sz=10.5, col=DK, lh=17)
    note_box(sl, note, bx+Inches(0.12), Inches(5.55), cw-Inches(0.24), Inches(0.52))
page_num(sl, 21)

# ═══════════════════════════════════════════════════════════════
# S22  結語
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
# 左側深藍欄
rect(sl, Inches(0), Inches(0), Inches(4.5), SLIDE_H, fill=HDB)
for i,(c,frac) in enumerate([(F_H,0.2),(C_H,0.4),(A_H,0.6),(B_H,0.8),(TE,1.0)]):
    rect(sl, Inches(0), Inches(7.5*frac-0.08), Inches(4.5), Inches(0.1), fill=c)

txt(sl, '感謝聆聽', Inches(0.2), Inches(1.2), Inches(4.1), Inches(0.9),
    sz=32, bold=True, col=W, align=PP_ALIGN.CENTER)
txt(sl, '金錢小達人', Inches(0.2), Inches(2.2), Inches(4.1), Inches(0.7),
    sz=22, bold=True, col=RGBColor(0xFD,0xBA,0x74), align=PP_ALIGN.CENTER)
txt(sl, '互動式特殊教育\n生活數學暨金融素養教學軟體',
    Inches(0.2), Inches(2.95), Inches(4.1), Inches(0.8),
    sz=13, col=RGBColor(0x9B,0xB8,0xD8), align=PP_ALIGN.CENTER)
txt(sl, '第八屆教育部獎助研發\n特殊教育教材教具競賽',
    Inches(0.2), Inches(3.95), Inches(4.1), Inches(0.6),
    sz=11, col=RGBColor(0x7A,0x98,0xBA), align=PP_ALIGN.CENTER)
txt(sl, '中華民國 115 年 3 月',
    Inches(0.2), Inches(6.8), Inches(4.1), Inches(0.32),
    sz=10, col=RGBColor(0x5A,0x78,0x9A), align=PP_ALIGN.CENTER)

# 右側重點摘要
section_label(sl, '核心訴求', Inches(4.9), Inches(0.55), col=B_H)
core_msgs = [
    (B_H, B_L, '🎯', '填補十年空白',
     '全台首個完整系統化特殊教育金融素養\n互動教學軟體套件（24個教學單元）'),
    (A_H, A_L, '🤝', '零門檻推廣',
     'NT$ 0 授權・完全離線・免安裝\n讓每位特教教師立即可用'),
    (C_H, C_L, '🌱', '真實成效',
     '三個個案的類化轉折紀錄\n從情境迴避到社區獨立操作的真實突破'),
    (F_H, F_L, '📋', '完整支持',
     'IEP 整合・作業單・獎勵系統・融合教育\n提供教師全方位教學支援'),
]
for i,(hc,lc,ic,title,desc) in enumerate(core_msgs):
    col = i % 2
    row = i // 2
    bx = Inches(4.9 + col*4.15)
    ry = Inches(0.98 + row*3.1)
    rect(sl, bx, ry, Inches(3.9), Inches(2.8), fill=lc, bc=hc, bpt=1.8, rad=10000)
    txt(sl, ic, bx+Inches(0.2), ry+Inches(0.25), Inches(0.7), Inches(0.6),
        sz=26, align=PP_ALIGN.CENTER)
    txt(sl, title, bx+Inches(1.0), ry+Inches(0.3), Inches(2.75), Inches(0.4),
        sz=14, bold=True, col=hc)
    hline(sl, ry+Inches(0.8), bx+Inches(0.18), bx+Inches(3.72), col=hc, pt=0.8)
    txt(sl, desc, bx+Inches(0.18), ry+Inches(0.95), Inches(3.55), Inches(1.6),
        sz=11, col=DK)

note_box(sl, '✦ 本軟體程式碼規模逾 15 萬行，完全由教師個人研發，零外部資金，零商業目的，只為讓特殊需求學生平等享有數位學習機會',
         Inches(4.9), Inches(7.07), Inches(8.15), Inches(0.36))
page_num(sl, 22)

# ── 儲存 ────────────────────────────────────────────────────
out = '金錢小達人_競賽申請簡報_V4.pptx'
prs.save(out)
print(f'✓ 已儲存：{out}')
