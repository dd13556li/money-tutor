"""
Claude Code 15 項功能 PPT 製作腳本
新手完整操作指南（繁體中文）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 色彩常數 ────────────────────────────────────────────────
C_PURPLE      = RGBColor(0x5B, 0x4F, 0xCF)   # Claude 主色
C_DARK_PURPLE = RGBColor(0x3B, 0x30, 0x9A)   # 深紫
C_AMBER       = RGBColor(0xF5, 0x9E, 0x0B)   # 琥珀（標題列）
C_GREEN       = RGBColor(0x10, 0xB9, 0x81)   # 綠色（成功/步驟）
C_BLUE        = RGBColor(0x3B, 0x82, 0xF6)   # 藍色（功能說明）
C_RED         = RGBColor(0xEF, 0x44, 0x44)   # 紅色（注意）
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK        = RGBColor(0x1F, 0x29, 0x37)
C_GRAY_BG     = RGBColor(0xF8, 0xF9, 0xFA)
C_LIGHT_PURPLE= RGBColor(0xED, 0xE9, 0xFF)
C_LIGHT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
C_LIGHT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)
C_LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)
C_CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)   # 深色代碼背景

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout

# ── 輔助函數 ─────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    if fill_color:
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, font_name='微軟正黑體',
                wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_text_para(tf, text, font_size=14, bold=False, color=C_DARK,
                  align=PP_ALIGN.LEFT, font_name='微軟正黑體', space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p

def slide_header(slide, title, subtitle='', bg_color=C_PURPLE):
    """標準標題列（頂部彩色橫條）"""
    add_rect(slide, 0, 0, 13.33, 1.2, fill_color=bg_color)
    add_textbox(slide, title, 0.3, 0.12, 10, 0.65,
                font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.75, 11, 0.42,
                    font_size=14, bold=False, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.LEFT)
    # 頁面背景淺灰
    add_rect(slide, 0, 1.2, 13.33, 6.3, fill_color=C_GRAY_BG)

def badge(slide, text, left, top, width=1.4, height=0.38,
          bg=C_PURPLE, fg=C_WHITE, font_size=13):
    r = add_rect(slide, left, top, width, height, fill_color=bg)
    add_textbox(slide, text, left, top, width, height,
                font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)

def code_box(slide, code_text, left, top, width, height, font_size=13):
    """深色代碼框"""
    add_rect(slide, left, top, width, height, fill_color=C_CODE_BG)
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.1),
        Inches(width - 0.3), Inches(height - 0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for i, line in enumerate(code_text.strip().split('\n')):
        if i == 0:
            run = p.add_run()
        else:
            p = tf.add_paragraph()
            run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = 'Consolas'
        run.font.color.rgb = RGBColor(0xA8, 0xCC, 0x8C)

def info_box(slide, title, content_lines, left, top, width, height,
             bg=C_LIGHT_BLUE, title_color=C_BLUE, icon='💡'):
    add_rect(slide, left, top, width, height, fill_color=bg)
    # title
    add_textbox(slide, f'{icon} {title}', left + 0.12, top + 0.1,
                width - 0.24, 0.38, font_size=13, bold=True, color=title_color)
    # content
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.48),
        Inches(width - 0.24), Inches(height - 0.58)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            run = tf.paragraphs[0].add_run()
        else:
            p = tf.add_paragraph()
            run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.name = '微軟正黑體'
        run.font.color.rgb = C_DARK

def step_box(slide, steps, left, top, width, height, bg=C_LIGHT_GREEN):
    """步驟清單框"""
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_textbox(slide, '📋 操作步驟', left + 0.12, top + 0.08,
                width - 0.24, 0.35, font_size=13, bold=True, color=C_GREEN)
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.45),
        Inches(width - 0.24), Inches(height - 0.55)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f'  {step}'
        run.font.size = Pt(12)
        run.font.name = '微軟正黑體'
        run.font.color.rgb = C_DARK

def example_box(slide, title, content_lines, left, top, width, height, bg=C_LIGHT_AMBER):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_textbox(slide, f'🌟 {title}', left + 0.12, top + 0.08,
                width - 0.24, 0.35, font_size=13, bold=True, color=C_AMBER)
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.45),
        Inches(width - 0.24), Inches(height - 0.55)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.name = '微軟正黑體'
        run.font.color.rgb = C_DARK

def card(slide, icon, title, desc, left, top, width=2.8, height=1.6,
         bg=C_LIGHT_PURPLE, title_color=C_PURPLE):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_textbox(slide, icon, left + 0.1, top + 0.1, 0.6, 0.5,
                font_size=24, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, left + 0.1, top + 0.55, width - 0.2, 0.4,
                font_size=13, bold=True, color=title_color)
    add_textbox(slide, desc, left + 0.1, top + 0.92, width - 0.2, 0.6,
                font_size=11, color=C_DARK)

# ══════════════════════════════════════════════════════════════
# SLIDE 1：封面
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
# 全頁深紫背景
add_rect(s, 0, 0, 13.33, 7.5, fill_color=C_DARK_PURPLE)
# 頂部裝飾條
add_rect(s, 0, 0, 13.33, 0.25, fill_color=C_AMBER)
# 底部裝飾條
add_rect(s, 0, 7.25, 13.33, 0.25, fill_color=C_AMBER)

add_textbox(s, '🤖 Claude Code', 1, 1.2, 11.33, 1,
            font_size=20, bold=False, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(s, '15 項被低估功能', 1, 2.0, 11.33, 1.2,
            font_size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, '新手完整操作指南', 1, 3.2, 11.33, 0.8,
            font_size=30, bold=False, color=C_AMBER, align=PP_ALIGN.CENTER)
add_textbox(s, '從安裝到實際應用，一步步帶你掌握每個功能', 1, 4.1, 11.33, 0.5,
            font_size=16, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(s, '來源：Boris Cherny（Claude Code 負責人）  ·  整理日期：2026-03-31',
            1, 6.5, 11.33, 0.4, font_size=13,
            color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2：什麼是 Claude Code？
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '什麼是 Claude Code？', '開始之前，先了解它的定位')

add_textbox(s, 'Claude Code 是 Anthropic 推出的 AI 編程工具，可以：',
            0.4, 1.45, 12.5, 0.45, font_size=16, bold=True, color=C_DARK)

features = [
    ('💬', '理解並執行自然語言指令（中文、英文都行）'),
    ('📝', '直接讀寫你電腦上的程式碼檔案'),
    ('🖥️', '在終端機執行指令、安裝套件、跑測試'),
    ('🌐', '操控瀏覽器進行 UI 測試與驗證'),
    ('🤖', '啟動多個 AI 代理人同時處理不同任務'),
]
for i, (icon, text) in enumerate(features):
    y = 1.95 + i * 0.62
    add_rect(s, 0.4, y, 12.4, 0.55, fill_color=C_WHITE)
    add_textbox(s, icon, 0.55, y + 0.08, 0.5, 0.4, font_size=18)
    add_textbox(s, text, 1.15, y + 0.1, 11.4, 0.38, font_size=14, color=C_DARK)

add_textbox(s, '📊  自 2026 年以來，Anthropic 工程師使用 Claude Code 後，程式碼產出成長超過 200%',
            0.4, 5.2, 12.4, 0.5, font_size=13, color=RGBColor(0x55, 0x55, 0x99))
add_textbox(s, '🔧  使用方式：CLI 終端機 / VS Code 擴充 / Claude Desktop App / Mobile App',
            0.4, 5.75, 12.4, 0.4, font_size=13, color=RGBColor(0x55, 0x55, 0x99))
add_textbox(s, '💡 新手提示：有 Claude.ai 帳號就能開始，免設定複雜環境！',
            0.4, 6.35, 12.4, 0.5, font_size=13, bold=True, color=C_AMBER)

# ══════════════════════════════════════════════════════════════
# SLIDE 3：安裝方式
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '如何開始使用 Claude Code？', '3 種安裝方式，選一個最適合你的')

methods = [
    ('方式 1', 'CLI 命令列（推薦新手）', C_GREEN,
     ['① 確認已安裝 Node.js 18+', '② 執行：npm install -g @anthropic-ai/claude-code', '③ 在專案資料夾執行：claude', '④ 輸入 Claude.ai 帳號登入即可']),
    ('方式 2', 'VS Code 擴充套件（最簡單）', C_BLUE,
     ['① 開啟 VS Code', '② 搜尋「Claude Code」擴充套件', '③ 點擊安裝', '④ 登入 Claude 帳號']),
    ('方式 3', 'Claude Desktop App（功能最完整）', C_PURPLE,
     ['① 前往 claude.ai 下載桌面版', '② 安裝並登入', '③ 開啟 Code 模式', '④ 支援內建瀏覽器、Cowork Dispatch']),
]

for i, (label, title, color, steps) in enumerate(methods):
    x = 0.3 + i * 4.35
    add_rect(s, x, 1.4, 4.15, 5.6, fill_color=C_WHITE)
    add_rect(s, x, 1.4, 4.15, 0.45, fill_color=color)
    add_textbox(s, label, x + 0.1, 1.4, 4.0, 0.45,
                font_size=14, bold=True, color=C_WHITE)
    add_textbox(s, title, x + 0.1, 1.9, 3.95, 0.55,
                font_size=13, bold=True, color=color)
    for j, step in enumerate(steps):
        add_textbox(s, step, x + 0.15, 2.5 + j * 0.72, 3.85, 0.65,
                    font_size=12, color=C_DARK)

add_textbox(s, '✅ 本指南以 CLI 方式為主要示範環境',
            0.4, 7.0, 12.4, 0.38, font_size=13,
            bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 4：功能總覽
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '15 項功能總覽', '5 大類別，逐一解鎖')

categories = [
    ('A', '跨裝置連續工作', '📱', ['① Mobile App', '② /teleport', '③ /remote-control', '④ Cowork Dispatch'], C_BLUE),
    ('B', '排程、並行與自動驗證', '⚙️', ['⑤ /loop + /schedule', '⑥ Hooks（4種）', '⑦ --worktree', '⑧ /batch'], C_GREEN),
    ('C', '開發工具鏈整合', '🔧', ['⑨ Chrome 擴充套件', '⑩ Desktop 瀏覽器', '⑪ /branch', '⑫ /btw'], C_PURPLE),
    ('D', 'SDK 與多 Repo 協作', '📦', ['⑬ --bare', '⑭ --add-dir', '⑮ --agent'], C_AMBER),
    ('E', '輸入體驗', '🎤', ['⑯ /voice 語音輸入'], RGBColor(0xEC, 0x48, 0x99)),
]

cols = [0.25, 2.65, 5.05, 7.45, 9.85]
for i, (letter, name, icon, items, color) in enumerate(categories):
    x = cols[i]
    add_rect(s, x, 1.35, 2.2, 5.8, fill_color=C_WHITE)
    add_rect(s, x, 1.35, 2.2, 0.55, fill_color=color)
    add_textbox(s, f'{icon} {letter}', x + 0.08, 1.36, 2.04, 0.52,
                font_size=20, bold=True, color=C_WHITE)
    add_textbox(s, name, x + 0.1, 1.95, 2.0, 0.55,
                font_size=12, bold=True, color=color)
    for j, item in enumerate(items):
        add_textbox(s, item, x + 0.12, 2.55 + j * 0.6, 1.96, 0.52,
                    font_size=11, color=C_DARK)

add_textbox(s, '💡 建議新手學習順序：①②③ → ⑦ → ⑥ → ⑤⑧ → 其餘進階功能',
            0.25, 7.05, 12.8, 0.38, font_size=13, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 5：功能① Mobile App
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ① — Mobile App（行動裝置存取）', '類別 A：跨裝置連續工作', bg_color=C_BLUE)

add_textbox(s, '📱 在手機上繼續你的 AI 開發工作，不需要開電腦',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

step_box(s, [
    '① 在手機安裝「Claude」App（iOS App Store / Android Google Play）',
    '② 登入與電腦相同的 Claude 帳號',
    '③ 點擊左側選單 → 切換至「Code」分頁',
    '④ 即可看到並繼續所有雲端工作階段',
    '⑤ 直接用手機鍵盤或語音輸入指令',
], 0.4, 1.85, 7.5, 3.2, bg=C_LIGHT_GREEN)

example_box(s, '實際應用場景',
    ['情境：你在咖啡廳，電腦沒帶，但突然想到一個 Bug',
     '',
     '→ 打開手機 Claude App → Code 分頁',
     '→ 找到昨天的工作階段',
     '→ 輸入「幫我修復 b3_savings_plan.js 第 245 行的 null 錯誤」',
     '→ Claude 修改後，回家用電腦確認即可',
     '',
     'Boris Cherny（Claude Code 負責人）表示：他大比例的程式修改直接從 iPhone 完成'],
    8.1, 1.85, 4.9, 3.2, bg=C_LIGHT_AMBER)

info_box(s, '新手提示',
    ['• 雲端工作階段才能跨裝置同步，本機 CLI 模式不適用',
     '• 手機版適合「查看進度」和「下簡單指令」',
     '• 複雜操作（如大量重構）建議回到電腦上執行'],
    0.4, 5.15, 12.2, 1.35, bg=C_LIGHT_BLUE, icon='💡')

add_textbox(s, '⭐ 難度：★☆☆☆☆（最簡單，零設定）',
            0.4, 6.6, 5, 0.45, font_size=14, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 6：功能② /teleport
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ② — /teleport（雲端⇄本機傳送）', '類別 A：跨裝置連續工作', bg_color=C_BLUE)

add_textbox(s, '☁️ 將雲端工作階段「傳送」回本機繼續執行，或反向上傳至雲端',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

add_textbox(s, '使用場景：手機開工 → 回到電腦上繼續（本機有完整檔案系統和工具）',
            0.4, 1.83, 12.4, 0.38, font_size=13, color=RGBColor(0x44, 0x44, 0x88))

step_box(s, [
    '方法 A：啟動時使用旗標',
    '  在終端機輸入：claude --teleport',
    '  → 自動列出雲端工作階段，選擇要拉回的那個',
    '',
    '方法 B：工作階段中途使用',
    '  在 Claude 對話中輸入：/teleport',
    '  → 當前工作階段立即切換至本機環境',
], 0.4, 2.28, 6.8, 3.3, bg=C_LIGHT_GREEN)

# Code box
add_textbox(s, '指令範例', 7.5, 2.28, 5.5, 0.38, font_size=13, bold=True, color=C_GREEN)
code_box(s, '''# 方法 A：在終端機啟動時
claude --teleport

# 方法 B：工作階段中直接輸入
/teleport

# 選擇要傳送的工作階段後
# Claude 自動在本機環境繼續執行''', 7.5, 2.66, 5.5, 2.95)

info_box(s, '新手提示',
    ['• 需要有 Claude.ai 帳號的「雲端工作階段」才能使用',
     '• 如果只用本機 CLI，沒有雲端階段可傳送',
     '• 傳送後本機有完整檔案讀寫權限，功能比手機強'],
    0.4, 5.65, 12.2, 1.35, bg=C_LIGHT_BLUE, icon='💡')

add_textbox(s, '⭐ 難度：★★☆☆☆', 0.4, 7.08, 5, 0.38, font_size=14, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 7：功能③ /remote-control
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ③ — /remote-control（遠端操控本機）', '類別 A：跨裝置連續工作', bg_color=C_BLUE)

add_textbox(s, '🖥️ 從手機或瀏覽器，即時操控你家電腦上正在執行的 Claude 工作階段',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

step_box(s, [
    '設定步驟（一次性）：',
    '① 電腦上開啟 Claude CLI → 輸入 /config',
    '② 找到「Enable Remote Control for all sessions」',
    '③ 設為「開啟（Enabled）」→ 之後每次自動生效',
    '',
    '每次使用時：',
    '④ 電腦上執行 claude（確保有在跑）',
    '⑤ 手機或其他裝置 → Claude App → 輸入 /remote-control',
    '⑥ 輸入電腦端顯示的連線碼，即可遠端操控',
], 0.4, 1.85, 7.4, 4.0)

example_box(s, '實際應用場景',
    ['情境：讓 Claude 在家裡電腦跑大型任務，你出門買東西',
     '',
     '電腦端（出門前）：',
     '  claude --remote-control --name "重構任務"',
     '',
     '手機端（路上查看）：',
     '  /remote-control → 輸入連線碼',
     '  → 看到 Claude 進度 → 插入新指令',
     '',
     '本專案的 cladue.bat 已整合此功能！'],
    7.9, 1.85, 5.1, 4.0, bg=C_LIGHT_AMBER)

info_box(s, '新手提示',
    ['• 電腦要保持開機且 Claude 在執行，手機才能連上',
     '• 非常適合「讓 AI 跑長時間任務」同時你去做其他事'],
    0.4, 5.95, 12.2, 1.1, bg=C_LIGHT_BLUE, icon='💡')

add_textbox(s, '⭐ 難度：★★☆☆☆', 0.4, 7.1, 5, 0.38, font_size=14, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 8：功能④ Cowork Dispatch
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ④ — Cowork Dispatch（MCP 遠端控制介面）', '類別 A：跨裝置連續工作', bg_color=C_BLUE)

add_textbox(s, '🤝 透過 Claude Desktop App，讓 AI 幫你操控瀏覽器、管理 Slack、處理電子郵件',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

step_box(s, [
    '① 安裝 Claude Desktop App（桌面版，非 CLI）',
    '② 在 Claude Desktop 的設定中，新增 MCP 伺服器',
    '   • Slack MCP：讓 Claude 讀寫 Slack 訊息',
    '   • Browser MCP：讓 Claude 操控 Chrome',
    '   • File MCP：讓 Claude 管理本機檔案',
    '③ 在對話中說：「幫我回覆今天所有 Slack 的問題」',
    '④ Claude 自動操作 Slack → 完成任務',
], 0.4, 1.85, 7.4, 3.8)

info_box(s, '⚠️ 注意事項（新手必讀）',
    ['• 需要 Claude Desktop App，CLI 版不支援',
     '• MCP 伺服器設定有一定技術門檻',
     '• 建議先熟悉其他功能，進階後再學此功能',
     '• Slack MCP 需要 Slack 管理員權限設定 Token'],
    7.9, 1.85, 5.1, 3.8, bg=RGBColor(0xFF, 0xEE, 0xEE),
    title_color=C_RED, icon='⚠️')

example_box(s, 'Boris Cherny 的使用方式',
    ['在非寫程式時間，Cherny 用 Cowork Dispatch：',
     '• 每天早上自動整理 Slack 未讀訊息',
     '• 起草並發送電子郵件回覆',
     '• 組織和命名桌面檔案',
     '→ 讓 Claude 幫處理非程式工作，',
     '  自己專注在更有創意的任務'],
    0.4, 5.75, 12.2, 1.5, bg=C_LIGHT_AMBER)

add_textbox(s, '⭐ 難度：★★★★☆（進階功能，建議有基礎後再嘗試）',
            0.4, 7.3, 9, 0.38, font_size=13, bold=True, color=C_RED)

# ══════════════════════════════════════════════════════════════
# SLIDE 9：功能⑤ /loop + /schedule
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑤ — /loop + /schedule（排程自動執行）', '類別 B：排程、並行與自動驗證', bg_color=C_GREEN)

add_textbox(s, '⏰ 讓 Claude 以固定時間間隔自動重複執行任務，最長可排程一週',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

add_textbox(s, '基本語法', 0.4, 1.85, 12, 0.38, font_size=14, bold=True, color=C_GREEN)
code_box(s, '''/loop [時間間隔] [要執行的指令或任務]

時間單位：
  s = 秒    m = 分鐘    h = 小時    d = 天

範例：
  /loop 5m /check-pr        # 每 5 分鐘檢查 PR 狀態
  /loop 30m /slack-feedback # 每 30 分鐘整理 Slack 回饋
  /loop 1h /run-tests       # 每 1 小時跑一次測試
  /loop 1d /generate-report # 每天生成報告''', 0.4, 2.25, 6.5, 3.2)

example_box(s, 'Boris Cherny 的實際用法',
    ['用法一：自動處理 Code Review',
     '  /loop 5m /babysit',
     '  → 每 5 分鐘自動：審查新 PR、處理 rebase 衝突',
     '',
     '用法二：自動整理 Slack 回饋',
     '  /loop 30m /slack-feedback',
     '  → 每 30 分鐘：讀取 Slack → 整理成 PR → 自動提交',
     '',
     '適合你的情境（Money Tutor 專案）：',
     '  /loop 1h /scan-console-errors',
     '  → 每小時掃描 JS 錯誤'],
    7.1, 1.85, 5.9, 3.6, bg=C_LIGHT_AMBER)

info_box(s, '新手提示',
    ['• 先測試指令能正確執行，再加 /loop 自動化',
     '• /schedule 可設定一次性的未來時間執行',
     '• 小心無限迴圈！加上 /stop 可以隨時停止'],
    0.4, 5.55, 12.2, 1.45, bg=C_LIGHT_BLUE, icon='💡')

add_textbox(s, '⭐ 難度：★★★☆☆', 0.4, 7.08, 5, 0.38, font_size=14, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 10：功能⑥ Hooks
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑥ — Hooks（週期節點注入邏輯）', '類別 B：排程、並行與自動驗證', bg_color=C_GREEN)

add_textbox(s, '🪝 在 Claude 工作的特定時機點，自動執行你預設好的程式或指令',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

hook_data = [
    ('SessionStart', '工作階段開始時', '自動載入 CLAUDE.md 說明\n→「每次開始工作先讀取專案說明」', C_BLUE),
    ('PreToolUse', '每次執行工具前', '自動寫日誌記錄\n→「每次 Claude 改檔案前先備份」', C_GREEN),
    ('PermissionRequest', '需要授權時', '將授權請求發到手機\n→「手機批准高風險操作」', C_AMBER),
    ('Stop', 'Claude 暫停等待時', '自動繼續推進任務\n→「無人值守時 Claude 自動繼續」', C_PURPLE),
]

for i, (name, timing, usage, color) in enumerate(hook_data):
    x = 0.3 + i * 3.17
    add_rect(s, x, 1.85, 3.1, 4.0, fill_color=C_WHITE)
    add_rect(s, x, 1.85, 3.1, 0.45, fill_color=color)
    add_textbox(s, name, x + 0.1, 1.87, 2.9, 0.42,
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(s, f'⏱️ 時機：{timing}', x + 0.1, 2.35, 2.9, 0.45,
                font_size=12, bold=True, color=color)
    add_textbox(s, f'💡 用途：\n{usage}', x + 0.1, 2.85, 2.9, 1.6,
                font_size=11, color=C_DARK)

# How to set up
add_textbox(s, '設定方式', 0.4, 5.95, 3, 0.38, font_size=13, bold=True, color=C_GREEN)
code_box(s, '''# 在 settings.json 中設定
{
  "hooks": {
    "SessionStart": [
      "cat CLAUDE.md"
    ],
    "PreToolUse": [
      "echo 'Tool: $TOOL' >> log.txt"
    ]
  }
}''', 0.4, 6.35, 5.8, 1.1, font_size=11)

info_box(s, '新手提示',
    ['• Hook 設定檔位於 ~/.claude/settings.json',
     '• 從最簡單的 SessionStart 開始學（載入說明文件）',
     '• PreToolUse 可在執行前驗證，防止意外操作'],
    6.4, 5.95, 6.5, 1.5, bg=C_LIGHT_BLUE, icon='💡')

add_textbox(s, '⭐ 難度：★★★☆☆', 0.4, 7.5, 5, 0.35, font_size=13, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 11：功能⑦ --worktree
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑦ — --worktree（隔離工作樹）', '類別 B：排程、並行與自動驗證', bg_color=C_GREEN)

add_textbox(s, '🌲 在同一個專案中，開啟多個完全獨立的工作環境，讓多個 Claude 同時互不干擾地工作',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

# Diagram
add_rect(s, 0.4, 1.88, 12.4, 0.5, fill_color=C_DARK_PURPLE)
add_textbox(s, '主要專案 (main branch)', 0.5, 1.9, 12.2, 0.45,
            font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

worktrees = [
    ('Worktree A', 'Claude #1 修復 Bug', '功能 A\nfeature-bug-fix', C_GREEN),
    ('Worktree B', 'Claude #2 新增功能', '功能 B\nfeature-new-ui', C_BLUE),
    ('Worktree C', 'Claude #3 寫測試', '功能 C\nfeature-tests', C_PURPLE),
]
for i, (wt, task, branch, color) in enumerate(worktrees):
    x = 0.4 + i * 4.1
    add_rect(s, x, 2.48, 3.9, 1.6, fill_color=color)
    add_textbox(s, wt, x + 0.1, 2.5, 3.7, 0.42,
                font_size=14, bold=True, color=C_WHITE)
    add_textbox(s, task, x + 0.1, 2.95, 3.7, 0.42,
                font_size=12, color=C_WHITE)
    add_textbox(s, branch, x + 0.1, 3.38, 3.7, 0.6,
                font_size=11, color=RGBColor(0xCC, 0xFF, 0xCC))

step_box(s, [
    '方法 A：CLI 啟動時',
    '  claude --worktree  → Claude 自動建立隔離環境',
    '',
    '方法 B：Agent 工具中（本專案已在用！）',
    '  isolation: "worktree"  → 子代理人在獨立分支作業',
    '',
    '方法 C：手動建立 Git Worktree',
    '  git worktree add ../feature-branch feature/my-task',
    '  cd ../feature-branch && claude',
], 0.4, 4.15, 7.0, 2.95)

example_box(s, '本專案的實際應用',
    ['Money Tutor 已在使用！',
     '',
     '當 Claude 需要同時：',
     '• 修復 B3 的 Bug',
     '• 重構 B6 的 CSS',
     '• 撰寫測試腳本',
     '',
     '→ 各自在不同 worktree 執行',
     '→ 不會互相污染程式碼'],
    7.5, 4.15, 5.1, 2.95, bg=C_LIGHT_AMBER)

add_textbox(s, '⭐ 難度：★★★☆☆  |  Boris Cherny 說：隨時都有幾十個 Claude 在跑，worktrees 是核心基礎設施',
            0.4, 7.15, 12.4, 0.38, font_size=12, color=C_PURPLE)

# ══════════════════════════════════════════════════════════════
# SLIDE 12：功能⑧ /batch
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑧ — /batch（大規模並行代理人）', '類別 B：排程、並行與自動驗證', bg_color=C_GREEN)

add_textbox(s, '⚡ 一個指令，Claude 自動拆解工作，同時開啟數十至數千個代理人並行完成',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

add_textbox(s, '使用方式', 0.4, 1.85, 5.5, 0.38, font_size=14, bold=True, color=C_GREEN)
code_box(s, '''/batch [你要做的大型任務]

實際範例：
  /batch migrate from jest to vite
  /batch add TypeScript types to all JS files
  /batch write unit tests for all components
  /batch translate all comments to English

Claude 會：
  1. 先問幾個問題確認需求
  2. 自動拆解成數十個子任務
  3. 開啟多個 worktree 並行執行
  4. 完成後合併結果''', 0.4, 2.28, 6.3, 3.2)

example_box(s, '實際應用場景',
    ['場景：Money Tutor 要把所有 B 系列',
     '單元從 JavaScript 改為 TypeScript',
     '',
     '舊做法：一個個手動改（幾天）',
     '',
     '新做法：',
     '  /batch migrate B series JS to TypeScript',
     '',
     '→ Claude 自動拆解：',
     '  • 代理人1: b1_daily_budget.js',
     '  • 代理人2: b2_allowance_diary.js',
     '  • ... 同時進行',
     '→ 幾分鐘內完成'],
    6.9, 1.85, 6.1, 3.65, bg=C_LIGHT_AMBER)

info_box(s, '⚠️ 新手注意事項',
    ['• 2025年2月正式推出，相對較新的功能',
     '• 「數千個 worktree」並行 = 大量 API 費用，先從小任務試',
     '• 先確認你的任務真的適合並行（有相依性的不適合）',
     '• 建議先用 /batch 搭配小型任務，熟悉後再擴大規模'],
    0.4, 5.55, 12.2, 1.5, bg=RGBColor(0xFF, 0xEE, 0xEE),
    title_color=C_RED, icon='⚠️')

add_textbox(s, '⭐ 難度：★★★★☆', 0.4, 7.1, 5, 0.38, font_size=14, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 13：功能⑨ Chrome 擴充套件
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑨ — Chrome 擴充套件（瀏覽器操控）', '類別 C：開發工具鏈整合', bg_color=C_PURPLE)

add_textbox(s, '🌐 讓 Claude 直接操控 Chrome，讀取 Console 錯誤，測試 UI 互動行為',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

step_box(s, [
    '① 在 Chrome 瀏覽器，前往 Chrome Web Store',
    '② 搜尋「Claude Code」擴充套件（注意：目前為 beta 版）',
    '③ 點擊「加到 Chrome」→ 安裝完成',
    '④ 點擊擴充套件圖示 → 連結到你的 Claude 工作階段',
    '⑤ 在 Claude 對話中說：「測試這個按鈕是否能正確點擊」',
    '⑥ Claude 自動操控瀏覽器完成測試',
], 0.4, 1.85, 7.2, 3.3)

example_box(s, '實際應用場景',
    ['情境：Money Tutor B1 頁面有個按鈕點不到',
     '',
     '舊做法：手動打開瀏覽器 → 開 DevTools → 看 Console',
     '',
     '新做法：',
     '  在 Claude 中說：',
     '  「打開 b1_daily_budget.html，',
     '    點擊確認按鈕，告訴我 Console 有沒有錯誤」',
     '',
     '→ Claude 自動開啟頁面 → 點擊 → 回報錯誤',
     '→ 同時提供修復建議'],
    7.8, 1.85, 5.2, 3.3, bg=C_LIGHT_AMBER)

add_textbox(s, 'Boris Cherny 的類比：「不讓工程師開瀏覽器，成果會打折；Chrome 擴充套件把驗證能力還給 Claude」',
            0.4, 5.25, 12.2, 0.5, font_size=13, color=RGBColor(0x44, 0x44, 0x88))

info_box(s, '新手提示',
    ['• 目前（2026-03）仍為 beta 版，功能可能持續更新',
     '• 非常適合前端開發者！可以邊改 CSS 邊讓 Claude 截圖驗證效果',
     '• 若用 Claude Desktop App，改用內建瀏覽器（功能10）效果更完整'],
    0.4, 5.8, 12.2, 1.3, bg=C_LIGHT_BLUE, icon='💡')

add_textbox(s, '⭐ 難度：★★☆☆☆', 0.4, 7.15, 5, 0.38, font_size=14, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 14：功能⑩⑪⑫ 合併
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑩⑪⑫ — Desktop 瀏覽器 / /branch / /btw', '類別 C：開發工具鏈整合', bg_color=C_PURPLE)

sections = [
    ('⑩ Desktop 內建瀏覽器', C_PURPLE,
     '讓 Claude Desktop App 自動啟動 web server 並以內建瀏覽器測試',
     ['Claude Desktop App 中自動生效', 'Claude 自動啟動本地伺服器', '完整閉環：改程式 → 測試 → 修復，不需切換視窗'],
     ['CLI 用戶→改用功能⑨ Chrome 擴充套件', '僅 Desktop 版支援（非 CLI）', '非常適合全端開發者']),
    ('⑪ /branch（工作節點分叉）', C_BLUE,
     '從同一工作點分叉出不同嘗試方向，隨時可切回原本主線',
     ['在 Claude 對話中輸入：/branch', '系統詢問新方向的名稱', '切換：/branch switch [名稱]', '回主線：/branch main'],
     ['不怕改壞！隨時可以切回去', '適合「不確定哪種方法比較好」的情況', '例：同時試兩種 CSS 解法再選一個']),
    ('⑫ /btw（不中斷插入問答）', C_GREEN,
     '代理人仍在執行長任務時，插入快速旁支問答而不中斷',
     ['Claude 正在跑大型任務時', '直接輸入：/btw [你的問題]', 'Claude 快速回答後繼續執行主任務', '不影響主任務進度'],
     ['讓 AI 一邊工作一邊回答你的問題', '例：Claude 重構程式時，你問「這個函數的意義是什麼？」', '不打斷主任務流程']),
]

for i, (title, color, desc, steps, tips) in enumerate(sections):
    x = 0.3 + i * 4.3
    add_rect(s, x, 1.38, 4.1, 5.8, fill_color=C_WHITE)
    add_rect(s, x, 1.38, 4.1, 0.48, fill_color=color)
    add_textbox(s, title, x + 0.1, 1.4, 3.9, 0.44,
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(s, desc, x + 0.1, 1.92, 3.9, 0.6,
                font_size=11, color=C_DARK)
    add_textbox(s, '操作步驟', x + 0.1, 2.58, 3.9, 0.32,
                font_size=11, bold=True, color=color)
    for j, step in enumerate(steps):
        add_textbox(s, f'• {step}', x + 0.1, 2.93 + j * 0.48,
                    3.9, 0.46, font_size=10, color=C_DARK)
    add_textbox(s, '💡 新手提示', x + 0.1, 4.7, 3.9, 0.3,
                font_size=11, bold=True, color=color)
    for j, tip in enumerate(tips):
        add_textbox(s, f'✓ {tip}', x + 0.1, 5.05 + j * 0.38,
                    3.9, 0.36, font_size=9, color=RGBColor(0x44, 0x44, 0x88))

add_textbox(s, '⭐ 難度：⑩★★★☆☆  ⑪★★☆☆☆  ⑫★☆☆☆☆',
            0.4, 7.15, 10, 0.38, font_size=13, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 15：功能⑬ --bare
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑬ — --bare（SDK 快速啟動旗標）', '類別 D：SDK 與多 Repo 協作優化', bg_color=C_AMBER)

add_textbox(s, '🚀 非互動式呼叫時，加上 --bare 可跳過 CLAUDE.md 等載入程序，啟動速度最多快 10 倍',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

add_textbox(s, '什麼情況需要用 --bare？', 0.4, 1.88, 6, 0.38, font_size=14, bold=True, color=C_AMBER)
scenarios = [
    ('通常模式（互動）', '每次啟動讀取 CLAUDE.md + 設定 + MCP', '= 較慢，但有完整上下文'),
    ('--bare 模式（腳本/CI）', '跳過 CLAUDE.md、設定檔、MCP 伺服器', '= 啟動快 10 倍，適合自動化'),
]
for i, (title, desc1, desc2) in enumerate(scenarios):
    add_rect(s, 0.4 + i * 6.1, 2.32, 5.8, 1.5, fill_color=C_WHITE)
    add_textbox(s, title, 0.55 + i * 6.1, 2.38, 5.5, 0.4,
                font_size=13, bold=True, color=C_AMBER)
    add_textbox(s, desc1, 0.55 + i * 6.1, 2.8, 5.5, 0.4,
                font_size=11, color=C_DARK)
    add_textbox(s, desc2, 0.55 + i * 6.1, 3.22, 5.5, 0.45,
                font_size=12, bold=True, color=C_GREEN if i == 1 else C_DARK)

add_textbox(s, '使用方式', 0.4, 3.95, 6, 0.38, font_size=14, bold=True, color=C_AMBER)
code_box(s, '''# 基本 --bare 用法
claude --bare -p "幫我分析這個函數有沒有 bug" < myfile.js

# 在 CI/CD 腳本中
echo "執行自動測試分析" | claude --bare -p "分析測試結果"

# Python SDK 中
import anthropic
claude = anthropic.Claude(bare=True)  # 快速啟動

# TypeScript SDK 中
const claude = new Claude({ bare: true })''', 0.4, 4.38, 7.5, 2.65)

example_box(s, '實際應用場景',
    ['情境：自動化 CI/CD 流程',
     '',
     '每次 git push 時自動分析：',
     '  #!/bin/bash',
     '  git diff HEAD~1 | claude --bare -p',
     '  "這次改動有沒有潛在的 Bug？"',
     '',
     '效果：',
     '  通常：等待3秒載入環境',
     '  --bare：0.3秒，快10倍！'],
    8.1, 3.95, 5.0, 3.08, bg=C_LIGHT_AMBER)

info_box(s, '新手提示',
    ['• --bare 適合腳本化、自動化場景，不適合需要上下文的互動開發',
     '• 啟動快的原因：不讀 CLAUDE.md（可能很大），不連 MCP 伺服器'],
    0.4, 7.08, 12.2, 0.42, bg=C_LIGHT_BLUE, icon='💡')

add_textbox(s, '⭐ 難度：★★★☆☆（需要了解 CLI/SDK 呼叫方式）',
            0.4, 7.4, 8, 0.38, font_size=13, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 16：功能⑭ --add-dir
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑭ — --add-dir（跨 Repo 操作）', '類別 D：SDK 與多 Repo 協作優化', bg_color=C_AMBER)

add_textbox(s, '📂 同時讓 Claude 操作多個專案資料夾，不只是「看到」，還能直接讀寫',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

step_box(s, [
    '方法 A：啟動時加入',
    '  claude --add-dir /path/to/other-project',
    '  claude --add-dir ../shared-components --add-dir ../api-server',
    '',
    '方法 B：工作階段中加入',
    '  /add-dir ../another-project',
    '',
    '方法 C：永久設定（團隊共用）',
    '  在 settings.json 中：',
    '  { "additionalDirectories": ["../shared-lib", "../backend"] }',
], 0.4, 1.85, 7.0, 3.5)

example_box(s, '實際應用場景',
    ['情境：Money Tutor 前端 + 後端 API',
     '',
     '目前結構：',
     '  money_tutor/（前端）',
     '  money_api/（後端，另一個資料夾）',
     '',
     '舊做法：需要開兩個 Claude 視窗分別處理',
     '',
     '新做法：',
     '  claude --add-dir ../money_api',
     '',
     '→ 一個 Claude 同時看兩個資料夾',
     '→ 「把後端的 API 路由同步到前端 JS」',
     '→ 一次完成！'],
    7.6, 1.85, 5.4, 3.5, bg=C_LIGHT_AMBER)

info_box(s, '新手提示',
    ['• --add-dir 賦予「完整讀寫權限」，不只是查看',
     '• 適合微服務架構、前後端分離的開發情境',
     '• 設定在 settings.json 可讓整個團隊共用相同的多目錄設定'],
    0.4, 5.45, 12.2, 1.35, bg=C_LIGHT_BLUE, icon='💡')

add_textbox(s, '⭐ 難度：★★☆☆☆', 0.4, 6.9, 5, 0.38, font_size=14, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 17：功能⑮ --agent
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑮ — --agent（自訂子代理人）', '類別 D：SDK 與多 Repo 協作優化', bg_color=C_AMBER)

add_textbox(s, '🤖 建立具有自訂系統提示與工具限制的專屬代理人，適合特定任務場景',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

step_box(s, [
    '① 在 .claude/agents/ 資料夾建立 Markdown 設定檔',
    '   例：.claude/agents/readonly.md',
    '',
    '② 在 Markdown 中定義代理人行為',
    '   （系統提示、允許的工具、限制條件）',
    '',
    '③ 啟動代理人：',
    '   claude --agent=readonly',
    '   claude --agent=tester',
    '   claude --agent=doc-writer',
], 0.4, 1.88, 6.2, 3.5)

add_textbox(s, '設定檔範例（.claude/agents/readonly.md）', 6.8, 1.88, 6.2, 0.38,
            font_size=13, bold=True, color=C_AMBER)
code_box(s, '''---
name: readonly
description: 只讀審查代理人（安全審查用）
tools: [read_file, search, list_files]
# 不允許：write_file, bash, delete
---

# 系統提示
你是一個程式碼審查代理人。
你只能讀取檔案，不能修改任何內容。
請仔細分析程式碼，找出潛在問題並提供建議。

分析重點：
- 安全性漏洞
- 效能問題
- 程式碼品質''', 6.8, 2.3, 6.2, 3.05, font_size=11)

example_box(s, '實用代理人配置',
    ['📖 readonly 代理人',
     '   只能讀檔 → 安全程式碼審查',
     '',
     '🧪 tester 代理人',
     '   只能跑測試 → 自動化 QA',
     '',
     '📝 doc-writer 代理人',
     '   只能寫文件 → 自動產生說明',
     '',
     '🔍 security-audit 代理人',
     '   讀取 + 分析 → 安全掃描',
     '',
     '本專案 CLAUDE.md 的 subagent_type',
     '就是同樣的概念延伸！'],
    0.4, 5.45, 6.2, 1.6, bg=C_LIGHT_AMBER)

info_box(s, '新手提示',
    ['• 從最簡單的 readonly 代理人開始學',
     '• 工具限制可以防止代理人做出意外操作',
     '• 可以建立多個不同用途的代理人備用'],
    6.8, 5.45, 6.2, 1.6, bg=C_LIGHT_BLUE, icon='💡')

add_textbox(s, '⭐ 難度：★★★☆☆', 0.4, 7.1, 5, 0.38, font_size=14, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 18：功能⑯ /voice
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 ⑯ — /voice（語音輸入）', '類別 E：輸入體驗', bg_color=RGBColor(0xEC, 0x48, 0x99))

add_textbox(s, '🎤 用說話代替打字，向 Claude 下指令！Boris Cherny 日常以說話為主、鍵盤為輔',
            0.4, 1.35, 12.4, 0.45, font_size=16, bold=True, color=C_DARK)

platforms = [
    ('CLI 終端機', '💻', ['① 在 Claude 對話中輸入：/voice', '② 按住「空白鍵」開始說話', '③ 放開空白鍵 → 自動轉文字', '④ Claude 執行你說的指令'], C_BLUE),
    ('Claude Desktop', '🖥️', ['① 開啟 Claude Desktop App', '② 點擊介面上的「語音按鈕」（麥克風圖示）', '③ 說出你的指令', '④ 放開後自動識別並執行'], C_PURPLE),
    ('iPhone / iPad', '📱', ['① 系統設定 → 一般 → 鍵盤 → 開啟「聽寫」', '② 在 Claude App 的輸入框', '③ 點擊鍵盤上的麥克風圖示', '④ 說出指令 → 自動輸入'], RGBColor(0xEC, 0x48, 0x99)),
]

for i, (platform, icon, steps, color) in enumerate(platforms):
    x = 0.3 + i * 4.2
    add_rect(s, x, 1.88, 4.0, 4.3, fill_color=C_WHITE)
    add_rect(s, x, 1.88, 4.0, 0.48, fill_color=color)
    add_textbox(s, f'{icon} {platform}', x + 0.1, 1.9, 3.8, 0.44,
                font_size=14, bold=True, color=C_WHITE)
    for j, step in enumerate(steps):
        add_textbox(s, step, x + 0.12, 2.43 + j * 0.78, 3.76, 0.72,
                    font_size=11, color=C_DARK)

example_box(s, '語音輸入的實際效益',
    ['你可以說：',
     '',
     '「幫我在 b3_savings_plan.js 的第 240 行',
     '  加一個 console.log，顯示 state 的值」',
     '',
     '「找出所有 B 系列 JS 中，',
     '  有 setInterval 的地方，列出來」',
     '',
     '「把剛才改的內容 commit 到 git，',
     '  訊息寫：修復 B3 月曆 UX 問題」',
     '',
     '→ 解放雙手，在站著、走路時也能工作！'],
    0.3, 6.28, 12.6, 1.18, bg=C_LIGHT_AMBER)

add_textbox(s, '⭐ 難度：★☆☆☆☆（iPhone 聽寫功能最簡單，零設定）',
            0.4, 7.48, 9, 0.38, font_size=13, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 19：新手學習路徑
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '新手學習路徑建議', '按照這個順序，由簡入難，穩步掌握所有功能')

add_textbox(s, '🔰 第一階段：入門（第1週）', 0.4, 1.38, 12, 0.42,
            font_size=16, bold=True, color=C_GREEN)

stage1 = [
    ('①', 'Mobile App', '下載 App 熟悉界面', '0️⃣ 零設定'),
    ('⑯', '/voice', '學會語音輸入，省打字', '⭐ 1分鐘'),
    ('⑫', '/btw', '工作中插入問答', '⭐ 立即可用'),
    ('⑪', '/branch', '分叉嘗試不怕改壞', '⭐⭐ 簡單'),
]
for i, (num, name, desc, diff) in enumerate(stage1):
    x = 0.35 + i * 3.15
    add_rect(s, x, 1.85, 3.0, 1.35, fill_color=C_LIGHT_GREEN)
    add_textbox(s, num, x + 0.08, 1.88, 0.5, 0.5, font_size=18, bold=True, color=C_GREEN)
    add_textbox(s, name, x + 0.6, 1.9, 2.35, 0.45, font_size=13, bold=True, color=C_DARK)
    add_textbox(s, desc, x + 0.1, 2.35, 2.85, 0.38, font_size=11, color=C_DARK)
    add_textbox(s, diff, x + 0.1, 2.73, 2.85, 0.38, font_size=11, color=C_GREEN)

add_textbox(s, '🔷 第二階段：進階（第2~3週）', 0.4, 3.28, 12, 0.42,
            font_size=16, bold=True, color=C_BLUE)
stage2 = [
    ('②③', '/teleport\n/remote-control', '跨裝置無縫工作', '⭐⭐ 需要雲端帳號'),
    ('⑦', '--worktree', '隔離環境並行開發', '⭐⭐⭐ 需懂 git'),
    ('⑭', '--add-dir', '跨資料夾操作', '⭐⭐ 設定簡單'),
    ('⑨', 'Chrome 擴充', '瀏覽器自動測試', '⭐⭐ Beta版'),
]
for i, (num, name, desc, diff) in enumerate(stage2):
    x = 0.35 + i * 3.15
    add_rect(s, x, 3.76, 3.0, 1.45, fill_color=C_LIGHT_BLUE)
    add_textbox(s, num, x + 0.08, 3.79, 0.6, 0.5, font_size=16, bold=True, color=C_BLUE)
    add_textbox(s, name, x + 0.68, 3.81, 2.27, 0.55, font_size=12, bold=True, color=C_DARK)
    add_textbox(s, desc, x + 0.1, 4.4, 2.85, 0.38, font_size=11, color=C_DARK)
    add_textbox(s, diff, x + 0.1, 4.78, 2.85, 0.38, font_size=11, color=C_BLUE)

add_textbox(s, '🔴 第三階段：高手（第4週起）', 0.4, 5.32, 12, 0.42,
            font_size=16, bold=True, color=C_RED)
stage3 = [
    ('⑥', 'Hooks', '自動化注入邏輯'),
    ('⑤', '/loop', '排程自動執行'),
    ('⑮', '--agent', '自訂代理人'),
    ('⑧', '/batch', '大規模並行任務'),
    ('④⑬', 'Cowork\n--bare', '進階整合'),
]
for i, (num, name, desc) in enumerate(stage3):
    x = 0.35 + i * 2.52
    add_rect(s, x, 5.8, 2.4, 1.35, fill_color=C_LIGHT_PURPLE)
    add_textbox(s, num, x + 0.08, 5.83, 0.5, 0.42, font_size=15, bold=True, color=C_PURPLE)
    add_textbox(s, name, x + 0.58, 5.85, 1.78, 0.5, font_size=12, bold=True, color=C_DARK)
    add_textbox(s, desc, x + 0.1, 6.35, 2.25, 0.38, font_size=11, color=C_DARK)

add_textbox(s, '💡 每掌握一個功能，就能節省大量重複性工作時間！',
            0.4, 7.2, 12.4, 0.38, font_size=14, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 20：本專案實際應用
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 'Money Tutor 專案實際應用案例', '這些功能我們已在使用，你也可以！')

add_textbox(s, '✅ 已在 Money Tutor 專案中使用的功能：',
            0.4, 1.38, 12, 0.45, font_size=16, bold=True, color=C_DARK)

used = [
    ('⑦ --worktree', 'Agent isolation', 'Agent 工具設定 isolation: "worktree"\n→ 每次子代理人任務自動建立隔離分支，互不干擾', C_GREEN),
    ('③ /remote-control', 'cladue.bat 整合', '雙視窗模式：一個跑 Claude，一個開啟遠端控制伺服器\n→ 可從另一台裝置連入監控進度', C_BLUE),
    ('⑮ --agent 概念', 'subagent_type', 'CLAUDE.md 定義多種代理人：Explore、Plan、general-purpose\n→ 不同任務自動選用最適合的 AI 角色', C_PURPLE),
    ('⑥ Hooks 概念', 'CLAUDE.md SessionStart', 'CLAUDE.md 在每次對話開始自動載入\n→ 等同 SessionStart Hook 的設計意圖', C_AMBER),
]

for i, (feat, impl, desc, color) in enumerate(used):
    y = 1.88 + i * 1.25
    add_rect(s, 0.4, y, 12.2, 1.18, fill_color=C_WHITE)
    add_rect(s, 0.4, y, 0.25, 1.18, fill_color=color)
    add_textbox(s, feat, 0.75, y + 0.08, 2.8, 0.45, font_size=13, bold=True, color=color)
    add_textbox(s, impl, 3.65, y + 0.08, 3.5, 0.45, font_size=13, bold=True, color=C_DARK)
    add_textbox(s, desc, 0.75, y + 0.56, 11.7, 0.55, font_size=11, color=C_DARK)

add_textbox(s, '🔮 下一步：可考慮加入的功能', 0.4, 7.0, 12, 0.38,
            font_size=14, bold=True, color=C_PURPLE)
add_textbox(s, '/loop 每日掃描 JS 錯誤  ·  --bare 整合 CI/CD 自動部署  ·  /btw 大型重構時同步問答',
            0.4, 7.38, 12.2, 0.38, font_size=13, color=C_DARK)

# ══════════════════════════════════════════════════════════════
# SLIDE 21：快速指令速查表
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '快速指令速查表', '把這頁截圖存手機，隨時查閱！')

commands = [
    ('/voice',           '語音輸入模式（CLI：按住空白鍵說話）'),
    ('/btw [問題]',      '不中斷主任務，插入快速問答'),
    ('/branch',          '從當前節點分叉，嘗試不同方向'),
    ('/teleport',        '將雲端工作階段拉回本機'),
    ('/remote-control',  '讓手機或其他裝置操控本機 Claude'),
    ('/loop [時間] [任務]', '排程重複執行（/loop 5m /check-pr）'),
    ('/batch [任務描述]', '大型任務自動拆解並行執行'),
    ('/add-dir [路徑]',  '加入其他資料夾到操作範圍'),
    ('claude --teleport','啟動時拉取雲端工作階段'),
    ('claude --worktree','在隔離工作樹中執行'),
    ('claude --bare -p "..."','快速非互動執行（啟動快10倍）'),
    ('claude --add-dir [路徑]','啟動時加入多個資料夾'),
    ('claude --agent=[名稱]','啟動自訂代理人'),
]

# Two columns
col_data = [commands[:7], commands[7:]]
for col, data in enumerate(col_data):
    x = 0.35 if col == 0 else 6.75
    add_rect(s, x, 1.38, 6.2, 0.38, fill_color=C_PURPLE)
    label = '工作階段中輸入' if col == 0 else '啟動 CLI 時使用'
    add_textbox(s, label, x + 0.1, 1.4, 6.0, 0.35,
                font_size=13, bold=True, color=C_WHITE)
    for i, (cmd, desc) in enumerate(data):
        y = 1.82 + i * 0.72
        bg = C_LIGHT_PURPLE if i % 2 == 0 else C_WHITE
        add_rect(s, x, y, 6.2, 0.68, fill_color=bg)
        # Command
        add_textbox(s, cmd, x + 0.1, y + 0.05, 2.5, 0.55,
                    font_size=11, bold=True, color=C_PURPLE,
                    font_name='Consolas')
        add_textbox(s, desc, x + 2.65, y + 0.05, 3.4, 0.6,
                    font_size=11, color=C_DARK)

add_textbox(s, '📌 完整文件：https://docs.anthropic.com/claude-code',
            0.4, 7.2, 12.4, 0.38, font_size=13, color=C_PURPLE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 22：總結
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=C_DARK_PURPLE)
add_rect(s, 0, 0, 13.33, 0.25, fill_color=C_AMBER)
add_rect(s, 0, 7.25, 13.33, 0.25, fill_color=C_AMBER)

add_textbox(s, '總結：掌握 Claude Code 的三個關鍵心態', 0.5, 0.6, 12.3, 0.7,
            font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

mindsets = [
    ('🤝', '讓 AI 做重複工作', '/loop、/batch、Hooks', '你只需要說「做什麼」，Claude 自己想「怎麼做」'),
    ('📱', '隨時隨地都能工作', 'Mobile App、/teleport、/remote-control', '不限裝置，不限地點，Claude 陪你跨越空間限制'),
    ('🌲', '大膽嘗試不怕失敗', '--worktree、/branch', '每次改動都有隔離環境，隨時可以回退'),
]

for i, (icon, title, features, desc) in enumerate(mindsets):
    y = 1.45 + i * 1.75
    add_rect(s, 0.5, y, 12.3, 1.6, fill_color=RGBColor(0x2D, 0x23, 0x60))
    add_textbox(s, icon, 0.65, y + 0.3, 0.8, 0.7, font_size=32)
    add_textbox(s, title, 1.55, y + 0.12, 4.5, 0.52,
                font_size=18, bold=True, color=C_AMBER)
    add_textbox(s, features, 1.55, y + 0.65, 4.5, 0.38,
                font_size=12, color=C_GREEN)
    add_textbox(s, desc, 6.3, y + 0.25, 6.3, 0.85,
                font_size=13, color=C_WHITE)

add_textbox(s, '🚀  現在就開始：安裝 Claude Code → 試試 /voice → 說出你的第一個指令！',
            0.5, 6.8, 12.3, 0.45, font_size=15, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)

# ── 儲存 ──────────────────────────────────────────────────────
output_path = 'C:/Users/super/Downloads/money_tutor/money_tutor/report/Claude_Code_15_Features_Guide.pptx'
prs.save(output_path)
print('PPT saved: ' + output_path)
print('Slides: ' + str(len(prs.slides)))
