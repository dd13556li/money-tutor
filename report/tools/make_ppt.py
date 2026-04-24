"""
Claude Code 15 項功能 PPT（v2 修正版）
新手完整操作指南 — 含 CLI 實際指令與輸出示範
修正：add_rect fill/line XML 衝突問題
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 色彩常數 ─────────────────────────────────────────────────
C_PURPLE       = RGBColor(0x5B, 0x4F, 0xCF)
C_DARK_PURPLE  = RGBColor(0x3B, 0x30, 0x9A)
C_AMBER        = RGBColor(0xF5, 0x9E, 0x0B)
C_GREEN        = RGBColor(0x10, 0xB9, 0x81)
C_BLUE         = RGBColor(0x3B, 0x82, 0xF6)
C_RED          = RGBColor(0xEF, 0x44, 0x44)
C_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK         = RGBColor(0x1F, 0x29, 0x37)
C_GRAY_BG      = RGBColor(0xF3, 0xF4, 0xF6)
C_LIGHT_PURPLE = RGBColor(0xED, 0xE9, 0xFF)
C_LIGHT_GREEN  = RGBColor(0xD1, 0xFA, 0xE5)
C_LIGHT_AMBER  = RGBColor(0xFE, 0xF3, 0xC7)
C_LIGHT_BLUE   = RGBColor(0xDB, 0xEA, 0xFE)
C_CODE_BG      = RGBColor(0x1E, 0x1E, 0x2E)
C_CODE_FG      = RGBColor(0xA8, 0xCC, 0x8C)
C_CODE_PROMPT  = RGBColor(0x89, 0xDD, 0xFF)
C_CODE_OUT     = RGBColor(0xCC, 0xCC, 0xCC)
C_PINK         = RGBColor(0xEC, 0x48, 0x99)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── 輔助函數 ─────────────────────────────────────────────────

def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=None):
    """修正版：只在有 fill_color 時呼叫 solid()，避免 XML 衝突"""
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, font_name='Microsoft JhengHei',
                wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name
    return tb


def code_block(slide, lines, left, top, width, height,
               prompt_lines=None, font_size=12):
    """深色代碼框，支援提示行（亮色）與輸出行（灰色）區分"""
    add_rect(slide, left, top, width, height, fill_color=C_CODE_BG)
    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12),
        Inches(width - 0.3), Inches(height - 0.24)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    prompt_set = set(prompt_lines or [])
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = 'Consolas'
        run.font.size = Pt(font_size)
        if i in prompt_set:
            run.font.color.rgb = C_CODE_PROMPT
        elif line.startswith('#'):
            run.font.color.rgb = RGBColor(0x77, 0x99, 0x66)
        else:
            run.font.color.rgb = C_CODE_FG if line.startswith('$') or line.startswith('>') else C_CODE_OUT


def slide_header(slide, title, subtitle='', bg=C_PURPLE):
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=bg)
    add_textbox(slide, title, 0.3, 0.1, 12.0, 0.6,
                font_size=24, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.7, 12.0, 0.36,
                    font_size=13, color=RGBColor(0xCC, 0xCC, 0xFF))
    add_rect(slide, 0, 1.1, 13.33, 6.4, fill_color=C_GRAY_BG)


def info_card(slide, title, lines, left, top, width, height,
              bg=C_LIGHT_BLUE, tc=C_BLUE, icon=''):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_textbox(slide, f'{icon} {title}' if icon else title,
                left+0.12, top+0.1, width-0.24, 0.38,
                font_size=13, bold=True, color=tc)
    tb = slide.shapes.add_textbox(
        Inches(left+0.12), Inches(top+0.52),
        Inches(width-0.24), Inches(height-0.62)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11.5)
        run.font.name = 'Microsoft JhengHei'
        run.font.color.rgb = C_DARK


def step_card(slide, steps, left, top, width, height, bg=C_LIGHT_GREEN):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_textbox(slide, '操作步驟', left+0.12, top+0.1, width-0.24, 0.38,
                font_size=13, bold=True, color=C_GREEN)
    tb = slide.shapes.add_textbox(
        Inches(left+0.12), Inches(top+0.52),
        Inches(width-0.24), Inches(height-0.62)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = step
        run.font.size = Pt(11.5)
        run.font.name = 'Microsoft JhengHei'
        run.font.color.rgb = C_DARK


def example_card(slide, title, lines, left, top, width, height, bg=C_LIGHT_AMBER):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_textbox(slide, f'實際場景：{title}',
                left+0.12, top+0.1, width-0.24, 0.38,
                font_size=13, bold=True, color=C_AMBER)
    tb = slide.shapes.add_textbox(
        Inches(left+0.12), Inches(top+0.52),
        Inches(width-0.24), Inches(height-0.62)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11.5)
        run.font.name = 'Microsoft JhengHei'
        run.font.color.rgb = C_DARK


# ══════════════════════════════════════════════════════════════
# SLIDE 1：封面
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=C_DARK_PURPLE)
add_rect(s, 0, 0,  13.33, 0.22, fill_color=C_AMBER)
add_rect(s, 0, 7.28, 13.33, 0.22, fill_color=C_AMBER)

add_textbox(s, 'Claude Code', 1, 1.1, 11.33, 1.0,
            font_size=56, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, '15 項被低估的隱藏功能', 1, 2.3, 11.33, 0.8,
            font_size=32, bold=False, color=C_AMBER, align=PP_ALIGN.CENTER)
add_textbox(s, '新手完整操作指南  —  從 CLI 指令到實際效果，一步步帶你上手',
            1, 3.25, 11.33, 0.55,
            font_size=16, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

add_rect(s, 2.5, 4.1, 8.33, 0.06, fill_color=RGBColor(0x5B, 0x4F, 0xCF))

add_textbox(s, '來源：Boris Cherny（Claude Code 負責人）推文整理',
            1, 4.4, 11.33, 0.4,
            font_size=13, color=RGBColor(0x99, 0x99, 0xBB), align=PP_ALIGN.CENTER)
add_textbox(s, '整理日期：2026-03-31',
            1, 4.85, 11.33, 0.35,
            font_size=12, color=RGBColor(0x77, 0x77, 0x99), align=PP_ALIGN.CENTER)

add_textbox(s, '適合對象：剛接觸 Claude Code 的新手 / 想深度使用的開發者',
            1.5, 5.7, 10.33, 0.45,
            font_size=14, bold=True, color=RGBColor(0xA8, 0xCC, 0x8C), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2：什麼是 Claude Code？
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '什麼是 Claude Code？', '它不只是「對話框」，而是你電腦上的 AI 代理人')

add_textbox(s, 'Claude Code 是 Anthropic 推出的 AI 編程工具：你用中文（或英文）下指令，它幫你寫程式、執行指令、管理檔案。',
            0.4, 1.25, 12.5, 0.55, font_size=14, bold=False, color=C_DARK)

items = [
    ('讀寫程式碼', '直接操作你電腦上的 .js .py .html 等檔案，不需要複製貼上'),
    ('執行終端機指令', '幫你跑 npm install、git commit、pytest 等指令'),
    ('操控瀏覽器', '打開網頁、點按鈕、截圖、讀取 console 錯誤'),
    ('啟動多個 AI 同時工作', '像指揮一個團隊，10 個 Claude 同時處理不同任務'),
    ('跨裝置連續工作', '電腦開始的任務，手機可以繼續下指令'),
]
for i, (title, desc) in enumerate(items):
    y = 1.95 + i * 0.58
    add_rect(s, 0.4, y, 12.4, 0.52, fill_color=C_WHITE)
    add_textbox(s, f'[{i+1}] {title}', 0.55, y + 0.08, 2.6, 0.38,
                font_size=13, bold=True, color=C_PURPLE)
    add_textbox(s, desc, 3.25, y + 0.1, 9.4, 0.36, font_size=13, color=C_DARK)

add_textbox(s, 'Boris Cherny（Claude Code 負責人）：Anthropic 工程師使用後，程式碼產出成長超過 200%',
            0.4, 5.0, 12.4, 0.42, font_size=12,
            color=RGBColor(0x55, 0x55, 0x99))
add_textbox(s, '安裝後只需在終端機輸入 claude  即可開始使用',
            0.4, 5.55, 12.4, 0.38, font_size=13, bold=True, color=C_GREEN)

code_block(s, [
    '$ claude',
    '> Claude Code v1.x.x  ( ctrl+c to exit )',
    '> 你好！我可以幫你做什麼？請用中文或英文下指令。',
], 0.4, 6.0, 12.4, 0.95, prompt_lines=[0])


# ══════════════════════════════════════════════════════════════
# SLIDE 3：如何安裝 Claude Code
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '如何安裝 Claude Code？', '3 種方式，10 分鐘內完成')

add_textbox(s, '方式 A — CLI（終端機）', 0.4, 1.22, 4.1, 0.4,
            font_size=15, bold=True, color=C_PURPLE)
code_block(s, [
    '# 前提：已安裝 Node.js 18+',
    '$ npm install -g @anthropic-ai/claude-code',
    '',
    '# 登入 Anthropic 帳號',
    '$ claude auth login',
    '  Opening browser for OAuth...',
    '  Logged in as user@example.com',
    '',
    '# 啟動！',
    '$ claude',
], 0.4, 1.65, 4.1, 2.65, prompt_lines=[1, 4, 9])

add_textbox(s, '方式 B — VS Code 擴充', 4.65, 1.22, 4.1, 0.4,
            font_size=15, bold=True, color=C_BLUE)
info_card(s, '步驟', [
    '1. 開啟 VS Code',
    '2. 按 Ctrl+Shift+X 開擴充市集',
    '3. 搜尋 "Claude Code"',
    '4. 點 Install',
    '5. 側欄出現 Claude 圖示',
    '',
    '優點：不需離開 VS Code，',
    '直接在編輯器旁對話',
], 4.65, 1.65, 4.1, 2.65, bg=C_LIGHT_BLUE, tc=C_BLUE, icon='')

add_textbox(s, '方式 C — Desktop App', 8.9, 1.22, 4.1, 0.4,
            font_size=15, bold=True, color=C_GREEN)
info_card(s, '步驟', [
    '1. 前往 claude.ai/download',
    '2. 下載 Mac / Windows 版',
    '3. 安裝後登入帳號',
    '4. 點選 "Code" 模式',
    '',
    '優點：整合內建瀏覽器，',
    '可自動測試網頁效果',
], 8.9, 1.65, 4.1, 2.65, bg=C_LIGHT_GREEN, tc=C_GREEN, icon='')

add_textbox(s, '每次開啟新的工作資料夾時，建議在該資料夾內啟動 claude，這樣 Claude 就知道你的專案在哪裡：',
            0.4, 4.45, 12.4, 0.48, font_size=13, color=C_DARK)
code_block(s, [
    '$ cd /Users/yourname/my-project',
    '$ claude --dangerously-skip-permissions',
    '> 已載入專案：my-project（共 47 個檔案）',
    '> 你好，請告訴我你想做什麼？',
], 0.4, 5.0, 12.4, 1.05, prompt_lines=[0, 1])

add_textbox(s, '--dangerously-skip-permissions  =  跳過每次工具呼叫的確認提示，適合熟悉後日常使用',
            0.4, 6.18, 12.4, 0.38, font_size=12,
            color=RGBColor(0x77, 0x77, 0x99))


# ══════════════════════════════════════════════════════════════
# SLIDE 4：15 項功能總覽
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '15 項被低估功能 — 總覽', 'Boris Cherny 親自點名，早已內建但多數人不知道')

cats = [
    (C_BLUE,   'A 跨裝置', ['Mobile App', '/teleport', '/remote-control', 'Cowork Dispatch']),
    (C_GREEN,  'B 排程並行', ['/loop /schedule', 'Hooks', '--worktree', '/batch']),
    (C_PURPLE, 'C 工具整合', ['Chrome 擴充', 'Desktop 瀏覽器', '/branch', '/btw']),
    (C_AMBER,  'D SDK 協作', ['--bare', '--add-dir', '--agent']),
    (C_PINK,   'E 輸入體驗', ['/voice']),
]
xs = [0.35, 3.0, 5.65, 8.3, 11.0]
for i, (color, cat, feats) in enumerate(cats):
    x = xs[i]
    w = 2.55
    add_rect(s, x, 1.18, w, 0.48, fill_color=color)
    add_textbox(s, cat, x, 1.18, w, 0.48,
                font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for j, feat in enumerate(feats):
        y = 1.76 + j * 0.72
        add_rect(s, x, y, w, 0.65, fill_color=C_WHITE)
        add_textbox(s, feat, x + 0.1, y + 0.13, w - 0.2, 0.4,
                    font_size=13, bold=False, color=C_DARK, align=PP_ALIGN.CENTER)

add_textbox(s, '本投影片接下來逐項說明，每一項都附「實際指令」與「執行結果」範例',
            0.35, 5.05, 12.6, 0.42, font_size=13, bold=True, color=C_DARK)
add_textbox(s, '建議學習順序：先掌握 A 類（跨裝置）和 B 類基礎，再進階到 C、D、E 類',
            0.35, 5.55, 12.6, 0.38, font_size=12, color=RGBColor(0x55, 0x55, 0x99))


# ══════════════════════════════════════════════════════════════
# SLIDE 5：功能 1 — Mobile App
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 1 — Mobile App（行動裝置存取）',
             '類別 A：跨裝置連續工作', bg=C_BLUE)

add_textbox(s, '在 iPhone 或 Android 上，直接存取雲端 Claude Code 工作階段，不需要打開電腦。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

info_card(s, '是什麼？', [
    '下載 Claude App（iOS / Android）',
    '切換至左側「Code」分頁',
    '就能看到雲端工作階段',
    '',
    'Boris Cherny 本人：',
    '「我很大比例的程式修改',
    '是直接從 iPhone 完成的」',
], 0.4, 1.75, 4.0, 3.4, bg=C_LIGHT_BLUE, tc=C_BLUE, icon='')

step_card(s, [
    '1. App Store / Google Play 搜尋 "Claude"',
    '2. 安裝官方 Claude App',
    '3. 登入你的 Anthropic 帳號',
    '4. 左側選單點選 "Code" 分頁',
    '5. 看到正在進行的工作階段',
    '6. 直接輸入指令繼續工作',
], 4.55, 1.75, 4.0, 3.4)

example_card(s, '通勤時繼續改 Bug', [
    '電腦留在辦公室，搭捷運時：',
    '',
    '手機 Claude App -> Code 分頁',
    '',
    '你：「把 login.js 的錯誤訊息改成中文」',
    '',
    'Claude：「已修改 login.js 第 47 行，',
    '     將 "Invalid password" 改為',
    '     "密碼錯誤，請重新輸入"」',
], 8.7, 1.75, 4.2, 3.4)

add_textbox(s, '前提條件：需要 Claude Pro 或 Max 訂閱計畫，工作階段需在電腦上先啟動',
            0.4, 5.3, 12.4, 0.4, font_size=12, color=RGBColor(0x77, 0x77, 0x99))
add_textbox(s, '搭配技巧：與 /remote-control（功能 3）組合使用，可直接操控本機電腦上的工作',
            0.4, 5.78, 12.4, 0.4, font_size=13, bold=True, color=C_BLUE)


# ══════════════════════════════════════════════════════════════
# SLIDE 6：功能 2 — --teleport
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 2 — --teleport（雲端傳回本機）',
             '類別 A：跨裝置連續工作', bg=C_BLUE)

add_textbox(s, '把在 claude.ai/code 網頁上開始的「雲端工作階段」，直接拉回你的本機電腦繼續執行。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

code_block(s, [
    '# 方法一：啟動 claude 時加上旗標',
    '$ claude --teleport',
    '',
    '# 方法二：在已開啟的工作階段中輸入',
    '> /teleport',
    '',
    '# 執行後出現：',
    'Scanning for cloud sessions...',
    'Found 1 session: "my-project" (started 2h ago)',
    'Resuming session locally...',
    'Loaded 12 files, 3 pending changes',
    'Ready!',
], 0.4, 1.75, 6.5, 3.5, prompt_lines=[1, 4])

info_card(s, '什麼時候用？', [
    '在咖啡廳用網頁版 claude.ai/code',
    '用筆電繼續做，想切回家裡桌機',
    '雲端工作階段遇到網路問題',
    '想把雲端 Claude 的結果下載到本機',
    '',
    '用一句話說：',
    '「雲端開始，本機執行」',
], 7.1, 1.75, 5.8, 3.5, bg=C_LIGHT_BLUE, tc=C_BLUE)

add_textbox(s, '與 /remote-control 的差異：teleport = 工作階段移動到本機；remote-control = 從其他裝置操控本機',
            0.4, 5.45, 12.4, 0.42, font_size=13, bold=False, color=C_DARK)
code_block(s, [
    '# 新手常見誤解：兩個功能方向相反',
    '--teleport   : 雲端 --> 本機   (把工作帶回來)',
    '/remote-control : 手機 --> 本機   (從外面操控本機)',
], 0.4, 5.98, 12.4, 0.85, prompt_lines=[])


# ══════════════════════════════════════════════════════════════
# SLIDE 7：功能 3 — /remote-control
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 3 — /remote-control（遠端操控本機）',
             '類別 A：跨裝置連續工作', bg=C_BLUE)

add_textbox(s, '在本機電腦上開啟一個「遠端控制伺服器」，讓你從手機或任何瀏覽器即時下指令。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

code_block(s, [
    '# 在本機電腦的終端機執行（先登入）',
    '$ claude auth login',
    '  Logged in as you@example.com',
    '',
    '# 啟動遠端控制伺服器',
    '$ claude remote-control --name "my-laptop"',
    '  Remote control server started',
    '  Session ID: abc123',
    '  Connect from: claude.ai/code -> Remote Sessions',
    '',
    '# 設定每次都自動開啟（在 claude 設定中）',
    '> /config',
    '  Enable Remote Control for all sessions: YES',
], 0.4, 1.75, 6.8, 3.55, prompt_lines=[1, 5, 11])

step_card(s, [
    '1. 本機執行 claude remote-control --name "xxx"',
    '2. 手機打開 claude.ai/code（瀏覽器）',
    '3. 選單 -> Remote Sessions',
    '4. 看到 "my-laptop" -> 點擊連線',
    '5. 現在手機輸入的指令，本機 Claude 執行',
    '6. 結果即時回傳到手機畫面',
], 7.4, 1.75, 5.5, 3.55)

add_textbox(s, 'Boris Cherny：「我在 /config 開啟 Enable Remote Control，確保每個工作階段都預設可以遠端接手」',
            0.4, 5.45, 12.4, 0.42, font_size=13, bold=True, color=C_BLUE)
add_textbox(s, '應用場景：出門前啟動本機 claude，外出時用手機指揮它繼續跑任務，回家直接看結果',
            0.4, 5.95, 12.4, 0.38, font_size=12, color=C_DARK)


# ══════════════════════════════════════════════════════════════
# SLIDE 8：功能 4 — Cowork Dispatch
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 4 — Cowork Dispatch（MCP 遠端介面）',
             '類別 A：跨裝置連續工作', bg=C_BLUE)

add_textbox(s, 'Claude Desktop App 的進階遠端控制介面，可透過 MCP 連接器操控瀏覽器、管理電子郵件與 Slack。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

info_card(s, 'MCP 是什麼？', [
    'MCP = Model Context Protocol',
    '模型上下文協定',
    '',
    '讓 Claude 連接外部系統的標準，',
    '例如：',
    '• Slack MCP -> 讀寫訊息',
    '• Gmail MCP -> 管理信件',
    '• Browser MCP -> 操控瀏覽器',
    '• File MCP -> 管理檔案',
], 0.4, 1.75, 3.8, 3.55, bg=C_LIGHT_BLUE, tc=C_BLUE)

info_card(s, 'Dispatch 能做什麼？', [
    'Boris Cherny 的日常使用：',
    '',
    '不在寫程式的時候...',
    '• 用 Dispatch 處理 Slack 訊息',
    '• 整理電子郵件',
    '• 管理電腦檔案',
    '',
    '等於把 Claude 變成',
    '你的私人 AI 助理',
], 4.35, 1.75, 3.8, 3.55, bg=C_LIGHT_PURPLE, tc=C_PURPLE)

step_card(s, [
    '1. 安裝 Claude Desktop App',
    '2. 設定 -> MCP Servers',
    '3. 加入想要的 MCP 連接器',
    '4. 重啟 App',
    '5. 在 Dispatch 介面下指令',
    '',
    '範例指令：',
    '「把今天 Slack 未讀訊息整理成摘要」',
], 8.3, 1.75, 4.6, 3.55)

add_textbox(s, '注意：MCP 連接器需要個別設定，初學者建議先掌握功能 1-3，再來研究 Dispatch',
            0.4, 5.45, 12.4, 0.42, font_size=12, color=RGBColor(0x77, 0x77, 0x99))


# ══════════════════════════════════════════════════════════════
# SLIDE 9：功能 5 — /loop + /schedule
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 5 — /loop + /schedule（排程自動執行）',
             '類別 B：排程、並行與自動驗證', bg=C_GREEN)

add_textbox(s, '讓 Claude 定期自動執行任務，你去睡覺它繼續工作。最長可排程一週。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

code_block(s, [
    '# 每 5 分鐘自動執行一次（帶看護模式）',
    '> /loop 5m /babysit',
    '  Loop started: every 5 minutes',
    '  Task: review code, handle PR feedback, rebase if needed',
    '',
    '# 每 30 分鐘整理 Slack 回饋成 PR',
    '> /loop 30m /slack-feedback',
    '  Loop started: every 30 minutes',
    '',
    '# 排程在指定時間執行',
    '> /schedule "2026-04-01 09:00" run-tests',
    '  Scheduled: 2026-04-01 09:00',
    '',
    '# 查看目前跑著的 loops',
    '> /loops list',
    '  [1] every 5m  /babysit   (next: 2m 30s)',
    '  [2] every 30m /slack     (next: 18m)',
], 0.4, 1.75, 7.0, 4.0, prompt_lines=[1, 6, 10, 14])

info_card(s, 'Boris Cherny 的實際設定', [
    '他同時跑著多個 loop：',
    '',
    '/loop 5m /babysit',
    '-> 自動處理 code review 回饋',
    '-> 自動 PR rebase',
    '',
    '/loop 30m /slack-feedback',
    '-> 把 Slack 使用者回饋',
    '   自動整理成 GitHub PR',
    '',
    '他說：「把工作流打包成 loop',
    '是 Claude Code 最值得探索的方向」',
], 7.6, 1.75, 5.3, 4.0, bg=C_LIGHT_GREEN, tc=C_GREEN)

add_textbox(s, '新手建議：先從簡單任務開始，例如 /loop 1h check-spelling 每小時檢查一次拼字',
            0.4, 5.92, 12.4, 0.42, font_size=13, bold=True, color=C_GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 10：功能 6 — Hooks
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 6 — Hooks（工作週期節點注入邏輯）',
             '類別 B：排程、並行與自動驗證', bg=C_GREEN)

add_textbox(s, '在 Claude Code 工作週期的特定「節點」，自動執行你定義的腳本。像事件監聽器，但針對 AI 行為。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

hooks = [
    ('SessionStart',     C_BLUE,   C_LIGHT_BLUE,   '工作階段啟動時', '自動把 CLAUDE.md 等上下文載入，\n讓 Claude 每次開工都有完整背景'),
    ('PreToolUse',       C_GREEN,  C_LIGHT_GREEN,  '執行任何 bash 指令前', '自動記錄操作日誌到 .claude_log.txt，\n完整追蹤 Claude 的每個動作'),
    ('PermissionRequest',C_AMBER,  C_LIGHT_AMBER,  '需要使用者授權時', '把通知轉發到 WhatsApp，\n手機直接批准或拒絕，不必守在電腦前'),
    ('Stop',             C_RED,    RGBColor(0xFE,0xE2,0xE2), 'Claude 停下來等待時', '自動 poke 它繼續執行，\n實現完全無人看管的持續推進'),
]

for i, (hook, tc, bg, when, what) in enumerate(hooks):
    x = 0.35 + (i % 2) * 6.45
    y = 1.75 + (i // 2) * 2.05
    add_rect(s, x, y, 6.2, 1.9, fill_color=bg)
    add_textbox(s, hook, x + 0.12, y + 0.1, 3.0, 0.42,
                font_size=14, bold=True, color=tc)
    add_textbox(s, f'觸發：{when}', x + 0.12, y + 0.55, 5.8, 0.35,
                font_size=12, bold=False, color=C_DARK)
    add_textbox(s, what, x + 0.12, y + 0.92, 5.8, 0.82,
                font_size=11.5, color=C_DARK, wrap=True)

add_textbox(s, 'Stop Hook 的意義：Claude 完成一個任務後停下來，Hook 自動叫它繼續，形成「永動機」式的自動化工作流',
            0.35, 5.95, 12.6, 0.42, font_size=13, bold=True, color=C_GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 11：功能 7 — --worktree
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 7 — --worktree（隔離工作樹）',
             '類別 B：排程、並行與自動驗證', bg=C_GREEN)

add_textbox(s, '在同一個 git repo 裡開多個「隔離的工作副本」，每個 Claude 各自在自己的分支工作，互不干擾。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

code_block(s, [
    '# 開啟一個新的隔離工作樹',
    '$ claude --worktree',
    '  Creating worktree: ../my-project-wt-1',
    '  Branch: claude/worktree-20260401-001',
    '  Claude is working in isolated copy',
    '',
    '# 同時開第二個工作樹（另一個終端機）',
    '$ claude --worktree',
    '  Creating worktree: ../my-project-wt-2',
    '  Branch: claude/worktree-20260401-002',
    '',
    '# 查看所有工作樹',
    '$ git worktree list',
    '  /Users/you/my-project        abc1234 [main]',
    '  /Users/you/my-project-wt-1   def5678 [claude/wt-001]',
    '  /Users/you/my-project-wt-2   ghi9012 [claude/wt-002]',
], 0.4, 1.75, 7.0, 4.05, prompt_lines=[1, 7, 12])

info_card(s, 'Boris Cherny 的使用方式', [
    '「我隨時都有幾十個 Claude 在跑，',
    'worktrees 是核心基礎設施」',
    '',
    '實際場景：',
    '- 同時修 3 個不同的 bug',
    '- 每個 Claude 在自己的分支',
    '- 互不影響，完成後 merge',
    '',
    '非 git 使用者：',
    '可透過 WorktreeCreate hook',
    '自訂工作樹建立邏輯',
], 7.6, 1.75, 5.3, 4.05, bg=C_LIGHT_GREEN, tc=C_GREEN)

add_textbox(s, '最大優點：改壞了不怕！每個工作樹都是獨立副本，隨時可以回退。主分支完全安全',
            0.4, 5.97, 12.4, 0.42, font_size=13, bold=True, color=C_GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 12：功能 8 — /batch
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 8 — /batch（大規模並行代理人）',
             '類別 B：排程、並行與自動驗證', bg=C_GREEN)

add_textbox(s, '2025 年 2 月底推出。自動拆解大型任務，同時開啟數十至數百個 worktree 代理人並行執行。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

code_block(s, [
    '# 大型程式碼遷移：Jest -> Vitest',
    '> /batch migrate from jest to vitest',
    '',
    '  Clarifying questions:',
    '  1. Which directories? [src/, tests/] > src/',
    '  2. Keep original files? [no] > no',
    '  3. Run tests after migration? [yes] > yes',
    '',
    '  Planning...',
    '  Found 147 test files to migrate',
    '  Spawning 20 parallel agents...',
    '',
    '  Agent 1:  migrating src/auth/login.test.js    [done]',
    '  Agent 2:  migrating src/api/user.test.js      [done]',
    '  Agent 3:  migrating src/utils/format.test.js  [running]',
    '  ...',
    '  All 147 files migrated in 4m 32s',
], 0.4, 1.75, 7.5, 4.15, prompt_lines=[1])

info_card(s, '適合用 /batch 的任務', [
    '大型重構（統一命名規則）',
    '跨檔案格式轉換',
    '批次加入測試或文件',
    '多語言翻譯所有頁面',
    '套用新版 API 到全部呼叫點',
    '',
    '不適合的任務：',
    '需要人工判斷的設計決策',
    '步驟之間有相依性的任務',
], 8.1, 1.75, 4.8, 4.15, bg=C_LIGHT_GREEN, tc=C_GREEN)

add_textbox(s, '注意：/batch 會使用大量 API tokens，執行前確認任務規模，避免意外費用',
            0.4, 6.07, 12.4, 0.38, font_size=12, color=RGBColor(0xEF, 0x44, 0x44))


# ══════════════════════════════════════════════════════════════
# SLIDE 13：功能 9 — Chrome 擴充套件
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 9 — Chrome 擴充套件（瀏覽器操控）',
             '類別 C：開發工具鏈整合', bg=C_PURPLE)

add_textbox(s, '讓 Claude 直接操控 Chrome 瀏覽器，讀取 console 錯誤、點擊按鈕、截圖測試 UI。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

add_textbox(s, 'Boris Cherny 的比喻：', 0.4, 1.78, 12.4, 0.38,
            font_size=13, bold=True, color=C_DARK)
add_textbox(s, '「如果你請一個工程師建網站，但不讓他開瀏覽器，成果會好嗎？」',
            0.4, 2.18, 12.4, 0.42, font_size=14, color=C_PURPLE)
add_textbox(s, '「給他瀏覽器，他就會不斷測試、迭代，直到結果正確。」',
            0.4, 2.62, 12.4, 0.38, font_size=13, color=C_DARK)

code_block(s, [
    '# Claude 看到你的網頁有 UI 問題後自動測試',
    '',
    'Claude: 我來打開瀏覽器測試登入流程...',
    '  Opening http://localhost:3000/login',
    '  Clicking "Login" button',
    '  Console error found:',
    '    TypeError: Cannot read property "token" of undefined',
    '    at auth.js:47',
    '  Fixing auth.js line 47...',
    '  Reloading page...',
    '  Login successful! No console errors.',
], 0.4, 3.12, 7.5, 2.78, prompt_lines=[])

step_card(s, [
    '1. Chrome Web Store 搜尋',
    '   "Claude Code" (選官方 beta 版)',
    '2. 安裝擴充套件',
    '3. 在 Claude Code 工作階段中',
    '   說：「用瀏覽器測試首頁」',
    '4. Claude 自動打開並測試',
    '',
    '目前狀態：Beta 版',
    '功能持續更新中',
], 8.1, 3.12, 4.8, 2.78)

add_textbox(s, 'CLI / VS Code 使用者：也可安裝 Chrome 擴充套件，與 CLI 工作階段配合使用',
            0.4, 6.05, 12.4, 0.38, font_size=12, color=RGBColor(0x77, 0x77, 0x99))


# ══════════════════════════════════════════════════════════════
# SLIDE 14：功能 10 — Desktop 瀏覽器 / 功能 11 — /branch
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 10 + 11 — Desktop 內建瀏覽器 & /branch',
             '類別 C：開發工具鏈整合', bg=C_PURPLE)

add_textbox(s, '功能 10：Claude Desktop App 內建瀏覽器', 0.4, 1.22, 6.2, 0.42,
            font_size=15, bold=True, color=C_PURPLE)
code_block(s, [
    '# Claude Desktop App 自動啟動 web server 測試',
    '',
    'Claude: 我來啟動開發伺服器並測試...',
    '  $ npm run dev',
    '  Server running at http://localhost:5173',
    '  Opening built-in browser...',
    '  Screenshot taken: homepage.png',
    '  Issue found: button misaligned on mobile',
    '  Fixing CSS...',
    '  Re-testing... All good!',
], 0.4, 1.72, 6.2, 2.75, prompt_lines=[])

info_card(s, '優點', [
    '不需要離開 Claude App',
    '整個開發 + 測試流程閉環',
    'CLI 用戶用 Chrome 擴充達成',
    '同樣效果',
], 0.4, 4.55, 6.2, 1.7, bg=C_LIGHT_PURPLE, tc=C_PURPLE)

add_textbox(s, '功能 11：/branch（工作節點分叉）', 7.0, 1.22, 5.9, 0.42,
            font_size=15, bold=True, color=C_PURPLE)
code_block(s, [
    '# 從現在的節點分叉出新嘗試方向',
    '> /branch',
    '',
    '  Saving checkpoint: "before-refactor"',
    '  Forked to new branch',
    '  You can try different approaches',
    '  and switch back anytime',
    '',
    '# 切回主線',
    '> claude --resume --fork-session',
    '  Restored checkpoint: "before-refactor"',
], 7.0, 1.72, 5.9, 2.75, prompt_lines=[1, 9])

info_card(s, '使用場景', [
    '「試試看這個寫法對不對」',
    '→ /branch 先存一個節點',
    '→ 試了發現不行',
    '→ 切回原本繼續',
    '',
    '解決「改壞了怕回不去」',
    '的心理負擔',
], 7.0, 4.55, 5.9, 1.7, bg=C_LIGHT_PURPLE, tc=C_PURPLE)


# ══════════════════════════════════════════════════════════════
# SLIDE 15：功能 12 — /btw
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 12 — /btw（不中斷插入旁支問題）',
             '類別 C：開發工具鏈整合', bg=C_PURPLE)

add_textbox(s, '在 Claude 還在執行長任務時，插入快速問題。Claude 回答後繼續原本工作，完全不中斷。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

code_block(s, [
    '# Claude 正在跑大型重構（需要 10 分鐘）',
    'Claude: [running] Refactoring authentication module...',
    '        Processed 12/47 files...',
    '',
    '# 你突然想知道某個函數的用途',
    '> /btw  什麼是 generateToken() 這個函數在做什麼？',
    '',
    'Claude: [btw] generateToken() 在 auth/utils.js 第 23 行，',
    '        作用是產生一個 32 位元的隨機 token 用於 session 驗證。',
    '        目前重構不會影響它的行為。',
    '',
    'Claude: [resuming main task]',
    '        Processed 13/47 files...',
], 0.4, 1.75, 8.5, 3.75, prompt_lines=[5])

info_card(s, '為什麼這很有用？', [
    '以前：必須等 Claude 做完才問',
    '或打斷任務（損失進度）',
    '',
    '現在：任何時候都能插入問題',
    '',
    '典型使用場景：',
    '• 確認某個設計決策',
    '• 快速查某個函數的作用',
    '• 詢問「這樣做會不會有問題？」',
    '',
    '答完後 Claude 自動繼續',
], 9.1, 1.75, 3.8, 3.75, bg=C_LIGHT_PURPLE, tc=C_PURPLE)

add_textbox(s, '新手提醒：/btw 問題要簡短，複雜問題等主任務結束再問，避免混淆 Claude 的上下文',
            0.4, 5.68, 12.4, 0.42, font_size=13, bold=True, color=C_PURPLE)
add_textbox(s, '搭配技巧：與 --worktree 配合，主工作樹跑長任務，另開工作樹做獨立問題',
            0.4, 6.18, 12.4, 0.38, font_size=12, color=C_DARK)


# ══════════════════════════════════════════════════════════════
# SLIDE 16：功能 13 — --bare
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 13 — --bare（SDK 快速啟動旗標）',
             '類別 D：SDK 與多 Repo 協作優化', bg=C_AMBER)

add_textbox(s, '在批次處理或 CI/CD 場景，加上 --bare 可跳過 CLAUDE.md 等設定載入，啟動時間縮短最多 10 倍。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

code_block(s, [
    '# 一般互動模式（載入所有設定）',
    '$ claude',
    '  Loading CLAUDE.md... (640 lines)',
    '  Loading MCP servers...',
    '  Loading settings...',
    '  Ready! (startup: 2.8s)',
    '',
    '# --bare 模式（跳過所有設定）',
    '$ claude -p --bare "檢查 login.js 有無 XSS 風險"',
    '  Ready! (startup: 0.3s)',
    '  login.js: No XSS vulnerabilities found.',
    '',
    '# 在批次腳本中使用',
    '$ for f in src/*.js; do',
    '$   claude -p --bare "簡述這個檔案的功能" < "$f"',
    '$ done',
], 0.4, 1.75, 7.5, 3.8, prompt_lines=[1, 8, 13, 14, 15])

info_card(s, '什麼時候用 --bare？', [
    '互動使用：不需要 --bare',
    '（設定只載入一次）',
    '',
    '適合 --bare 的場景：',
    '• Shell 腳本批次處理',
    '• CI/CD 自動化流程',
    '• TypeScript/Python SDK 呼叫',
    '• 快速的單次問答 (-p 模式)',
    '',
    'Boris Cherny：',
    '「可縮短 SDK 啟動最多 10 倍」',
], 8.1, 1.75, 4.8, 3.8, bg=C_LIGHT_AMBER, tc=C_AMBER)

code_block(s, [
    '# Python SDK 使用範例',
    'import anthropic',
    'client = anthropic.Anthropic()',
    '# 加 --bare 等效：建立 client 時跳過 claude 設定',
    '# 適合大量批次 API 呼叫，回應更快',
], 0.4, 5.72, 12.4, 1.05, prompt_lines=[])


# ══════════════════════════════════════════════════════════════
# SLIDE 17：功能 14 — --add-dir
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 14 — --add-dir（跨 Repo 操作）',
             '類別 D：SDK 與多 Repo 協作優化', bg=C_AMBER)

add_textbox(s, '啟動 claude 時加入其他 repo 的目錄，讓 Claude 能讀寫多個不同的專案資料夾。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

code_block(s, [
    '# 在主要 repo 啟動，同時加入其他 repo',
    '$ cd /projects/frontend',
    '$ claude --add-dir /projects/backend --add-dir /projects/shared-lib',
    '',
    '  Loaded: /projects/frontend  (312 files)',
    '  Added:  /projects/backend   (189 files) [full write access]',
    '  Added:  /projects/shared-lib (64 files) [full write access]',
    '',
    '# 工作階段中動態加入',
    '> /add-dir /projects/design-tokens',
    '  Added: /projects/design-tokens (12 files)',
    '',
    '# 或寫入 settings.json 讓整個團隊共用',
    '# .claude/settings.json',
    '# { "additionalDirectories": ["/projects/shared-lib"] }',
], 0.4, 1.75, 7.5, 3.75, prompt_lines=[1, 2, 9])

info_card(s, '重要：不只是「看到」', [
    '一般以為 --add-dir 只是讓',
    'Claude 能讀取其他目錄',
    '',
    '實際上：',
    'Claude 在這些目錄有',
    '完整的讀寫操作權限',
    '',
    'Boris Cherny：',
    '「它不只讓 Claude 看到',
    '那個目錄，還賦予完整',
    '操作權限」',
], 8.1, 1.75, 4.8, 3.75, bg=C_LIGHT_AMBER, tc=C_AMBER)

example_card(s, '前後端同步修改', [
    '前端發現 API 格式改變了',
    '需要同時改前端呼叫 + 後端回傳',
    '',
    '以前：開兩個終端機、兩個 Claude',
    '現在：一個 Claude 同時搞定兩個 repo',
], 0.4, 5.67, 12.4, 1.3, bg=C_LIGHT_AMBER)


# ══════════════════════════════════════════════════════════════
# SLIDE 18：功能 15 — --agent
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 15 — --agent（自訂子代理人）',
             '類別 D：SDK 與多 Repo 協作優化', bg=C_AMBER)

add_textbox(s, '在 .claude/agents/ 資料夾定義代理人設定檔，建立具有不同行為、不同工具限制的 Claude 子代理人。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

code_block(s, [
    '# 建立 .claude/agents/ 資料夾',
    '$ mkdir -p .claude/agents',
    '',
    '# 建立 readonly 代理人（只能讀，不能寫）',
    '$ cat .claude/agents/readonly.md',
    '---',
    'name: readonly',
    'description: 只能讀取程式碼，無法修改或執行',
    'allowedTools: [ReadFile, ListFiles, SearchCode]',
    '---',
    'You are a code reviewer. Only read and analyze code.',
    'Never modify files or run commands.',
    '',
    '# 啟動自訂代理人',
    '$ claude --agent=readonly',
    '  Loaded agent: readonly',
    '  Tools available: ReadFile, ListFiles, SearchCode',
    '  [Edit and Bash tools are DISABLED]',
], 0.4, 1.75, 7.2, 4.15, prompt_lines=[1, 4, 14])

info_card(s, 'Boris Cherny 的 ReadOnly 代理人', [
    '設定：',
    '只開放讀取工具',
    '無法編輯檔案',
    '無法執行 bash 指令',
    '',
    '用途：',
    '• 安全地做程式碼審查',
    '• 讓不信任的人使用',
    '• CI 環境只讀分析',
    '',
    '你也可以建立：',
    '• 只能寫特定資料夾',
    '• 只能用某些工具',
], 7.7, 1.75, 5.2, 4.15, bg=C_LIGHT_AMBER, tc=C_AMBER)

add_textbox(s, '進階：可以在 --agent 中設定不同的系統提示（system prompt），讓 Claude 扮演不同的專業角色',
            0.4, 6.08, 12.4, 0.42, font_size=13, bold=True, color=C_AMBER)


# ══════════════════════════════════════════════════════════════
# SLIDE 19：功能 16 — /voice
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '功能 16 — /voice（語音輸入）',
             '類別 E：輸入體驗', bg=C_PINK)

add_textbox(s, '不用打字，直接說話下指令。Boris Cherny 本人親口確認：「我日常以說話為主、鍵盤為輔」。',
            0.4, 1.22, 12.4, 0.45, font_size=14, color=C_DARK)

methods = [
    ('CLI 終端機', C_LIGHT_PURPLE, C_PURPLE, [
        '$ claude',
        '> /voice',
        '',
        '  [Voice mode activated]',
        '  Hold SPACE to speak',
        '  Release to send',
        '',
        '  你說：「把 CSS 的顏色',
        '         全部換成設計系統的',
        '         CSS 變數」',
        '',
        '  Claude 執行中...',
    ], [1]),
    ('Desktop App', C_LIGHT_GREEN, C_GREEN, [
        '1. 開啟 Claude Desktop App',
        '2. 點擊右下角麥克風按鈕',
        '3. 說出你的指令',
        '4. 放開 / 點停止按鈕',
        '',
        '優點：介面更直覺，',
        '有視覺回饋顯示',
        '語音辨識狀態',
    ], None),
    ('iPhone / iPad', C_LIGHT_AMBER, C_AMBER, [
        '1. 設定 -> 一般 -> 鍵盤',
        '2. 開啟「啟用聽寫」',
        '',
        '使用時：',
        '鍵盤上按麥克風圖示',
        '對著手機說話',
        '',
        '不需要安裝任何額外',
        '軟體，系統層級整合',
    ], None),
]
xs2 = [0.35, 4.5, 8.65]
for i, (title, bg, tc, lines, plines) in enumerate(methods):
    x = xs2[i]
    add_rect(s, x, 1.75, 4.05, 0.42, fill_color=bg)
    add_textbox(s, title, x, 1.75, 4.05, 0.42,
                font_size=13, bold=True, color=tc, align=PP_ALIGN.CENTER)
    if i == 0:
        code_block(s, lines, x, 2.2, 4.05, 3.35, prompt_lines=plines or [], font_size=11)
    else:
        add_rect(s, x, 2.2, 4.05, 3.35, fill_color=bg)
        tb = s.shapes.add_textbox(
            Inches(x + 0.15), Inches(2.32), Inches(3.75), Inches(3.1)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for j, line in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(12)
            run.font.name = 'Microsoft JhengHei'
            run.font.color.rgb = C_DARK

add_textbox(s, '語音輸入特別適合：長段說明型指令（「幫我寫一個函數，功能是…，輸入是…，要注意…」）',
            0.35, 5.68, 12.6, 0.42, font_size=13, bold=True, color=C_PINK)
add_textbox(s, '鍵盤輸入特別適合：精確的程式碼片段、檔案路徑、需要特定格式的指令',
            0.35, 6.18, 12.6, 0.38, font_size=12, color=C_DARK)


# ══════════════════════════════════════════════════════════════
# SLIDE 20：CLI 速查表
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 'CLI 快速指令速查表', '把這頁截圖存手機，隨時查閱！')

cols = [
    ('跨裝置', C_BLUE, C_LIGHT_BLUE, [
        ('claude --teleport',              '雲端階段移到本機'),
        ('/teleport',                      '（工作階段中輸入）'),
        ('claude remote-control',          '開啟遠端控制伺服器'),
        ('/remote-control',                '（工作階段中輸入）'),
        ('/config',                        '設定預設開啟遠端控制'),
    ]),
    ('排程並行', C_GREEN, C_LIGHT_GREEN, [
        ('/loop 5m task-name',             '每 5 分鐘執行'),
        ('/loop 1h task-name',             '每 1 小時執行'),
        ('/schedule "時間" task',           '指定時間執行'),
        ('/loops list',                    '查看執行中的 loops'),
        ('claude --worktree',              '開新隔離工作樹'),
        ('/batch 任務說明',                 '大規模並行執行'),
    ]),
    ('工具整合', C_PURPLE, C_LIGHT_PURPLE, [
        ('/branch',                        '從目前節點分叉'),
        ('claude --resume --fork-session', '切回主線'),
        ('/btw 快速問題',                  '不中斷插入旁支問答'),
    ]),
    ('SDK 協作', C_AMBER, C_LIGHT_AMBER, [
        ('claude -p --bare "問題"',        '快速單次問答（無設定）'),
        ('claude --add-dir /path',         '加入其他目錄'),
        ('/add-dir /path',                 '（工作階段中動態加入）'),
        ('claude --agent=name',            '啟動自訂代理人'),
        ('/voice',                         '語音輸入模式'),
    ]),
]
xs3 = [0.3, 3.45, 6.6, 9.65]
widths3 = [3.05, 3.05, 2.95, 3.33]
for i, (cat, tc, bg, cmds) in enumerate(cols):
    x, w = xs3[i], widths3[i]
    add_rect(s, x, 1.18, w, 0.4, fill_color=tc)
    add_textbox(s, cat, x, 1.18, w, 0.4,
                font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, 1.6, w, 5.3, fill_color=bg)
    tb = s.shapes.add_textbox(
        Inches(x + 0.1), Inches(1.68), Inches(w - 0.2), Inches(5.1)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for j, (cmd, desc) in enumerate(cmds):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(5 if j > 0 else 0)
        run = p.add_run()
        run.text = cmd
        run.font.name = 'Consolas'
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = tc

        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = f'  {desc}'
        run2.font.name = 'Microsoft JhengHei'
        run2.font.size = Pt(10.5)
        run2.font.color.rgb = C_DARK


# ══════════════════════════════════════════════════════════════
# SLIDE 21：新手學習路徑
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, '新手學習路徑建議', '由淺入深，穩步掌握所有 16 項功能')

phases = [
    ('第 1 週\n基礎入門', C_GREEN, C_LIGHT_GREEN, [
        '安裝 Claude Code（CLI 方式）',
        '練習用中文下指令',
        '試試 /voice 語音輸入',
        '用 claude 讀懂一個你不熟的程式碼',
    ]),
    ('第 2 週\n跨裝置', C_BLUE, C_LIGHT_BLUE, [
        '安裝 Mobile App，試試跨裝置',
        '啟動 claude remote-control',
        '用手機下一個指令給電腦執行',
        '嘗試 /branch 建立一個分叉點',
    ]),
    ('第 3 週\n自動化', C_AMBER, C_LIGHT_AMBER, [
        '設定 SessionStart Hook（自動載入上下文）',
        '試試 claude --worktree 隔離工作',
        '用 /loop 30m 跑一個自動任務',
        '體驗 /btw 不中斷插入問題',
    ]),
    ('第 4 週\n進階', C_PURPLE, C_LIGHT_PURPLE, [
        '嘗試 /batch 執行一個批次任務',
        '設定 --add-dir 跨 repo 操作',
        '建立一個 --agent 自訂代理人',
        '配置 Chrome 擴充套件測試 UI',
    ]),
]

for i, (phase, tc, bg, items) in enumerate(phases):
    x = 0.35 + i * 3.2
    add_rect(s, x, 1.18, 3.05, 0.72, fill_color=tc)
    add_textbox(s, phase, x, 1.18, 3.05, 0.72,
                font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, 1.92, 3.05, 3.6, fill_color=bg)
    tb = s.shapes.add_textbox(
        Inches(x + 0.12), Inches(2.02), Inches(2.82), Inches(3.4)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(4 if j > 0 else 0)
        run = p.add_run()
        run.text = f'{j+1}. {item}'
        run.font.size = Pt(12)
        run.font.name = 'Microsoft JhengHei'
        run.font.color.rgb = C_DARK

add_textbox(s, '重要原則：每項功能先「能用就好」，不需要一次掌握所有細節。用到了再深入研究。',
            0.35, 5.68, 12.6, 0.42, font_size=13, bold=True, color=C_DARK)
add_textbox(s, 'Cowork Dispatch 和 Hooks 較複雜，建議第 3-4 週以後再深入研究',
            0.35, 6.15, 12.6, 0.38, font_size=12, color=RGBColor(0x77, 0x77, 0x99))


# ══════════════════════════════════════════════════════════════
# SLIDE 22：總結
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=C_DARK_PURPLE)
add_rect(s, 0, 0, 13.33, 0.22, fill_color=C_AMBER)
add_rect(s, 0, 7.28, 13.33, 0.22, fill_color=C_AMBER)

add_textbox(s, '掌握 Claude Code 的三個關鍵心態',
            0.5, 0.45, 12.3, 0.65,
            font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

mindsets = [
    ('AI 做重複工作', '/loop  /batch  Hooks',
     '你只需要說「做什麼」，Claude 自己想「怎麼做」。\n把重複任務打包成 loop，下班後繼續跑。'),
    ('隨時隨地工作', 'Mobile App  /teleport  /remote-control',
     '不限裝置，不限地點。搭捷運可以下指令，\n手機操控家裡電腦，雲端與本機無縫切換。'),
    ('大膽嘗試不怕失敗', '--worktree  /branch',
     '每次改動都有隔離環境，隨時可以回退。\n試錯的成本幾乎是零，勇敢嘗試新方向。'),
]
for i, (title, feats, desc) in enumerate(mindsets):
    y = 1.3 + i * 1.78
    add_rect(s, 0.4, y, 12.5, 1.65, fill_color=RGBColor(0x2D, 0x23, 0x60))
    add_textbox(s, title, 0.6, y + 0.1, 4.2, 0.52,
                font_size=18, bold=True, color=C_AMBER)
    add_textbox(s, feats, 0.6, y + 0.65, 4.2, 0.42,
                font_size=12, color=C_GREEN)
    add_textbox(s, desc, 5.1, y + 0.2, 7.6, 0.88,
                font_size=13, color=C_WHITE, wrap=True)

add_textbox(s, '現在就開始：安裝 Claude Code  ->  輸入 /voice  ->  說出你的第一個指令！',
            0.5, 6.68, 12.3, 0.48,
            font_size=16, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)


# ── 儲存 ──────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   'Claude_Code_15_Features_Guide.pptx')
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
